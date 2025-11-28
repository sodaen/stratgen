from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

def generate_pptx(output_path: str, title: str = "Marketingstrategie", sections: list[str] | None = None):
    prs = Presentation()
    # Titelfolie
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Auto-generated MVP · Bitte Zahlen/Quellen prüfen"

    # Inhaltsfolien
    sections = sections or [
        "Executive Summary",
        "Markt & Kategorie",
        "Zielgruppe & Archetypen",
        "Brand Narrative",
        "Kanäle & Content",
        "Maßnahmenplan",
        "KPIs",
        "Roadmap",
        "Budgetrahmen"
    ]
    for sec in sections:
        s = prs.slides.add_slide(prs.slide_layouts[1])  # Titel + Inhalt
        s.shapes.title.text = sec
        tf = s.placeholders[1].text_frame
        tf.text = "🟡 TODO: Inhalte werden hier eingefügt.\nPlatzhalter für Zahlen/Quellen."
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
