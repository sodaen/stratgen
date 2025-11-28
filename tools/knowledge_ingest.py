#!/usr/bin/env python3
import os, sys, json, csv, time, uuid, sqlite3
from pathlib import Path

try:
    from pptx import Presentation
except Exception:
    Presentation = None

ROOT = Path(__file__).resolve().parents[1]
DATA = ROOT / "data"
RAW_DIR = DATA / "raw"
EXP_DIR = DATA / "exports"
DB = DATA / "stratgen.db"

def read_text_file(p: Path) -> str:
    try: return p.read_text(encoding="utf-8", errors="ignore")
    except: return ""

def read_json_file(p: Path) -> str:
    try:
        obj = json.loads(p.read_text(encoding="utf-8", errors="ignore"))
        # extrahiere alle string values grob
        acc = []
        def walk(x):
            if isinstance(x, dict):
                for v in x.values(): walk(v)
            elif isinstance(x, list):
                for v in x: walk(v)
            elif isinstance(x, str):
                acc.append(x)
        walk(obj)
        return "\n".join(acc)
    except:
        return ""

def read_csv_file(p: Path) -> str:
    try:
        out = []
        with p.open("r", encoding="utf-8", errors="ignore", newline="") as f:
            for row in csv.reader(f):
                out.append(" | ".join(row))
        return "\n".join(out)
    except:
        return ""

def read_pptx_file(p: Path) -> str:
    if Presentation is None:
        return ""
    try:
        prs = Presentation(str(p))
        parts = []
        for i, s in enumerate(prs.slides, 1):
            buf = []
            for shape in s.shapes:
                if hasattr(shape, "text") and shape.text:
                    buf.append(shape.text)
            if buf:
                parts.append(f"[Slide {i}]\n" + "\n".join(buf))
        return "\n\n".join(parts)
    except:
        return ""

def read_any(p: Path) -> str:
    ext = p.suffix.lower()
    if ext in (".txt", ".md"): return read_text_file(p)
    if ext == ".json": return read_json_file(p)
    if ext == ".csv":  return read_csv_file(p)
    if ext == ".pptx": return read_pptx_file(p)
    return ""

def chunk_text(txt: str, max_len=1200) -> list[str]:
    txt = " ".join(txt.split())
    chunks, i = [], 0
    while i < len(txt):
        chunks.append(txt[i:i+max_len])
        i += max_len
    return [c for c in chunks if c.strip()]

def upsert_doc(db, doc_id, source, path):
    db.execute("INSERT OR IGNORE INTO knowledge_docs(id,source,path,created_at) VALUES (?,?,?,?)",
               (doc_id, source, path, int(time.time())))

def upsert_chunks(db, doc_id, chunks):
    # simple: clear and insert (idempotent genug bei gleicher doc_id)
    db.execute("DELETE FROM knowledge_fts WHERE doc_id=?", (doc_id,))
    db.executemany("INSERT INTO knowledge_fts(doc_id,content) VALUES (?,?)",
                   [(doc_id, c) for c in chunks])

def scan_dir(db, base: Path, source: str):
    inserted = skipped = 0
    for p in sorted(base.rglob("*")):
        if not p.is_file(): continue
        if p.suffix.lower() not in (".txt",".md",".json",".csv",".pptx"): continue
        rel = str(p.relative_to(ROOT))
        doc_id = f"{source}:{rel}"
        text = read_any(p)
        if not text.strip():
            skipped += 1; continue
        chunks = chunk_text(text)
        upsert_doc(db, doc_id, source, rel)
        upsert_chunks(db, doc_id, chunks)
        inserted += 1
    return inserted, skipped

def main():
    os.makedirs(DATA, exist_ok=True)
    db = sqlite3.connect(DB)
    ins1, sk1 = scan_dir(db, RAW_DIR, "raw")
    ins2, sk2 = scan_dir(db, EXP_DIR, "exports")
    db.commit(); db.close()
    print(json.dumps({"ok": True, "inserted": ins1+ins2, "skipped": sk1+sk2}))

if __name__ == "__main__":
    main()
