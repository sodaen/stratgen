#!/usr/bin/env python3
import json, sys
from pathlib import Path
try:
    from pptx import Presentation
except Exception:
    print("python-pptx fehlt.")
    sys.exit(2)

def sniff(pptx_path: Path) -> dict:
    try:
        prs = Presentation(str(pptx_path))
    except Exception:
        return {}
    out = {"file": pptx_path.name, "slides": len(prs.slides), "layouts": []}
    for s in prs.slides:
        layout = {"placeholders": [], "shapes": len(s.shapes)}
        try:
            for ph in s.placeholders:
                layout["placeholders"].append({
                    "type": getattr(ph, "placeholder_format", None) and ph.placeholder_format.type,
                    "left": ph.left, "top": ph.top, "width": ph.width, "height": ph.height
                })
        except Exception:
            pass
        out["layouts"].append(layout)
    return out

base = Path("data/raw")
man = []
for pptx in list(base.glob("**/*.pptx"))[:200]:
    info = sniff(pptx)
    if info: man.append(info)

outp = Path("data/style/archetypes.json")
outp.parent.mkdir(parents=True, exist_ok=True)
outp.write_text(json.dumps({"files": man}, indent=2), encoding="utf-8")
print(f"wrote {outp} (files={len(man)})")
