import json
from pathlib import Path
from pptx import Presentation

MASTER = Path("styles/master.pptx")
OUT = Path("styles/style.json")

def rgb_from_scheme(srgb):
    # srgb is like <a:srgbClr val="FFFFFF"> wrapper
    try:
        val = srgb.get("val")
        if val:
            return f"#{val}"
    except Exception:
        pass
    return None

def main():
    style = {
        "font_title": "Arial",
        "font_body": "Arial",
        "size_title_pt": 36,
        "size_body_pt": 16,
        "primary": "#111111",
        "secondary": "#555555",
        "accent1": "#1F77B4",
        "accent2": "#FF7F0E",
        "accent3": "#2CA02C"
    }
    if MASTER.exists():
        prs = Presentation(MASTER.as_posix())
        try:
            theme = prs.slide_master.part.theme
            cs = theme.theme_elements.clr_scheme
            # Versuch, 2–3 Akzentfarben aus dem Schema zu ziehen
            style["primary"]   = rgb_from_scheme(cs.dk1.srgbClr) or style["primary"]
            style["secondary"] = rgb_from_scheme(cs.lt1.srgbClr) or style["secondary"]
            style["accent1"]   = rgb_from_scheme(cs.accent1.srgbClr) or style["accent1"]
            style["accent2"]   = rgb_from_scheme(cs.accent2.srgbClr) or style["accent2"]
            style["accent3"]   = rgb_from_scheme(cs.accent3.srgbClr) or style["accent3"]
        except Exception:
            pass
        # Schriftarten aus erster Layout-Textbox heuristisch erfassen
        try:
            for layout in prs.slide_layouts:
                for shp in layout.shapes:
                    if hasattr(shp, "text_frame") and shp.text_frame and shp.text_frame.paragraphs:
                        run = shp.text_frame.paragraphs[0].runs[0]
                        if run and run.font:
                            if run.font.name:
                                style["font_body"] = run.font.name
                            break
        except Exception:
            pass

    OUT.parent.mkdir(parents=True, exist_ok=True)
    OUT.write_text(json.dumps(style, indent=2), encoding="utf-8")
    print(f"Styles gespeichert → {OUT.resolve()}")

if __name__ == "__main__":
    main()
