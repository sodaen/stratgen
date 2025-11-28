from pptx import Presentation
from pathlib import Path

tmpl = Path("styles/master.pptx")
assert tmpl.exists(), "styles/master.pptx fehlt"
prs = Presentation(tmpl.as_posix())

def tokens_in_layout(layout):
    found = set()
    for shp in getattr(layout, "shapes", []):
        if hasattr(shp, "text_frame") and shp.text_frame and shp.text_frame.text:
            txt = shp.text_frame.text
            for t in ["#KUNDENLOGO","#TITEL","#THEMA","#DATUM","#ZWISCHENFOLIE","#TOPIC","#SUBTOPIC","#INHALT","#ENDE"]:
                if t in txt: found.add(t)
    return found

for i, layout in enumerate(prs.slide_layouts):
    print(f"Layout {i}: {sorted(tokens_in_layout(layout))}")
