#!/usr/bin/env python3
import json, sys
from pathlib import Path
import httpx
from pptx import Presentation
from pptx.util import Pt

BASE = "http://127.0.0.1:8000"
API_KEY = "changeme-local"  # falls du API-Key-Check aktiv hast
OUTLINE_FILE = Path("/tmp/outline_raw.json")
OUT = Path("/tmp/strategy_deck.pptx")

if not OUTLINE_FILE.exists():
    print("❌ /tmp/outline_raw.json fehlt. Bitte zuerst: make outline")
    sys.exit(2)

o = json.loads(OUTLINE_FILE.read_text())
agenda = o.get("agenda", [])
brief = {"scope":"Marketingstrategie & Social Media","market":"Elektronikhandel","region":"DACH"}

prs = Presentation()
def add_slide(title, bullets):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for b in bullets:
        p = body.add_paragraph() if body.text else body.paragraphs[0]
        p.text = b
        p.font.size = Pt(18)

with httpx.Client(timeout=30) as client:
    headers = {"Content-Type":"application/json","X-API-Key":API_KEY}
    # Titelfolie
    tslide = prs.slides.add_slide(prs.slide_layouts[0])
    tslide.shapes.title.text = o.get("project_title","Strategy")
    tslide.placeholders[1].text = o.get("customer_name","")

    for sec in agenda:
        topic = sec.get("topic","")
        subs = sec.get("subtopics") or [None]
        for sub in subs:
            req = {
                "customer_name": o.get("customer_name","Acme GmbH"),
                "topic": topic,
                "mode": "facts",
                "brief": brief
            }
            if sub: req["subtopic"] = sub
            r = client.post(f"{BASE}/content/preview", headers=headers, json=req)
            r.raise_for_status()
            data = r.json()
            bullets = data.get("bullets") or []
            title = topic if not sub else f"{topic} — {sub}"
            add_slide(title, bullets[:6])

prs.save(OUT)
print(f"✅ Deck geschrieben: {OUT}")
