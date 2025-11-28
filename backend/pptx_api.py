# -*- coding: utf-8 -*-
from __future__ import annotations

import os, json, time, urllib.request, urllib.error
from pathlib import Path
from fastapi import APIRouter

router = APIRouter(prefix="/pptx", tags=["pptx"])

def _api_base() -> str:
    return os.environ.get("STRATGEN_INTERNAL_URL", "http://127.0.0.1:8011").rstrip("/")

def _jget(path: str, timeout: int = 10):
    try:
        with urllib.request.urlopen(_api_base() + path, timeout=timeout) as r:
            return json.loads(r.read().decode("utf-8"))
    except Exception:
        return None

def _jpost(path: str, payload: dict | None = None, timeout: int = 30):
    try:
        req = urllib.request.Request(
            _api_base() + path,
            data=json.dumps(payload or {}).encode("utf-8"),
            headers={"content-type": "application/json"},
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=timeout) as r:
            return json.loads(r.read().decode("utf-8"))
    except Exception:
        return None

def _write_sidecar(project_id: str, out_path: Path) -> str | None:
    """Best-effort: schreibt <pptx>.json neben die Datei und gibt json_url zurück."""
    json_url = None
    try:
        base = _api_base()
        with urllib.request.urlopen(f"{base}/projects/{project_id}") as r:
            proj = json.loads(r.read().decode("utf-8"))
        meta = (proj or {}).get("project", {}).get("meta", {}) or {}
        sidecar = {
            "project_id": project_id,
            "export_file": out_path.name,
            "slide_plan": meta.get("slide_plan", []),
            "slide_plan_len": meta.get("slide_plan_len", 0),
            "created_at": int(time.time()),
        }
        side_path = out_path.with_suffix(out_path.suffix + ".json")
        side_path.write_text(json.dumps(sidecar, ensure_ascii=False, indent=2), encoding="utf-8")
        json_url = f"/exports/download/{side_path.name}"
    except Exception as e:
        try:
            import logging
            logging.getLogger("stratgen.api").warning("sidecar failed: %s", e)
        except Exception:
            pass
    return json_url

@router.post("/render_from_project/{project_id}")
def render_from_project(project_id: str):
    """Projekt laden, ggf. Plan erzeugen, einfache PPTX bauen und speichern."""
    # 1) Projekt holen
    data = _jget(f"/projects/{project_id}") or {}
    if not data or not data.get("ok"):
        return {"ok": False, "error": "project not found", "project_id": project_id}
    project = data.get("project") or {}
    meta = project.get("meta") or {}

    # 2) Plan sicherstellen
    plan = meta.get("slide_plan") or []
    if not plan:
        gen = _jpost(f"/projects/{project_id}/generate", {}) or {}
        plan = ((((gen.get("project") or {}).get("meta") or {}).get("slide_plan")) or [])

    # 3) Präsentation bauen
    try:
        from pptx import Presentation
        from pptx.util import Pt  # optional fürs Textframe
    except Exception:
        return {"ok": False, "error": "python-pptx not available"}

    prs = Presentation()
    if not plan:
        plan = [{"title": project.get("topic") or "Presentation", "bullets": []}]

    for item in list(plan)[:100]:
        title = str(item.get("title") or "Slide").strip() or "Slide"
        bullets = item.get("bullets") or []
        layout = prs.slide_layouts[1 if bullets else 0]
        slide = prs.slides.add_slide(layout)
        # Titel
        try:
            slide.shapes.title.text = title
        except Exception:
            pass
        # Bullets
        if bullets:
            try:
                body = slide.shapes.placeholders[1].text_frame
                body.clear()
                body.text = str(bullets[0])
                for b in bullets[1:]:
                    body.add_paragraph().text = str(b)
            except Exception:
                pass

    # 4) Speichern
    exports_dir = Path("data/exports")
    exports_dir.mkdir(parents=True, exist_ok=True)
    ts = int(time.time())
    safe_pid = project_id.replace("/", "-")
    out_name = f"project-{safe_pid}-{ts}.pptx"
    out_path = exports_dir / out_name
    prs.save(out_path)

    # 5) Sidecar schreiben
    json_url = _write_sidecar(project_id, out_path)

    return {"ok": True, "path": str(out_path), "url": f"/exports/download/{out_name}", "json_url": json_url}
