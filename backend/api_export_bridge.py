"""
API Export Bridge - Verbindet Frontend mit Export-Services.
Frontend ruft: /api/export/{format}/{session_id}
"""

from fastapi import APIRouter, HTTPException
from pathlib import Path
import json
import time
import logging

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api", tags=["export-bridge"])

EXPORTS_DIR = Path("data/exports")
EXPORTS_DIR.mkdir(parents=True, exist_ok=True)

LIVE_SESSIONS_DIR = Path("data/live_sessions")
LIVE_SESSIONS_DIR.mkdir(parents=True, exist_ok=True)


@router.get("/export/{format}/{session_id}")
@router.post("/export/{format}/{session_id}")
async def export_session(format: str, session_id: str):
    """
    Exportiert eine Session zu PPTX/PDF/etc.
    Wird vom Frontend (Editor, Generator) aufgerufen.
    """
    
    if format not in ["pptx", "pdf", "json", "md"]:
        raise HTTPException(400, f"Unsupported format: {format}")
    
    slides = []
    title = "Strategie Präsentation"
    colors = None
    company = ""
    
    # 1. Versuche aus Live-Generator API (aktive Session)
    try:
        import httpx
        resp = httpx.get(f"http://127.0.0.1:8011/live/slides/{session_id}", timeout=5)
        if resp.status_code == 200:
            live_data = resp.json()
            slides = live_data.get("slides", [])
            if slides:
                logger.info(f"Found {len(slides)} slides from live API")
    except Exception as e:
        logger.debug(f"Live API: {e}")
    
    # 2. Fallback: Gespeicherte Session-Datei
    if not slides:
        session_file = LIVE_SESSIONS_DIR / f"{session_id}.json"
        if session_file.exists():
            try:
                with open(session_file) as f:
                    session_data = json.load(f)
                slides = session_data.get("slides", [])
                title = session_data.get("title", title)
                logger.info(f"Found {len(slides)} slides from file")
            except:
                pass
    
    # 3. Versuche aus Project
    if not slides:
        try:
            import httpx
            resp = httpx.get(f"http://127.0.0.1:8011/projects/{session_id}", timeout=5)
            if resp.status_code == 200:
                proj_data = resp.json()
                if proj_data.get("ok"):
                    project = proj_data.get("project", {})
                    meta = project.get("meta", {})
                    slides = meta.get("slide_plan", [])
                    title = project.get("topic", title)
        except:
            pass
    
    if not slides:
        raise HTTPException(404, f"No slides found for session: {session_id}")
    
    # 4. Export durchführen
    ts = int(time.time())
    
    if format == "pptx":
        try:
            from services.pptx_designer_v2 import PPTXDesignerV2
            
            # Default Corporate Farben
            if not colors:
                colors = {
                    "primary": "#1E40AF",
                    "secondary": "#3B82F6",
                    "accent": "#10B981",
                    "background": "#FFFFFF",
                    "text": "#111827"
                }
            
            designer = PPTXDesignerV2(colors=colors)
            pptx_bytes = designer.create_presentation(
                slides=slides,
                title=title,
                company=company,
                include_sources_slide=True,
                customer_name=company,
                use_images=True
            )
            
            filename = f"export-{session_id[:8]}-{ts}.pptx"
            out_path = EXPORTS_DIR / filename
            
            with open(out_path, 'wb') as f:
                f.write(pptx_bytes)

            # ── SELF-LEARNING HOOK ──────────────────────────
            try:
                from services.self_learning import self_learning as _sl
                import threading as _threading
                _export_path_for_sl = str(out_path)
                _session_id_for_sl = session_id
                def _run_learning():
                    _sl.on_export_complete(
                        export_path=_export_path_for_sl,
                        session_id=_session_id_for_sl,
                        export_type="pptx"
                    )
                _threading.Thread(target=_run_learning, daemon=True).start()
                logger.info("Self-learning hook triggered for session %s", session_id)
            except Exception as _sl_err:
                logger.debug("Self-learning hook skipped: %s", _sl_err)
            # ── /SELF-LEARNING HOOK ─────────────────────────

            
            return {
                "ok": True,
                "url": f"/exports/download/{filename}",
                "path": str(out_path),
                "slides_count": len(slides),
                "knowledge_slides": sum(1 for s in slides if s.get('knowledge_used', False))
            }
            
        except Exception as e:
            logger.error(f"PPTX Designer failed: {e}")
            # Fallback zu einfachem Export
            try:
                from pptx import Presentation
                
                prs = Presentation()
                for slide_data in slides[:100]:
                    layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(layout)
                    if slide.shapes.title:
                        slide.shapes.title.text = slide_data.get("title", "Slide")
                    if slide_data.get("bullets"):
                        try:
                            body = slide.shapes.placeholders[1].text_frame
                            body.clear()
                            body.text = str(slide_data["bullets"][0])
                            for b in slide_data["bullets"][1:]:
                                body.add_paragraph().text = str(b)
                        except:
                            pass
                
                filename = f"export-{session_id[:8]}-{ts}.pptx"
                out_path = EXPORTS_DIR / filename
                prs.save(str(out_path))
                
                return {
                    "ok": True,
                    "url": f"/exports/download/{filename}",
                    "slides_count": len(slides)
                }
            except Exception as e2:
                raise HTTPException(500, f"PPTX export failed: {str(e2)}")
    
    elif format == "json":
        filename = f"export-{session_id[:8]}-{ts}.json"
        out_path = EXPORTS_DIR / filename
        
        with open(out_path, 'w') as f:
            json.dump({"title": title, "slides": slides}, f, indent=2, ensure_ascii=False)
        
        return {
            "ok": True,
            "url": f"/exports/download/{filename}",
            "slides_count": len(slides)
        }
    
    elif format == "md":
        filename = f"export-{session_id[:8]}-{ts}.md"
        out_path = EXPORTS_DIR / filename
        
        md_content = f"# {title}\n\n"
        for i, slide in enumerate(slides, 1):
            md_content += f"## Slide {i}: {slide.get('title', 'Untitled')}\n\n"
            if slide.get('content'):
                md_content += f"{slide['content']}\n\n"
            if slide.get('bullets'):
                for b in slide['bullets']:
                    md_content += f"- {b}\n"
                md_content += "\n"
        
        with open(out_path, 'w') as f:
            f.write(md_content)
        
        return {
            "ok": True,
            "url": f"/exports/download/{filename}",
            "slides_count": len(slides)
        }
    
    raise HTTPException(400, f"Format {format} not implemented")


@router.post("/export/learn/{session_id}")
async def trigger_learning_manually(session_id: str, export_type: str = "pptx"):
    """Manuell Self-Learning für eine bereits exportierte Session triggern."""
    try:
        from services.self_learning import self_learning as _sl
        candidates = list(EXPORTS_DIR.glob(f"export-{session_id[:8]}-*.{export_type}"))
        if not candidates:
            candidates = list(EXPORTS_DIR.glob(f"*{session_id[:8]}*.{export_type}"))
        if not candidates:
            return {"ok": False, "error": "Keine Export-Datei gefunden"}
        export_path = str(sorted(candidates, key=lambda f: f.stat().st_mtime, reverse=True)[0])
        result = _sl.on_export_complete(
            export_path=export_path, session_id=session_id, export_type=export_type
        )
        return {"ok": True, "learning_result": result}
    except Exception as e:
        return {"ok": False, "error": str(e)}

