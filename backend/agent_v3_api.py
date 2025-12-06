# -*- coding: utf-8 -*-
"""
backend/agent_v3_api.py
=======================
Agent V3.1 - Mit Quick Wins + Agent Intelligence (Stufe 1)

QUICK WINS:
1. NLG-Module Integration (12 vorhandene Module aktiviert)
2. Template-Learner Integration (lernt aus /raw)
3. Feedback-Loop Integration (speichert Qualitätsdaten)
4. Multi-Model Support (verschiedene Ollama-Modelle für verschiedene Tasks)

STUFE 1 - AGENT INTELLIGENCE:
- Multi-Step Reasoning Chain
- Kontextuelles Slide-Generieren
- Dynamische Slide-Anzahl
- Iterative Verbesserung
- Selbstkritik und Auto-Revision

Pipeline:
1. ANALYZE   → Briefing verstehen, Komplexität einschätzen
2. PLAN      → Struktur dynamisch planen
3. RESEARCH  → RAG + Knowledge sammeln
4. DRAFT     → Content generieren (kontextbewusst)
5. CRITIQUE  → Selbstkritik
6. REVISE    → Verbesserung (iterativ)
7. VISUALIZE → Charts + Assets
8. RENDER    → PPTX Export
"""
from __future__ import annotations
import os
import time
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from pydantic import BaseModel, Field
from enum import Enum

from fastapi import APIRouter, HTTPException, BackgroundTasks

# ============================================
# IMPORTS - Services
# ============================================

# LLM Content Generation
try:
    from services.llm_content import (
        generate_bullets,
        generate_summary,
        generate_persona,
        generate_critique,
        generate_metrics,
        generate_slide_content,
        generate_from_template,
        check_ollama,
    )
    HAS_LLM_CONTENT = True
except ImportError:
    HAS_LLM_CONTENT = False

# Direkter LLM-Zugriff über services/llm.py
try:
    from services.llm import generate as llm_generate, is_enabled as llm_enabled
    HAS_LLM = True
except ImportError:
    llm_generate = None
    llm_enabled = None
    HAS_LLM = False

# NLG Modules (Quick Win 1)
try:
    from services.nlg import (
        mod_exec_summary,
        mod_audience,
        mod_insights,
        mod_positioning,
        mod_messaging,
        mod_channels,
        mod_metrics,
        mod_risks,
        mod_next_steps
    )
    HAS_NLG = True
except ImportError:
    HAS_NLG = False

# Spezifische NLG-Module
NLG_MODULES = {}
try:
    from services.nlg.personas import generate as nlg_personas
    NLG_MODULES["personas"] = nlg_personas
except ImportError:
    pass
try:
    from services.nlg.gtm_basics import generate as nlg_gtm
    NLG_MODULES["gtm"] = nlg_gtm
except ImportError:
    pass
try:
    from services.nlg.funnel import generate as nlg_funnel
    NLG_MODULES["funnel"] = nlg_funnel
except ImportError:
    pass
try:
    from services.nlg.kpis import generate as nlg_kpis
    NLG_MODULES["kpis"] = nlg_kpis
except ImportError:
    pass
try:
    from services.nlg.market_sizing import generate as nlg_market
    NLG_MODULES["market_sizing"] = nlg_market
except ImportError:
    pass
try:
    from services.nlg.competitive import generate as nlg_competitive
    NLG_MODULES["competitive"] = nlg_competitive
except ImportError:
    pass
try:
    from services.nlg.value_proof import generate as nlg_value
    NLG_MODULES["value_proof"] = nlg_value
except ImportError:
    pass
try:
    from services.nlg.risks_mitigations import generate as nlg_risks
    NLG_MODULES["risks"] = nlg_risks
except ImportError:
    pass

# Asset Tagger
try:
    from services.asset_tagger import (
        scan_uploads_directory,
        match_assets_to_slides,
        get_asset_suggestions
    )
    HAS_ASSET_TAGGER = True
except ImportError:
    HAS_ASSET_TAGGER = False

# Chart Generator
try:
    from services.chart_generator import (
        create_bar_chart,
        create_pie_chart,
        create_line_chart,
        create_timeline,
        create_funnel_chart,
        create_gauge_chart,
        create_comparison_matrix,
        auto_create_chart
    )
    HAS_CHART_GEN = True
except ImportError:
    HAS_CHART_GEN = False

# Visual Intelligence (Stufe 3)
try:
    from services.visual_intelligence import (
        enhance_slide_visuals,
        enhance_all_slides,
        generate_chart_for_slide,
        recommend_images_for_slide,
        recommend_layout,
        check_status as visual_status
    )
    HAS_VISUAL_INTELLIGENCE = True
except ImportError:
    HAS_VISUAL_INTELLIGENCE = False
    enhance_all_slides = None

# Learning & Adaptation (Stufe 4)
try:
    from services.learning_adaptation import (
        record_feedback,
        predict_quality,
        record_quality_result,
        get_merged_style,
        get_improvement_suggestions,
        learn_from_all_templates,
        get_feedback_stats,
        check_status as learning_status
    )
    HAS_LEARNING = True
except ImportError:
    HAS_LEARNING = False
    predict_quality = None
    get_merged_style = None

# Killer-Features (Stufe 6)
try:
    from services.persona_engine import generate_persona, generate_personas, analyze_audience, list_archetypes
    from services.competitive_intelligence import analyze_competition, generate_swot, generate_battle_card
    from services.roi_calculator import generate_business_case, calculate_project_roi
    from services.story_engine import create_story_structure, build_narrative_arc, generate_hooks
    from services.briefing_analyzer import analyze as analyze_briefing, BriefingQuality
    HAS_KILLER_FEATURES = True
except ImportError as e:
    HAS_KILLER_FEATURES = False
    print(f"Killer-Features nicht verfügbar: {e}")

# Advanced Features (Stufe 7)
try:
    from services.slide_dna_analyzer import extract_slide_dna, get_optimal_structure
    from services.semantic_slide_matcher import get_slide_suggestions
    from services.brand_voice_extractor import get_writing_guidelines
    from services.argument_engine import check_deck_consistency
    from services.content_intelligence import score_deck_complexity, detect_knowledge_gaps
    from services.live_generator import live_generator
    HAS_ADVANCED_FEATURES = True
except ImportError as e:
    HAS_ADVANCED_FEATURES = False
    print(f"Advanced Features nicht verfügbar: {e}")

# Feature Orchestrator (Integration Layer)
try:
    from services.feature_orchestrator import (
        orchestrator,
        orchestrate_analysis,
        orchestrate_quality_check,
        check_status as orchestrator_status
    )
    HAS_ORCHESTRATOR = True
except ImportError as e:
    HAS_ORCHESTRATOR = False
    orchestrator = None
    print(f"Feature Orchestrator nicht verfügbar: {e}")

# Multi-Modal Export (Stufe 5)
try:
    from services.multimodal_export import (
        export_to_html,
        export_to_pdf,
        export_to_markdown,
        export_to_json,
        export_presentation,
        get_available_formats,
        check_status as export_status
    )
    HAS_MULTIMODAL_EXPORT = True
except ImportError:
    HAS_MULTIMODAL_EXPORT = False
    export_to_html = None
    export_to_pdf = None

# Template Learner (Quick Win 2)
try:
    from services.template_learner import (
        suggest_structure,
        scan_templates,
        get_statistics as get_template_stats
    )
    HAS_TEMPLATE_LEARNER = True
except ImportError:
    HAS_TEMPLATE_LEARNER = False

# Feedback Loop (Quick Win 3)
try:
    from services.feedback_loop import (
        get_quality_score,
        get_improvement_suggestions,
        record_feedback,
        analyze_patterns
    )
    HAS_FEEDBACK = True
except ImportError:
    HAS_FEEDBACK = False

# Knowledge/RAG (Basic)
try:
    from services.knowledge import search as knowledge_search, scan_dir as scan_knowledge
    HAS_KNOWLEDGE = True
except ImportError:
    knowledge_search = None
    HAS_KNOWLEDGE = False

# Knowledge Enhanced (Stufe 2)
try:
    from services.knowledge_enhanced import (
        multi_source_search,
        extract_facts_from_results,
        build_research_context,
        research_for_slide,
        CitationManager,
        check_status as knowledge_status
    )
    HAS_KNOWLEDGE_ENHANCED = True
except ImportError:
    HAS_KNOWLEDGE_ENHANCED = False
    multi_source_search = None
    CitationManager = None

# ============================================
# KONFIGURATION
# ============================================

EXPORTS_DIR = os.getenv("STRATGEN_EXPORTS_DIR", "data/exports")
UPLOADS_DIR = os.getenv("STRATGEN_UPLOADS_DIR", "data/uploads")
RAW_DIR = os.getenv("STRATGEN_RAW_DIR", "data/raw")

# Multi-Model Konfiguration (Quick Win 4)
OLLAMA_MODELS = {
    "default": os.getenv("LLM_MODEL", "mistral"),
    "fast": "mistral",           # Schnell für einfache Tasks
    "quality": "llama3:8b",      # Qualität für wichtige Inhalte
    "creative": "mistral",       # Kreativ für Ideen
    "analysis": "qwen2.5:7b-instruct",  # Analyse/Reasoning
}

# Agent Intelligence Konfiguration
AGENT_CONFIG = {
    "max_iterations": 3,         # Max Verbesserungsschleifen
    "quality_threshold": 7.5,    # Mindestqualität (1-10)
    "min_slides": 5,
    "max_slides": 50,
    "context_window_slides": 3,  # Slides vor/nach für Kontext
}

# ============================================
# ROUTER & MODELS
# ============================================

router = APIRouter(prefix="/agent", tags=["agent-v3"])


class DeckSize(str, Enum):
    SHORT = "short"
    MEDIUM = "medium"
    LARGE = "large"
    AUTO = "auto"  # Agent entscheidet


class AgentV3Request(BaseModel):
    """Request für Agent V3 Run."""
    # Briefing
    topic: str
    brief: str = ""
    customer_name: str = ""
    industry: str = ""
    audience: str = ""
    
    # Konfiguration
    deck_size: DeckSize = DeckSize.MEDIUM
    language: str = "de"
    style: str = "professional"
    
    # Optionen
    use_rag: bool = True
    generate_charts: bool = True
    match_assets: bool = True
    include_critique: bool = True
    export_pptx: bool = True
    export_html: bool = False      # HTML/Reveal.js Export
    export_pdf: bool = False       # PDF Export
    export_markdown: bool = False  # Markdown Export
    export_json: bool = False      # JSON Export
    
    # Agent Intelligence
    auto_improve: bool = True    # Iterative Verbesserung
    max_iterations: int = 2      # Max Verbesserungsschleifen
    use_nlg_modules: bool = True # NLG-Module nutzen
    learn_from_templates: bool = True  # Aus /raw lernen
    
    # Advanced
    k: int = 5  # RAG top-k
    use_cases: List[str] = Field(default_factory=list)
    custom_sections: List[str] = Field(default_factory=list)
    llm_model: str = ""     # Spezifisches Modell erzwingen (z.B. "llama3:8b")


class SlideContent(BaseModel):
    """Ein einzelner Slide."""
    type: str
    title: str
    bullets: List[str] = Field(default_factory=list)
    notes: str = ""
    layout_hint: str = "Title and Content"
    citations: List[str] = Field(default_factory=list)
    chart: Optional[str] = None
    image: Optional[str] = None
    has_chart: bool = False
    confidence: float = 0.8  # Wie sicher ist der Agent über diesen Slide


class AgentPlan(BaseModel):
    """Der Plan des Agents."""
    complexity: str = "medium"  # low, medium, high
    estimated_slides: int = 15
    key_topics: List[str] = Field(default_factory=list)
    recommended_sections: List[str] = Field(default_factory=list)
    research_queries: List[str] = Field(default_factory=list)
    rationale: str = ""


class AgentV3Response(BaseModel):
    """Response von Agent V3."""
    ok: bool
    run_id: str
    project_id: Optional[str] = None
    
    # Ergebnisse
    slides: List[Dict[str, Any]] = Field(default_factory=list)
    slide_count: int = 0
    
    # Agent Intelligence
    plan: Optional[Dict[str, Any]] = None
    iterations: int = 1
    final_quality: float = 0.0
    
    # Exports
    pptx_url: Optional[str] = None
    pdf_url: Optional[str] = None
    html_url: Optional[str] = None      # Stufe 5
    markdown_url: Optional[str] = None  # Stufe 5
    json_url: Optional[str] = None      # Stufe 5
    exports: Optional[Dict[str, Any]] = None  # Alle Exports
    
    # Metriken
    quality_score: Optional[float] = None
    duration_s: float = 0
    
    # Details
    phases: Dict[str, Any] = Field(default_factory=dict)
    warnings: List[str] = Field(default_factory=list)
    improvements: List[str] = Field(default_factory=list)


# ============================================
# HELPER FUNCTIONS
# ============================================

def _generate_run_id() -> str:
    """Generiert eine eindeutige Run-ID."""
    return f"v3-{datetime.now().strftime('%Y%m%d-%H%M%S')}-{os.urandom(4).hex()}"


def _select_model(task: str, override: str = "") -> str:
    """Wählt das beste Modell für einen Task (Quick Win 4)."""
    if override:
        return override
    return OLLAMA_MODELS.get(task, OLLAMA_MODELS["default"])


def _llm_call(prompt: str, task: str = "default", model_override: str = "", max_tokens: int = 1000) -> str:
    """Zentraler LLM-Call mit Model-Selection."""
    # Nutze services/llm.py generate() direkt
    if not HAS_LLM or not llm_generate:
        return ""
    
    if llm_enabled and not llm_enabled():
        return ""
    
    model = _select_model(task, model_override)
    
    try:
        result = llm_generate(
            prompt=prompt,
            model=model,
            max_tokens=max_tokens
        )
        if isinstance(result, dict):
            if result.get("ok"):
                return result.get("response", "")
            return ""
        return str(result)
    except Exception as e:
        return ""


# ============================================
# PHASE 1: ANALYZE (Agent Intelligence)
# ============================================

def phase_analyze(req: AgentV3Request) -> AgentPlan:
    """
    Phase 1: Analysiert das Briefing und erstellt einen Plan.
    Der Agent "denkt" über die Aufgabe nach.
    """
    # Komplexität einschätzen
    brief_length = len(req.brief)
    has_use_cases = len(req.use_cases) > 0
    has_custom = len(req.custom_sections) > 0
    
    if brief_length < 100 and not has_use_cases:
        complexity = "low"
        estimated = 7
    elif brief_length > 500 or has_use_cases or has_custom:
        complexity = "high"
        estimated = 25
    else:
        complexity = "medium"
        estimated = 15
    
    # Deck Size Override
    if req.deck_size == DeckSize.SHORT:
        estimated = min(estimated, 10)
    elif req.deck_size == DeckSize.LARGE:
        estimated = max(estimated, 25)
    elif req.deck_size == DeckSize.AUTO:
        pass  # Agent entscheidet
    
    # LLM-basierte Analyse (wenn verfügbar)
    key_topics = []
    research_queries = []
    rationale = ""
    
    if HAS_LLM_CONTENT and req.brief:
        analysis_prompt = f"""Analysiere dieses Briefing für eine Strategie-Präsentation:

Thema: {req.topic}
Briefing: {req.brief}
Kunde: {req.customer_name}
Branche: {req.industry}
Zielgruppe: {req.audience}

Gib zurück (JSON):
{{
    "key_topics": ["Liste der 3-5 wichtigsten Themen"],
    "research_queries": ["3-5 Suchbegriffe für Knowledge Base"],
    "recommended_sections": ["empfohlene Slide-Typen"],
    "rationale": "Kurze Begründung für die Struktur"
}}

NUR JSON, keine Erklärung."""

        try:
            result = _llm_call(analysis_prompt, task="analysis", max_tokens=500)
            # JSON extrahieren
            import re
            json_match = re.search(r'\{.*\}', result, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group())
                key_topics = data.get("key_topics", [])[:5]
                research_queries = data.get("research_queries", [])[:5]
                rationale = data.get("rationale", "")
        except Exception:
            pass
    
    # Fallback Key Topics
    if not key_topics:
        key_topics = [req.topic]
        if req.industry:
            key_topics.append(f"{req.industry} Strategie")
        if req.customer_name:
            key_topics.append(f"Lösung für {req.customer_name}")
    
    # Fallback Research Queries
    if not research_queries:
        research_queries = [
            req.topic,
            f"{req.topic} {req.industry}",
            f"{req.topic} Best Practices"
        ]
    
    return AgentPlan(
        complexity=complexity,
        estimated_slides=estimated,
        key_topics=key_topics,
        research_queries=research_queries,
        rationale=rationale or f"Basierend auf Briefing-Länge ({brief_length} Zeichen) und Komplexität"
    )


# ============================================
# PHASE 2: RESEARCH (Enhanced RAG)
# ============================================

def phase_research(req: AgentV3Request, plan: AgentPlan) -> Dict[str, Any]:
    """
    Phase 2: Enhanced Research mit Multi-Source RAG und Fact Extraction.
    Nutzt knowledge_enhanced.py für bessere Ergebnisse.
    """
    result = {
        "sources": [],
        "facts": [],
        "citations": [],
        "template_insights": [],
        "citation_manager": None,
    }
    
    if not req.use_rag:
        return result
    
    context = {
        "topic": req.topic,
        "industry": req.industry,
        "customer_name": req.customer_name,
        "brief": req.brief
    }
    
    # === ENHANCED RESEARCH (Stufe 2) ===
    if HAS_KNOWLEDGE_ENHANCED and multi_source_search:
        try:
            # Citation Manager erstellen
            result["citation_manager"] = CitationManager() if CitationManager else None
            
            # Multi-Source-Suche
            for query in plan.research_queries[:5]:
                research = multi_source_search(
                    query=query,
                    context=context,
                    sources=["knowledge", "templates", "uploads"],
                    k=5
                )
                
                # Ergebnisse sammeln
                for res in research.results:
                    if res.path not in [s.get("path") for s in result["sources"]]:
                        result["sources"].append({
                            "path": res.path,
                            "title": res.title,
                            "snippet": res.snippet[:400],
                            "score": res.score,
                            "source_type": res.source_type
                        })
                        
                        # Citation hinzufügen
                        if result["citation_manager"]:
                            result["citation_manager"].add_source(
                                res.path, res.title, res.snippet
                            )
                
                # Fakten extrahieren
                if research.results:
                    extracted_facts = extract_facts_from_results(research.results)
                    for fact in extracted_facts:
                        if fact.text not in result["facts"]:
                            result["facts"].append(fact.text)
            
        except Exception as e:
            # Fallback zu Basic Search
            pass
    
    # === FALLBACK: Basic Knowledge Search ===
    if not result["sources"] and HAS_KNOWLEDGE and knowledge_search:
        all_results = []
        for query in plan.research_queries[:5]:
            try:
                search_result = knowledge_search(query, limit=3, semantic=1)
                if search_result.get("ok"):
                    all_results.extend(search_result.get("results", []))
            except Exception:
                pass
        
        # Deduplizieren
        seen_paths = set()
        for item in all_results:
            path = item.get("path", "")
            if path not in seen_paths:
                seen_paths.add(path)
                result["sources"].append({
                    "path": path,
                    "title": item.get("title") or Path(path).stem,
                    "snippet": item.get("snippet", "")[:300]
                })
                if item.get("snippet"):
                    result["facts"].append(item["snippet"][:400])
    
    # === Template-Insights (Quick Win 2) ===
    if req.learn_from_templates and HAS_TEMPLATE_LEARNER:
        try:
            scan_templates(RAW_DIR)
            stats = get_template_stats()
            if stats.get("ok") and stats.get("patterns"):
                patterns = stats["patterns"]
                result["template_insights"] = [
                    f"Durchschnittlich {patterns.get('avg_bullets_per_slide', 4)} Bullets pro Slide",
                    f"Durchschnittlich {patterns.get('avg_slides_per_deck', 15)} Slides pro Deck",
                ]
        except Exception:
            pass
    
    return result


# ============================================
# PHASE 3: STRUCTURE (Dynamic Planning)
# ============================================

def phase_structure(req: AgentV3Request, plan: AgentPlan, research: Dict[str, Any]) -> List[Dict[str, Any]]:
    """
    Phase 3: Plant die Deck-Struktur dynamisch basierend auf Analyse.
    """
    slides = []
    
    # Template-Learner Vorschlag (Quick Win 2)
    if req.learn_from_templates and HAS_TEMPLATE_LEARNER:
        try:
            size_map = {"short": "short", "medium": "medium", "large": "large", "auto": "medium"}
            suggestion = suggest_structure(
                deck_size=size_map.get(req.deck_size.value, "medium"),
                topic=req.topic,
                industry=req.industry
            )
            if suggestion.get("ok") and suggestion.get("slides"):
                return suggestion["slides"]
        except Exception:
            pass
    
    # Dynamische Struktur basierend auf Komplexität
    base_modules = [
        {"type": "title", "title": req.topic, "required": True},
        {"type": "executive_summary", "title": "Executive Summary", "required": True},
        {"type": "problem", "title": "Herausforderung", "required": True},
        {"type": "solution", "title": "Unser Ansatz", "required": True},
        {"type": "benefits", "title": "Ihr Nutzen", "required": True},
        {"type": "next_steps", "title": "Nächste Schritte", "required": True},
        {"type": "contact", "title": "Kontakt", "required": True},
    ]
    
    optional_modules = [
        {"type": "agenda", "title": "Agenda", "complexity": "medium"},
        {"type": "use_case", "title": "Use Case", "complexity": "medium"},
        {"type": "roi", "title": "ROI & Business Case", "complexity": "medium"},
        {"type": "roadmap", "title": "Roadmap", "complexity": "low"},
        {"type": "team", "title": "Unser Team", "complexity": "high"},
        {"type": "competitive", "title": "Marktpositionierung", "complexity": "high"},
        {"type": "risks", "title": "Risiken & Mitigation", "complexity": "high"},
        {"type": "personas", "title": "Zielgruppen", "complexity": "high"},
        {"type": "kpis", "title": "KPIs & Metriken", "complexity": "medium"},
    ]
    
    complexity_order = {"low": 0, "medium": 1, "high": 2}
    plan_complexity = complexity_order.get(plan.complexity, 1)
    
    # Required Slides
    for mod in base_modules:
        slides.append({
            "type": mod["type"],
            "title": mod["title"],
            "order": len(slides)
        })
    
    # Optional Slides basierend auf Komplexität
    for mod in optional_modules:
        mod_complexity = complexity_order.get(mod.get("complexity", "medium"), 1)
        if mod_complexity <= plan_complexity:
            # Vor "next_steps" einfügen
            insert_idx = len(slides) - 2
            slides.insert(insert_idx, {
                "type": mod["type"],
                "title": mod["title"],
                "order": insert_idx
            })
    
    # Use Cases hinzufügen
    if req.use_cases:
        insert_idx = next((i for i, s in enumerate(slides) if s["type"] == "benefits"), len(slides) - 2)
        for i, uc in enumerate(req.use_cases[:5]):
            slides.insert(insert_idx + i, {
                "type": "use_case_detail",
                "title": f"Use Case: {uc}",
                "use_case": uc,
                "order": insert_idx + i
            })
    
    # Custom Sections
    if req.custom_sections:
        insert_idx = len(slides) - 2
        for section in req.custom_sections:
            slides.insert(insert_idx, {
                "type": "custom",
                "title": section,
                "order": insert_idx
            })
            insert_idx += 1
    
    # Auf Plan-Größe begrenzen
    max_slides = min(plan.estimated_slides + 5, AGENT_CONFIG["max_slides"])
    if len(slides) > max_slides:
        # Entferne optionale Slides von hinten
        required_types = {"title", "executive_summary", "solution", "next_steps", "contact"}
        slides = [s for s in slides if s["type"] in required_types][:max_slides]
    
    return slides


# ============================================
# PHASE 4: DRAFT (Contextual Content Generation)
# ============================================

def phase_draft(
    req: AgentV3Request,
    structure: List[Dict[str, Any]],
    research: Dict[str, Any],
    iteration: int = 1
) -> List[Dict[str, Any]]:
    """
    Phase 4: Generiert Content für jeden Slide.
    Kontextbewusst: Jeder Slide kennt seine Nachbarn.
    """
    slides = []
    facts_context = "\n".join(research.get("facts", [])[:5])
    total_slides = len(structure)
    
    for idx, slide_def in enumerate(structure):
        slide_type = slide_def.get("type", "content")
        title = slide_def.get("title", "")
        
        # Kontext aufbauen (Agent Intelligence)
        prev_slide = structure[idx - 1] if idx > 0 else None
        next_slide = structure[idx + 1] if idx < total_slides - 1 else None
        
        context_info = f"""
Position: Slide {idx + 1} von {total_slides}
Vorheriger Slide: {prev_slide['title'] if prev_slide else 'Keiner (Start)'}
Nächster Slide: {next_slide['title'] if next_slide else 'Keiner (Ende)'}
Gesamtthema: {req.topic}
"""
        
        slide = {
            "type": slide_type,
            "title": title,
            "bullets": [],
            "notes": "",
            "layout_hint": _get_layout_hint(slide_type),
            "citations": [],
            "confidence": 0.8,
        }
        
        # NLG-Module nutzen (Quick Win 1)
        if req.use_nlg_modules and slide_type in NLG_MODULES:
            try:
                nlg_result = NLG_MODULES[slide_type](
                    {"customer_name": req.customer_name, "topic": req.topic, "brief": req.brief},
                    {}
                )
                if nlg_result:
                    slide["bullets"] = nlg_result.get("bullets", [])[:6]
                    slide["notes"] = nlg_result.get("notes", "")
                    slide["confidence"] = 0.9
            except Exception:
                pass
        
        # LLM-basierte Content-Generierung (falls NLG nicht genutzt/fehlgeschlagen)
        if not slide["bullets"] and HAS_LLM_CONTENT:
            try:
                # Kontextbewusster Prompt
                content_prompt = f"""Erstelle Inhalt für diesen Präsentations-Slide:

KONTEXT:
{context_info}

SLIDE-DETAILS:
Typ: {slide_type}
Titel: {title}

BRIEFING:
Thema: {req.topic}
Details: {req.brief}
Kunde: {req.customer_name}
Branche: {req.industry}

FAKTEN AUS RECHERCHE:
{facts_context[:500] if facts_context else 'Keine zusätzlichen Fakten'}

ITERATION: {iteration} (höher = mehr Fokus auf Qualität)

Erstelle:
1. 3-5 prägnante Bullet-Points (je max 15 Wörter)
2. Speaker Notes (2-4 Sätze zur Erläuterung)

Der Inhalt muss nahtlos zum vorherigen und nächsten Slide passen.

Output als JSON:
{{"bullets": ["..."], "notes": "..."}}

NUR JSON, keine Erklärung."""

                result = _llm_call(
                    content_prompt, 
                    task="quality" if iteration > 1 else "default",
                    model_override=req.llm_model,
                    max_tokens=600
                )
                
                # JSON parsen
                import re
                json_match = re.search(r'\{.*\}', result, re.DOTALL)
                if json_match:
                    data = json.loads(json_match.group())
                    slide["bullets"] = data.get("bullets", [])[:6]
                    slide["notes"] = data.get("notes", "")
                    slide["confidence"] = 0.85
                    
            except Exception as e:
                slide["notes"] = f"[Generierung: {str(e)[:50]}]"
        
        # Fallback
        if not slide["bullets"]:
            slide["bullets"] = [f"• Punkt zu {title}"]
            slide["confidence"] = 0.5
        
        # Citations aus Research
        if research.get("sources") and slide_type in ["executive_summary", "problem", "roi"]:
            slide["citations"] = [s["title"] for s in research["sources"][:2]]
        
        slides.append(slide)
    
    return slides


def _get_layout_hint(slide_type: str) -> str:
    """Gibt den passenden Layout-Hint für einen Slide-Typ."""
    layouts = {
        "title": "Title",
        "agenda": "Title and Content",
        "executive_summary": "Title and Content",
        "problem": "Title and Content",
        "solution": "Title and Content",
        "use_case": "Title and Content",
        "use_case_detail": "Two Content",
        "benefits": "Title and Content",
        "roi": "Title and Content",
        "roadmap": "Title and Content",
        "team": "Title and Content",
        "competitive": "Comparison",
        "risks": "Title and Content",
        "kpis": "Title and Content",
        "personas": "Title and Content",
        "next_steps": "Title and Content",
        "contact": "Title and Content",
    }
    return layouts.get(slide_type, "Title and Content")


# ============================================
# PHASE 5: CRITIQUE (Self-Evaluation)
# ============================================

def phase_critique(
    req: AgentV3Request,
    slides: List[Dict[str, Any]],
    plan: AgentPlan
) -> Dict[str, Any]:
    """
    Phase 5: Agent bewertet seinen eigenen Output.
    """
    result = {
        "score": 7.0,
        "issues": [],
        "suggestions": [],
        "should_revise": False,
    }
    
    # Basis-Bewertung
    total_bullets = sum(len(s.get("bullets", [])) for s in slides)
    avg_bullets = total_bullets / len(slides) if slides else 0
    avg_confidence = sum(s.get("confidence", 0.5) for s in slides) / len(slides) if slides else 0
    
    # Score berechnen
    score = 5.0
    
    # Bullet-Count Check
    if 3 <= avg_bullets <= 5:
        score += 1.0
    elif avg_bullets < 2:
        result["issues"].append("Zu wenig Inhalt pro Slide")
        result["suggestions"].append("Mehr Details hinzufügen")
    elif avg_bullets > 6:
        result["issues"].append("Zu viel Text pro Slide")
        result["suggestions"].append("Auf Kernaussagen fokussieren")
    
    # Confidence Check
    score += avg_confidence * 2  # Max +2
    
    # Struktur Check
    required_types = {"title", "executive_summary", "next_steps"}
    present_types = {s.get("type") for s in slides}
    missing = required_types - present_types
    if missing:
        result["issues"].append(f"Fehlende Slides: {missing}")
        score -= len(missing) * 0.5
    
    # LLM-basierte Kritik (wenn verfügbar und gewünscht)
    if req.include_critique and HAS_LLM_CONTENT:
        try:
            critique_prompt = f"""Bewerte diese Präsentation kritisch:

THEMA: {req.topic}
SLIDES: {len(slides)}
INHALT-ÜBERSICHT:
{chr(10).join([f"- {s.get('title')}: {len(s.get('bullets', []))} Bullets" for s in slides[:10]])}

Bewerte auf einer Skala von 1-10:
- Ist die Struktur logisch?
- Sind die Inhalte relevant?
- Fehlt etwas Wichtiges?
- Ist die Präsentation überzeugend?

Output als JSON:
{{"score": 7, "issues": ["..."], "suggestions": ["..."]}}

NUR JSON."""

            llm_result = _llm_call(critique_prompt, task="analysis", max_tokens=400)
            
            import re
            json_match = re.search(r'\{.*\}', llm_result, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group())
                llm_score = data.get("score", 7)
                # Kombiniere Scores
                score = (score + llm_score) / 2
                result["issues"].extend(data.get("issues", [])[:3])
                result["suggestions"].extend(data.get("suggestions", [])[:3])
                
        except Exception:
            pass
    
    # Feedback-Loop Integration (Quick Win 3)
    if HAS_FEEDBACK:
        try:
            fb_suggestions = get_improvement_suggestions(
                content={"slides": slides},
                slide_types=[s.get("type") for s in slides]
            )
            for s in fb_suggestions[:3]:
                if s.get("priority") == "high":
                    result["suggestions"].append(s.get("message", ""))
        except Exception:
            pass
    
    result["score"] = round(max(1, min(10, score)), 1)
    result["should_revise"] = result["score"] < AGENT_CONFIG["quality_threshold"] and req.auto_improve
    
    return result


# ============================================
# PHASE 6: REVISE (Iterative Improvement)
# ============================================

def phase_revise(
    req: AgentV3Request,
    slides: List[Dict[str, Any]],
    critique: Dict[str, Any],
    research: Dict[str, Any]
) -> List[Dict[str, Any]]:
    """
    Phase 6: Verbessert Slides basierend auf Kritik.
    """
    if not critique.get("should_revise"):
        return slides
    
    improved_slides = []
    
    for slide in slides:
        # Slides mit niedriger Confidence verbessern
        if slide.get("confidence", 1.0) < 0.7:
            if HAS_LLM_CONTENT:
                try:
                    improve_prompt = f"""Verbessere diesen Slide:

AKTUELL:
Titel: {slide.get('title')}
Bullets: {slide.get('bullets')}

KRITIK:
{chr(10).join(critique.get('issues', []))}

VERBESSERUNGSVORSCHLÄGE:
{chr(10).join(critique.get('suggestions', []))}

Erstelle verbesserte Bullets (3-5 Stück, prägnant).

Output als JSON:
{{"bullets": ["..."], "notes": "..."}}"""

                    result = _llm_call(improve_prompt, task="quality", max_tokens=400)
                    
                    import re
                    json_match = re.search(r'\{.*\}', result, re.DOTALL)
                    if json_match:
                        data = json.loads(json_match.group())
                        slide["bullets"] = data.get("bullets", slide["bullets"])
                        slide["notes"] = data.get("notes", slide["notes"])
                        slide["confidence"] = 0.85
                        slide["revised"] = True
                except Exception:
                    pass
        
        improved_slides.append(slide)
    
    return improved_slides


# ============================================
# PHASE 7: VISUALIZE (Charts & Assets)
# ============================================

def phase_visualize(
    req: AgentV3Request,
    slides: List[Dict[str, Any]],
    context: Dict[str, Any] = None
) -> List[Dict[str, Any]]:
    """
    Phase 7: Enhanced Visuals mit Visual Intelligence.
    Fügt Charts, Images und Layout-Optimierung hinzu.
    """
    if not req.generate_charts:
        return slides
    
    # === VISUAL INTELLIGENCE (Stufe 3) ===
    if HAS_VISUAL_INTELLIGENCE and enhance_all_slides:
        try:
            enhanced = enhance_all_slides(
                slides=slides,
                context=context or {
                    "topic": req.topic,
                    "industry": req.industry,
                    "customer_name": req.customer_name
                },
                generate_charts=True,
                recommend_images_flag=req.match_assets,
                use_llm=True
            )
            return enhanced
        except Exception as e:
            pass  # Fallback zu altem Code
    
    # === FALLBACK: Alter Code ===
    if not HAS_CHART_GEN:
        return slides
    
    for slide in slides:
        slide_type = slide.get("type", "")
        
        # Roadmap → Timeline
        if slide_type == "roadmap":
            try:
                result = create_timeline(
                    phases=[
                        {"name": "Phase 1", "duration": "2 Wo", "description": "Discovery"},
                        {"name": "Phase 2", "duration": "4 Wo", "description": "Pilot"},
                        {"name": "Phase 3", "duration": "8 Wo", "description": "Rollout"},
                        {"name": "Phase 4", "duration": "Ongoing", "description": "Optimierung"},
                    ],
                    title=""
                )
                if result.get("ok"):
                    slide["chart"] = result.get("path")
                    slide["has_chart"] = True
            except Exception:
                pass
        
        # ROI → Bar Chart
        elif slide_type == "roi":
            try:
                result = create_bar_chart(
                    labels=["Investition", "Jahr 1", "Jahr 2", "Jahr 3"],
                    values=[-100, 30, 80, 150],
                    title="ROI Projektion",
                    horizontal=True
                )
                if result.get("ok"):
                    slide["chart"] = result.get("path")
                    slide["has_chart"] = True
            except Exception:
                pass
        
        # KPIs → Gauge oder Bar
        elif slide_type == "kpis":
            try:
                result = create_bar_chart(
                    labels=["Conversion", "Retention", "NPS", "ROI"],
                    values=[15, 85, 72, 45],
                    title="Key Performance Indicators"
                )
                if result.get("ok"):
                    slide["chart"] = result.get("path")
                    slide["has_chart"] = True
            except Exception:
                pass
        
        # Funnel
        elif "funnel" in slide_type.lower() or "pipeline" in slide.get("title", "").lower():
            try:
                result = create_funnel_chart(
                    stages=["Awareness", "Interest", "Consideration", "Intent", "Purchase"],
                    values=[1000, 600, 400, 200, 80],
                    title="Sales Funnel"
                )
                if result.get("ok"):
                    slide["chart"] = result.get("path")
                    slide["has_chart"] = True
            except Exception:
                pass
        
        # Competitive → Matrix
        elif slide_type == "competitive":
            try:
                result = create_comparison_matrix(
                    items=["Wir", "Wettbewerber A", "Wettbewerber B"],
                    criteria=["Preis", "Qualität", "Service", "Innovation"],
                    scores=[
                        [8, 9, 9, 8],
                        [6, 7, 5, 6],
                        [7, 6, 7, 5]
                    ],
                    title="Marktvergleich"
                )
                if result.get("ok"):
                    slide["chart"] = result.get("path")
                    slide["has_chart"] = True
            except Exception:
                pass
    
    # Asset Matching
    if req.match_assets and HAS_ASSET_TAGGER:
        try:
            scan_uploads_directory(UPLOADS_DIR)
            slides = match_assets_to_slides(slides)
        except Exception:
            pass
    
    return slides


# ============================================
# PHASE 8: RENDER (PPTX Export)
# ============================================

def phase_render(
    req: AgentV3Request,
    slides: List[Dict[str, Any]],
    run_id: str
) -> Dict[str, Any]:
    """
    Phase 8: Rendert das Deck zu PPTX.
    """
    result = {
        "pptx_url": None,
        "pdf_url": None,
        "project_id": run_id,
    }
    
    if not req.export_pptx:
        return result
    
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
    except ImportError:
        return result
    
    prs = Presentation()
    
    for slide_data in slides:
        title = str(slide_data.get("title", "Slide")).strip() or "Slide"
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("notes", "")
        slide_type = slide_data.get("type", "content")
        
        # Layout wählen
        if slide_type == "title":
            layout = prs.slide_layouts[0]
        elif not bullets:
            layout = prs.slide_layouts[5]
        else:
            layout = prs.slide_layouts[1]
        
        slide = prs.slides.add_slide(layout)
        
        # Titel
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Bullets
        if bullets and len(slide.shapes.placeholders) > 1:
            try:
                body = slide.shapes.placeholders[1].text_frame
                body.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        body.text = str(bullet)
                    else:
                        p = body.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
            except Exception:
                pass
        
        # Notes
        if notes and slide.notes_slide:
            try:
                slide.notes_slide.notes_text_frame.text = str(notes)
            except Exception:
                pass
        
        # Chart
        chart_path = slide_data.get("chart")
        if chart_path and os.path.exists(chart_path):
            try:
                slide.shapes.add_picture(
                    chart_path,
                    Inches(5.5),
                    Inches(2.5),
                    width=Inches(4)
                )
            except Exception:
                pass
    
    # Speichern
    exports_dir = Path(EXPORTS_DIR)
    exports_dir.mkdir(parents=True, exist_ok=True)
    
    ts = int(time.time())
    safe_topic = "".join(c if c.isalnum() else "_" for c in req.topic[:30])
    out_name = f"v3_{safe_topic}_{ts}.pptx"
    out_path = exports_dir / out_name
    
    prs.save(str(out_path))
    
    result["pptx_url"] = f"/exports/download/{out_name}"
    result["pptx_path"] = str(out_path)
    
    return result


# ============================================
# MAIN ENDPOINT
# ============================================

@router.post("/run_v3", response_model=AgentV3Response)
async def run_v3(req: AgentV3Request, background_tasks: BackgroundTasks) -> AgentV3Response:
    """
    Agent V3.1: Intelligente Deck-Generierung mit Multi-Step Reasoning.
    
    Features:
    - Analyse & Planning Phase
    - Kontextbewusste Content-Generierung
    - Iterative Selbstverbesserung
    - NLG-Module Integration
    - Template Learning
    - Feedback-Loop Integration
    """
    run_id = _generate_run_id()
    t0 = time.time()
    warnings = []
    improvements = []
    phases = {}
    
    # === PHASE 1: ANALYZE ===
    t1 = time.time()
    plan = phase_analyze(req)
    phases["analyze"] = {
        "duration_ms": int((time.time() - t1) * 1000),
        "complexity": plan.complexity,
        "estimated_slides": plan.estimated_slides
    }
    
    # === PHASE 2: RESEARCH ===
    t2 = time.time()
    research = phase_research(req, plan)
    phases["research"] = {
        "duration_ms": int((time.time() - t2) * 1000),
        "sources_found": len(research.get("sources", [])),
        "facts_gathered": len(research.get("facts", [])),
        "enhanced": HAS_KNOWLEDGE_ENHANCED
    }
    
    # === PHASE 3: STRUCTURE ===
    t3 = time.time()
    structure = phase_structure(req, plan, research)
    phases["structure"] = {
        "duration_ms": int((time.time() - t3) * 1000),
        "planned_slides": len(structure)
    }
    
    # === ITERATIVE IMPROVEMENT LOOP ===
    slides = []
    iteration = 1
    final_critique = {"score": 0}
    
    for iteration in range(1, req.max_iterations + 1):
        # === PHASE 4: DRAFT ===
        t4 = time.time()
        slides = phase_draft(req, structure, research, iteration)
        phases[f"draft_{iteration}"] = {
            "duration_ms": int((time.time() - t4) * 1000),
            "generated_slides": len(slides)
        }
        
        # === PHASE 5: CRITIQUE ===
        t5 = time.time()
        critique = phase_critique(req, slides, plan)
        phases[f"critique_{iteration}"] = {
            "duration_ms": int((time.time() - t5) * 1000),
            "score": critique.get("score")
        }
        final_critique = critique
        
        # Qualität gut genug?
        if not critique.get("should_revise") or iteration >= req.max_iterations:
            break
        
        # === PHASE 6: REVISE ===
        t6 = time.time()
        slides = phase_revise(req, slides, critique, research)
        phases[f"revise_{iteration}"] = {
            "duration_ms": int((time.time() - t6) * 1000),
            "slides_revised": sum(1 for s in slides if s.get("revised"))
        }
        improvements.append(f"Iteration {iteration}: Score {critique.get('score')} → Überarbeitung")
    
    # === PHASE 7: VISUALIZE ===
    t7 = time.time()
    slides = phase_visualize(req, slides, context={
        "topic": req.topic,
        "industry": req.industry,
        "customer_name": req.customer_name,
        "brief": req.brief
    })
    charts_count = sum(1 for s in slides if s.get("has_chart"))
    phases["visualize"] = {
        "duration_ms": int((time.time() - t7) * 1000),
        "charts_generated": charts_count
    }
    
    # === PHASE 8: RENDER ===
    t8 = time.time()
    render_result = phase_render(req, slides, run_id)
    phases["render"] = {
        "duration_ms": int((time.time() - t8) * 1000),
        "exported": render_result.get("pptx_url") is not None
    }
    
    # Feedback speichern (Quick Win 3)
    if HAS_FEEDBACK:
        try:
            background_tasks.add_task(
                record_feedback,
                project_id=run_id,
                overall_score=int(final_critique.get("score", 7)),
                comments=f"Auto-generated, iterations={iteration}"
            )
        except Exception:
            pass
    
    # Warnings aus Critique
    warnings.extend(final_critique.get("issues", [])[:5])
    
    # Response
    duration = round(time.time() - t0, 2)
    
    return AgentV3Response(
        ok=True,
        run_id=run_id,
        project_id=render_result.get("project_id"),
        slides=slides,
        slide_count=len(slides),
        plan=plan.model_dump(),
        iterations=iteration,
        final_quality=final_critique.get("score", 7.0),
        pptx_url=render_result.get("pptx_url"),
        pdf_url=render_result.get("pdf_url"),
        quality_score=final_critique.get("score"),
        duration_s=duration,
        phases=phases,
        warnings=warnings,
        improvements=improvements
    )


# ============================================
# ZUSÄTZLICHE ENDPOINTS
# ============================================

@router.get("/v3/status")
def agent_v3_status():
    """Gibt den Status aller Services zurück."""
    
    # Knowledge Enhanced Status
    ke_status = {}
    if HAS_KNOWLEDGE_ENHANCED:
        try:
            ke_status = knowledge_status()
        except Exception:
            ke_status = {"ok": False}
    
    return {
        "version": "3.8",
        "features": {
            "agent_intelligence": True,
            "iterative_improvement": True,
            "nlg_modules": len(NLG_MODULES) > 0,
            "template_learning": HAS_TEMPLATE_LEARNER,
            "feedback_loop": HAS_FEEDBACK,
            "multi_model": True,
        },
        "services": {
            "llm_content": HAS_LLM_CONTENT,
            "nlg_modules_count": len(NLG_MODULES),
            "asset_tagger": HAS_ASSET_TAGGER,
            "chart_generator": HAS_CHART_GEN,
            "template_learner": HAS_TEMPLATE_LEARNER,
            "feedback_loop": HAS_FEEDBACK,
            "knowledge": HAS_KNOWLEDGE,
            "knowledge_enhanced": HAS_KNOWLEDGE_ENHANCED,
            "visual_intelligence": HAS_VISUAL_INTELLIGENCE,
            "learning_adaptation": HAS_LEARNING,
            "multimodal_export": HAS_MULTIMODAL_EXPORT,
            "killer_features": HAS_KILLER_FEATURES,
            "advanced_features": HAS_ADVANCED_FEATURES,
            "live_generator": HAS_ADVANCED_FEATURES,
            "orchestrator": HAS_ORCHESTRATOR,
        },
        "models": OLLAMA_MODELS,
        "config": AGENT_CONFIG,
        "ollama": check_ollama() if HAS_LLM_CONTENT else {"ok": False}
    }


@router.post("/v3/analyze")
async def analyze_briefing(req: AgentV3Request):
    """
    Nur Analyse-Phase: Gibt Plan zurück ohne Content zu generieren.
    Für Vorschau/Planung.
    """
    plan = phase_analyze(req)
    return {
        "ok": True,
        "plan": plan.model_dump(),
        "deck_size": req.deck_size.value,
    }


@router.post("/v3/preview")
async def preview_structure(req: AgentV3Request):
    """
    Gibt nur die geplante Struktur zurück (ohne Content-Generierung).
    """
    plan = phase_analyze(req)
    research = phase_research(req, plan)
    structure = phase_structure(req, plan, research)
    
    return {
        "ok": True,
        "plan": plan.model_dump(),
        "slides": structure,
        "slide_count": len(structure),
        "research_sources": len(research.get("sources", []))
    }


@router.get("/v3/models")
def get_available_models():
    """Gibt verfügbare Ollama-Modelle zurück."""
    result = {
        "configured": OLLAMA_MODELS,
        "available": []
    }
    
    if HAS_LLM_CONTENT:
        try:
            status = check_ollama()
            result["available"] = status.get("available_models", [])
        except Exception:
            pass
    
    return result


# Alias für Frontend-Kompatibilität
@router.get("/status")
def agent_status_alias():
    """Alias für /agent/v3/status - Frontend-Kompatibilität."""
    return agent_v3_status()
