# -*- coding: utf-8 -*-
from __future__ import annotations
from typing import Any, Dict, List, Tuple
from pathlib import Path
from pptx import Presentation

TEMPLATE_DIR = Path("data/templates")

def save_template(file_path: str|Path, name: str) -> Dict[str, Any]:
    TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)
    dst = TEMPLATE_DIR / f"{name}.pptx"
    Path(file_path).replace(dst)
    return {"ok": True, "name": name, "path": str(dst)}

def list_templates() -> List[Dict[str, Any]]:
    out = []
    for p in TEMPLATE_DIR.glob("*.pptx"):
        out.append({"name": p.stem, "path": str(p)})
    return sorted(out, key=lambda x: x["name"])

def inspect_template(name: str) -> Dict[str, Any]:
    p = TEMPLATE_DIR / f"{name}.pptx"
    if not p.exists():
        return {"ok": False, "error": "not_found"}
    prs = Presentation(str(p))
    layouts = []
    for i, layout in enumerate(prs.slide_layouts):
        phs = []
        for shp in layout.shapes:
            if not getattr(shp, "has_text_frame", False):
                # Placeholder-Infos, so gut es geht
                title = getattr(shp, "name", "")
                phs.append({"name": title, "text_like": False})
            else:
                phs.append({"name": getattr(shp, "name", ""), "text_like": True})
        layouts.append({"index": i, "name": layout.name, "placeholders": phs})
    return {"ok": True, "name": name, "layouts": layouts}
