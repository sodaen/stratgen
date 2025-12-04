
from __future__ import annotations
import qdrant_client


def _encode_vectors(emb, texts):
    """
    Normalisiert Zugriff auf den Embedder:
    - Wenn 'emb' callable ist: emb(texts)
    - Sonst: emb.encode(texts)
    Gibt immer eine Python-Liste zurück.
    """
    out = emb(texts) if callable(emb) else emb.encode(texts)
    try:
        return out.tolist()
    except AttributeError:
        # falls schon Liste/Numpy-Array ohne .tolist()
        try:
            return list(out)
        except TypeError:
            # letzte Fallback-Konvertierung
            return [*out]


import re, uuid, httpx
from pathlib import Path
from typing import Any
from services.rag_pipeline import get_qdrant, get_embedder, COLL
from qdrant_client.models import PointStruct

def _ensure_collection(qdr, emb, coll):
    """
    Stellt sicher, dass die Qdrant-Collection existiert.
    Ermittelt die Embedding-Dimension dynamisch (callable vs. .encode()).
    """
    try:
        qdr.get_collection(coll)
        return
    except Exception:
        pass
    # Dimension bestimmen
    vec = (emb(["ping"])[0] if callable(emb) else emb.encode(["ping"])[0])
    size = len(vec)
    from qdrant_client.http.models import Distance, VectorParams
    qdr.create_collection(
        collection_name=coll,
        vectors_config=VectorParams(size=size, distance=Distance.COSINE),
    )

def _read_file(p: str) -> str:
    """Liest verschiedene Dateiformate und extrahiert Text."""
    path = Path(p)
    ext = path.suffix.lower()
    
    # Text-Dateien
    if ext in ('.txt', '.md', '.csv', '.json', '.xml', '.html'):
        return path.read_text(encoding="utf-8", errors="ignore")
    
    # PDF
    if ext == '.pdf':
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(path))
            text = "\n".join([page.get_text() for page in doc])
            doc.close()
            return text
        except ImportError:
            raise ValueError("PyMuPDF nicht installiert - pip install pymupdf")
        except Exception as e:
            raise ValueError(f"PDF-Fehler: {e}")
    
    # Word DOCX
    if ext == '.docx':
        try:
            from docx import Document
            doc = Document(str(path))
            return "\n".join([p.text for p in doc.paragraphs])
        except ImportError:
            raise ValueError("python-docx nicht installiert")
        except Exception as e:
            raise ValueError(f"DOCX-Fehler: {e}")
    
    # Word DOC (alt)
    if ext == '.doc':
        try:
            import subprocess
            result = subprocess.run(['antiword', str(path)], capture_output=True, text=True)
            if result.returncode == 0:
                return result.stdout
            raise ValueError("antiword fehlgeschlagen")
        except FileNotFoundError:
            raise ValueError("antiword nicht installiert - sudo apt install antiword")
    
    # PowerPoint PPTX
    if ext == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            raise ValueError("python-pptx nicht installiert")
        except Exception as e:
            raise ValueError(f"PPTX-Fehler: {e}")
    
    # Excel XLSX
    if ext in ('.xlsx', '.xls'):
        try:
            import pandas as pd
            df = pd.read_excel(str(path), sheet_name=None)
            texts = []
            for sheet_name, sheet_df in df.items():
                texts.append(f"=== {sheet_name} ===")
                texts.append(sheet_df.to_string())
            return "\n".join(texts)
        except ImportError:
            raise ValueError("pandas/openpyxl nicht installiert")
        except Exception as e:
            raise ValueError(f"Excel-Fehler: {e}")
    
    # Fallback: versuche als Text
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except:
        raise ValueError(f"Nicht unterstütztes Format: {ext}")

def _fetch_url(url: str, timeout: float = 10.0) -> str:
    with httpx.Client(timeout=timeout, follow_redirects=True) as c:
        r = c.get(url)
        r.raise_for_status()
        txt = r.text
    # sehr einfache HTML→Text-Extraktion
    txt = re.sub(r"(?is)<(script|style).*?</\\1>", " ", txt)
    txt = re.sub(r"(?is)<[^>]+>", " ", txt)
    txt = re.sub(r"[\\t\\r\\f]+", " ", txt)
    txt = re.sub(r"\\s{2,}", " ", txt)
    return txt.strip()

def _chunk_text(text: str, chunk_size: int = 900, overlap: int = 100) -> list[str]:
    words = text.split()
    if not words:
        return []
    chunks = []
    i = 0
    while i < len(words):
        chunk = words[i:i+chunk_size]
        chunks.append(" ".join(chunk))
        i += chunk_size - overlap if chunk_size > overlap else chunk_size
    return chunks

def _upsert(customer: str, source: dict[str, Any], texts: list[str]) -> int:
    if not texts:
        return 0
    emb = get_embedder()
    qdr = get_qdrant()
    vecs = _encode_vectors(emb, texts)
    points = []
    for v, t in zip(vecs, texts):
        pid = str(uuid.uuid4())
        payload = {
            # Kompatibilität + neue Felder fürs Filtern:
            "customer": customer,
            "customer_name": customer,
            "title": source.get("title"),
            # bestehende Meta
            "source_id": source.get("id"),
            "source_type": source.get("type"),
            "tokens": source.get("tokens", []),
            "topics": source.get("topics", []),
            "subtopics": source.get("subtopics", []),
            # Inhalt
            "text": t,
        }
        points.append(PointStruct(id=pid, vector=v, payload=payload))
    qdr.upsert(collection_name=COLL, points=points)
    return len(points)

def ingest_entry(customer: str, entry: dict[str, Any]) -> dict[str, Any]:
    typ = entry.get("type", "file")
    if typ == "file":
        text = _read_file(entry.get("path", ""))
    elif typ == "web":
        text = _fetch_url(entry.get("url", ""))
    else:
        return {"ok": False, "error": f"unsupported type: {typ}"}
    chunks = _chunk_text(text)
    n = _upsert(customer, entry, chunks)
    return {"ok": True, "count": n}

def ingest_entries(customer: str, entries: list[dict[str, Any]]) -> dict[str, Any]:
    total = 0
    details = []
    for e in entries:
        res = ingest_entry(customer, e)
        details.append({"id": e.get("id"), **res})
        if res.get("ok"):
            total += res.get("count", 0)
    return {"ok": True, "total": total, "details": details}

# === [DQ EXTENSIONS: dedupe + auto-tagging] ===================================
# Leichtgewichtige Normalisierung, Hash, Duplikat-Check & Auto-Tagging.
# Nicht-invasiv: wir wrappen nur ingest_entry / _upsert, ohne bestehende Logik zu löschen.

import re as _dq_re
import hashlib as _dq_hashlib
from typing import List as _dq_List

try:
    from services.rag_pipeline import get_qdrant, COLL
except Exception:
    # Fallback: späte Bindung in Wrappern
    get_qdrant = None
    COLL = "stratgen_docs"

# --- Text-Normalisierung & Hash ------------------------------------------------
def _dq_normalize_text(txt: str) -> str:
    if not isinstance(txt, str): 
        return ""
    t = txt.strip().lower()
    t = _dq_re.sub(r'\s+', ' ', t)
    return t

def _dq_doc_hash(customer: str, txt: str) -> str:
    base = (customer or "") + "|" + _dq_normalize_text(txt or "")
    return _dq_hashlib.sha1(base.encode("utf-8")).hexdigest()

# --- Sehr einfache Keyword-Extraktion (de/en Stopwörter, Häufigkeit) -----------
_STOP = set("""
a an and are as at be by for from has have if in into is it its of on or
that the to up was were will with und der die das ist im den des auf für mit
ein eine einem einen einer als bei aus dem zu vom zur nicht noch schon sehr
""".split())

def _dq_auto_tags(txt: str, max_k: int = 8) -> _dq_List[str]:
    if not isinstance(txt, str) or not txt.strip():
        return []
    words = _dq_re.findall(r"[a-zA-ZäöüÄÖÜß\-]{3,}", txt.lower())
    freq = {}
    for w in words:
        if w in _STOP: 
            continue
        freq[w] = freq.get(w, 0) + 1
    # Top-N nach Häufigkeit (stabil)
    tags = sorted(freq.items(), key=lambda kv: (-kv[1], kv[0]))[:max_k]
    return [t for t,_ in tags]

# --- Bestehende Symbole sichern (Core-Refs) ------------------------------------
_ingest_entry_core = None
_upsert_core = None

# Wir suchen nach bestehenden Funktions-Definitionen, um sie später aufzurufen.
# (Der Host-Code ruft weiterhin ingest_entry/ _upsert auf – nur jetzt mit DQ)
try:
    # Namen existieren zur Laufzeit im Modul-Namespace
    _ingest_entry_core = ingest_entry  # type: ignore
except Exception:
    pass

try:
    _upsert_core = _upsert  # type: ignore
except Exception:
    pass

# --- Wrapper _upsert: Payload anreichern (doc_hash, deleted, tokens) -----------
def _dq_upsert_wrapper(customer, entry, chunks):
    # Falls der Originalcode entry-Keys in Payload übernimmt (wie tokens/topics/...),
    # dann reicht es, entry zu erweitern.
    doc_hash = entry.get("doc_hash")
    if not doc_hash:
        # wenn Text vorhanden, Hash nachziehen
        txt = entry.get("text") or (chunks[0] if chunks else "")
        if isinstance(txt, str) and txt:
            entry = dict(entry)
            entry["doc_hash"] = _dq_doc_hash(customer, txt)
    # deleted-Flag standardmäßig false
    if "deleted" not in entry:
        entry = dict(entry)
        entry["deleted"] = False
    # sicherstellen, dass tokens eine Liste ist
    if "tokens" in entry and not isinstance(entry["tokens"], list):
        entry = dict(entry)
        entry["tokens"] = list(entry["tokens"])

    if _upsert_core is None:
        raise RuntimeError("DQ wrapper: _upsert_core nicht gefunden – bitte prüfen, ob _upsert existiert.")
    return _upsert_core(customer, entry, chunks)

# --- Qdrant-Duplikat-Check -----------------------------------------------------
def _dq_exists_hash(customer: str, h: str) -> bool:
    if not h:
        return False
    try:
        qdr = get_qdrant() if callable(get_qdrant) else None
        if qdr is None:
            return False
        from qdrant_client import models as qm
        flt = qdrant_client.models.Filter(
            must=[
                qdrant_client.models.FieldCondition(key="customer", match=qdrant_client.models.MatchValue(value=customer)),
                qdrant_client.models.FieldCondition(key="doc_hash", match=qdrant_client.models.MatchValue(value=h)),
                qdrant_client.models.FieldCondition(key="deleted", match=qdrant_client.models.MatchValue(value=False)),
            ]
        )
        r = qdr.scroll(collection_name=COLL, scroll_filter=flt, limit=1)
        pts, _ = r
        return bool(pts)
    except Exception:
        # Fallback: im Zweifel kein Skip (fail-open)
        return False

# --- Wrapper ingest_entry: Dedupe + Auto-Tags + Upsert-Wrapper -----------------
def ingest_entry(customer, entry, *args, **kwargs):  # type: ignore
    """
    Dedupe & Auto-Tagging:
      - doc_hash = sha1(customer|normalized(text))
      - Duplikate (gleicher Hash + customer + deleted=False) werden übersprungen
      - tokens werden um Auto-Tags ergänzt
      - _upsert wird durch DQ-Wrapper aufgerufen (Payload-Anreicherung)
    """
    global _ingest_entry_core

    if _ingest_entry_core is None:
        raise RuntimeError("DQ wrapper: _ingest_entry_core nicht gefunden – ursprüngliche ingest_entry existiert nicht?")

    e = dict(entry or {})
    txt = e.get("text") if isinstance(e.get("text"), str) else ""

    # Auto-Tags
    auto = _dq_auto_tags(txt, max_k=8)
    if auto:
        toks = list(e.get("tokens") or [])
        for t in auto:
            if t not in toks:
                toks.append(t)
        e["tokens"] = toks

    # Hash
    h = _dq_doc_hash(customer, txt) if txt else None
    if h:
        e["doc_hash"] = h

    # Duplikate vermeiden
    if h and _dq_exists_hash(customer, h):
        # Konsistentes Return-Format: wie ein ingest, aber "skipped"
        return {"ok": True, "skipped": "duplicate", "hash": h, "points": 0}

    # _upsert hooken: Wir ersetzen lokal die Referenz, rufen Original ingest_entry auf,
    # und setzen danach wieder zurück.
    global _upsert
    if _upsert_core is None and '_upsert' in globals():
        # capture original now
        globals()['_upsert_core'] = globals()['_upsert']

    if '_upsert' in globals():
        original_upsert = globals()['_upsert']
        globals()['_upsert'] = _dq_upsert_wrapper
        try:
            return _ingest_entry_core(customer, e, *args, **kwargs)
        finally:
            globals()['_upsert'] = original_upsert
    else:
        # Falls kein _upsert existiert, rufen wir direkt den Core (best effort)
        return _ingest_entry_core(customer, e, *args, **kwargs)

# === [/DQ EXTENSIONS] ==========================================================

def ingest_entries_with_report(customer_name: str, entries):
    """
    Wrapper um ingest_entries(), der sicher immer eine Report-Struktur mitgibt.
    Struktur:
    {
      "ok": bool,
      "items": [... optional ...],
      "ids": [<uuid|str>],
      "message": "... optional ..."
    }
    """
    try:
        res = ingest_entries(customer_name, entries)
    except Exception as e:
        return {"ok": False, "error": {"code": "INGEST_FAILED", "message": str(e)}}

    # Bestehende Antworten übernehmen und absichern
    if isinstance(res, dict):
        out = dict(res)
        # Fallbacks: ids & items
        out.setdefault("ids", [])
        out.setdefault("items", [])
        out.setdefault("ok", True)
        return out

    # unbekannter Rückgabetyp -> minimaler Erfolg ohne IDs
    return {"ok": True, "ids": [], "items": []}

def query_semantic(customer_name: str, query: str, limit: int = 5, filters: dict | None = None):
    """
    Semantische Suche in Qdrant:
    - filtert immer auf payload["customer"] == customer_name
    - optionale Filter: tokens/topics/subtopics (MatchAny), source_type (MatchValue),
      title_contains (MatchText, falls konfiguriert; ansonsten wird kein Titeltext-Match erzwungen)
    """
    qdr = get_qdrant()
    emb = get_embedder()

    # Vektor für Query
    vec = _encode_vectors(emb, [query])[0]

    # Qdrant-Filter zusammenbauen
    must = [qdrant_client.models.FieldCondition(key="customer", match=qdrant_client.models.MatchValue(value=customer_name))]
    # Tippfehler fix (Leerzeichen entfernen)
    must[-1] = qdrant_client.models.FieldCondition(key="customer", match=qdrant_client.models.MatchValue(value=customer_name))

    filters = filters or {}

    def add_match_any(key: str):
        vals = filters.get(key)
        if vals:
            must.append(qdrant_client.models.FieldCondition(key=key, match=qdrant_client.models.MatchAny(any=vals)))

    add_match_any("tokens")
    add_match_any("topics")
    add_match_any("subtopics")

    # source_type == "text"/"file"/"web"/...
    if filters.get("source_type"):
        must.append(qdrant_client.models.FieldCondition(key="source_type", match=qdrant_client.models.MatchValue(value=filters["source_type"])))

    # "title contains ..." (nur wenn es den payload-key "title" gibt; MatchText erfordert Text-Index)
    if filters.get("title_contains"):
        must.append(qdrant_client.models.FieldCondition(key="title", match=qdrant_client.models.MatchText(text=filters["title_contains"])))

    qfilter = qdrant_client.models.Filter(must=must) if must else None

    # Suche absetzen
    hits = qdr.search(
        collection_name=COLL,
        query_vector=vec,
        limit=limit,
        query_filter=qfilter,
        with_payload=True,
        with_vectors=False,
    )

    items = []
    for h in hits or []:
        payload = h.payload or {}
        items.append({
            "id": str(h.id),
            "score": float(h.score) if getattr(h, "score", None) is not None else None,
            "payload": payload,
            "text": payload.get("text"),
            "title": payload.get("title"),
            "customer_name": payload.get("customer"),
        })

    return {"ok": True, "items": items, "next_offset": None}
