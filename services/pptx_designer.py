"""
PPTX Designer Service für Stratgen.
Erstellt professionelle PowerPoint-Präsentationen mit echtem Design.
"""

import os
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import json

logger = logging.getLogger(__name__)

DATA_ROOT = Path(os.getenv("STRATGEN_DATA", "/home/sodaen/stratgen/data"))


class PPTXDesigner:
    """Erstellt professionelle PPTX-Präsentationen."""
    
    # Farbpaletten
    PALETTES = {
        "corporate": {
            "primary": "1E40AF",
            "secondary": "3B82F6", 
            "accent": "10B981",
            "dark": "1F2937",
            "light": "F3F4F6",
            "text": "111827"
        },
        "modern": {
            "primary": "7C3AED",
            "secondary": "A78BFA",
            "accent": "F59E0B",
            "dark": "18181B",
            "light": "FAFAFA",
            "text": "27272A"
        },
        "minimal": {
            "primary": "000000",
            "secondary": "525252",
            "accent": "3B82F6",
            "dark": "171717",
            "light": "FFFFFF",
            "text": "262626"
        },
        "vibrant": {
            "primary": "EC4899",
            "secondary": "8B5CF6",
            "accent": "06B6D4",
            "dark": "1E1B4B",
            "light": "FDF4FF",
            "text": "581C87"
        }
    }
    
    # Slide Layouts
    LAYOUTS = {
        "title": {
            "title_top": 2.5,
            "title_size": 44,
            "subtitle_size": 24
        },
        "content": {
            "title_top": 0.5,
            "title_size": 32,
            "content_top": 1.5,
            "content_size": 18
        },
        "two_column": {
            "title_top": 0.5,
            "title_size": 32,
            "col_width": 4.5
        },
        "image_left": {
            "image_width": 5,
            "content_width": 4.5
        },
        "full_image": {
            "overlay": True
        }
    }
    
    def __init__(self, palette: str = "corporate"):
        self.palette_name = palette
        self.palette = self.PALETTES.get(palette, self.PALETTES["corporate"])
        self.prs = None
    
    def _hex_to_rgb(self, hex_color: str) -> RgbColor:
        """Konvertiert Hex zu RGB."""
        hex_color = hex_color.lstrip('#')
        return RgbColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def create_presentation(self, slides: List[Dict], 
                           title: str = "Präsentation",
                           palette: str = None) -> bytes:
        """
        Erstellt eine komplette PPTX-Präsentation.
        
        Args:
            slides: Liste von Slide-Dicts mit type, title, content, etc.
            title: Titel der Präsentation
            palette: Farbpalette (corporate, modern, minimal, vibrant)
        
        Returns:
            PPTX als Bytes
        """
        if palette:
            self.palette = self.PALETTES.get(palette, self.PALETTES["corporate"])
        
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)
        
        for i, slide_data in enumerate(slides):
            slide_type = slide_data.get("type", "content")
            
            if slide_type == "title" or i == 0:
                self._add_title_slide(slide_data)
            elif slide_type == "section":
                self._add_section_slide(slide_data)
            elif slide_type == "two_column":
                self._add_two_column_slide(slide_data)
            elif slide_type == "bullets":
                self._add_bullet_slide(slide_data)
            elif slide_type == "quote":
                self._add_quote_slide(slide_data)
            else:
                self._add_content_slide(slide_data)
        
        # Speichere als Bytes
        from io import BytesIO
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.read()
    
    def _add_background(self, slide, color_key: str = "light"):
        """Fügt Hintergrundfarbe hinzu."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(self.palette[color_key])
    
    def _add_title_slide(self, data: Dict):
        """Erstellt Title-Slide."""
        layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(layout)
        
        # Hintergrund
        self._add_background(slide, "primary")
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.palette["light"])
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if data.get("subtitle"):
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(1)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = data.get("subtitle", "")
            p.font.size = Pt(24)
            p.font.color.rgb = self._hex_to_rgb(self.palette["light"])
            p.alignment = PP_ALIGN.CENTER
        
        # Accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.5), Inches(4), Inches(2.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.palette["accent"])
        line.line.fill.background()
    
    def _add_section_slide(self, data: Dict):
        """Erstellt Section-Divider-Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide, "secondary")
        
        # Section Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.palette["light"])
        p.alignment = PP_ALIGN.CENTER
    
    def _add_content_slide(self, data: Dict):
        """Erstellt Content-Slide mit Titel und Text."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide, "light")
        
        # Accent bar oben
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(self.palette["primary"])
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.palette["text"])
        
        # Content
        content = data.get("content", "") or data.get("body", "")
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            # Split by newlines or bullets
            lines = content.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                line = line.strip()
                if line.startswith('- ') or line.startswith('• '):
                    p.text = "• " + line[2:]
                    p.level = 0
                else:
                    p.text = line
                
                p.font.size = Pt(18)
                p.font.color.rgb = self._hex_to_rgb(self.palette["text"])
                p.space_after = Pt(12)
    
    def _add_bullet_slide(self, data: Dict):
        """Erstellt Bullet-Point-Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide, "light")
        
        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.15), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(self.palette["primary"])
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.palette["text"])
        
        # Bullets
        bullets = data.get("bullets", [])
        if not bullets and data.get("content"):
            bullets = [l.strip().lstrip('- •') for l in data["content"].split('\n') if l.strip()]
        
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5)
            )
            tf = content_box.text_frame
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = "• " + bullet
                p.font.size = Pt(20)
                p.font.color.rgb = self._hex_to_rgb(self.palette["text"])
                p.space_after = Pt(16)
    
    def _add_two_column_slide(self, data: Dict):
        """Erstellt Two-Column-Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide, "light")
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.palette["text"])
        
        # Left Column
        left_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(5.5), Inches(5.5)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("left_content", data.get("content", "")[:len(data.get("content", ""))//2])
        p.font.size = Pt(16)
        p.font.color.rgb = self._hex_to_rgb(self.palette["text"])
        
        # Right Column
        right_box = slide.shapes.add_textbox(
            Inches(7.083), Inches(1.5), Inches(5.5), Inches(5.5)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("right_content", "")
        p.font.size = Pt(16)
        p.font.color.rgb = self._hex_to_rgb(self.palette["text"])
        
        # Divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.583), Inches(1.5), Inches(0.02), Inches(5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.palette["secondary"])
        line.line.fill.background()
    
    def _add_quote_slide(self, data: Dict):
        """Erstellt Quote-Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide, "dark")
        
        # Quote marks
        quote_mark = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(2), Inches(1.5)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = "\""
        p.font.size = Pt(120)
        p.font.color.rgb = self._hex_to_rgb(self.palette["accent"])
        
        # Quote text
        quote_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(3)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("quote", data.get("content", ""))
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = self._hex_to_rgb(self.palette["light"])
        p.alignment = PP_ALIGN.CENTER
        
        # Attribution
        if data.get("author"):
            attr_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5), Inches(11.333), Inches(0.5)
            )
            tf = attr_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {data['author']}"
            p.font.size = Pt(18)
            p.font.color.rgb = self._hex_to_rgb(self.palette["secondary"])
            p.alignment = PP_ALIGN.CENTER


# Singleton
_designer = None

def get_designer(palette: str = "corporate") -> PPTXDesigner:
    global _designer
    if _designer is None or _designer.palette_name != palette:
        _designer = PPTXDesigner(palette)
    return _designer


def create_presentation(slides: List[Dict], title: str = "Präsentation", 
                       palette: str = "corporate") -> bytes:
    """Erstellt PPTX aus Slides."""
    designer = PPTXDesigner(palette)
    return designer.create_presentation(slides, title, palette)
