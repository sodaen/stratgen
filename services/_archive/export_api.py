# backend/data/data/data/data/exportssss_api.py
from __future__ import annotations
from io import BytesIO
from typing import Optional, Any, Dict, List

from fastapi import APIRouter
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel

router = APIRouter(prefix="/data/data/data/data/exportssss", tags=["export"])

class DeckExportReq(BaseModel):
    customer_name: str
    topic: str
    outline: Dict[str, Any]          # exakt das JSON von /outline/generate
    template_id: Optional[str] = None
    style: Optional[Dict[str, Any]] = None
    filename: Optional[str] = None   # optionaler Dateiname (ohne Pfad)

def _iter_outline_sections(outline: Dict[str, Any]) -> List[Dict[str, Any]]:
    sections = outline.get("sections") or []
    # Defensive: Nur Dicts mit title/bullets
    out = []
    for s in sections:
        if isinstance(s, dict):
            title = s.get("title") or "Section"
            bullets = s.get("bullets") or []
            out.append({"title": title, "bullets": [str(b) for b in bullets]})
    return out

def _render_pptx_fallback(req: DeckExportReq) -> bytes:
    """
    Minimal-Renderer mit python-pptx (Fallback), falls du keinen eigenen Renderer nutzen willst.
    Anforderungen:
      pip install python-pptx
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()

        # Titelfolie
        slide_layout = prs.slide_layouts[0]  # Title + Subtitle
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = req.topic
        subtitle = slide.placeholders[1]
        subtitle.text = req.customer_name

        # Agenda
        sections = _iter_outline_sections(req.outline)
        if sections:
            layout = prs.slide_layouts[1]  # Title + Content
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = "Agenda"
            tf = slide.shapes.placeholders[1].text_frame
            tf.clear()
            for s in sections:
                p = tf.add_paragraph()
                p.text = s["title"]
                p.level = 0

        # Pro Section eine Folie
        for s in sections:
            layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = s["title"]
            tf = slide.shapes.placeholders[1].text_frame
            tf.clear()
            # erste Bullet direkt auf text_frame setzen (pptx hat schon einen Absatz)
            first = True
            for b in s["bullets"][:6]:
                if first:
                    tf.text = b
                    first = False
                else:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0

        bio = BytesIO()
        prs.save(bio)
        return bio.getvalue()
    except Exception as e:
        # Sauberer Fehler zurück (z. B. wenn python-pptx nicht installiert ist)
        raise RuntimeError(f"Fallback-PPTX-Render fehlgeschlagen: {e}")

@router.post("/deck")
def export_deck(req: DeckExportReq):
    """
    V0: Nimmt Outline-JSON und gibt eine PPTX-Datei zurück.
    Falls du services.ppt_renderer.generate_pptx_from_template verwenden willst,
    kannst du das hier anschließen. Sonst nutzt der Fallback python-pptx.
    """
    # 1) Versuch: Eigenen Renderer nutzen (falls vorhanden)
    try:
        from services.ppt_renderer import generate_pptx_from_template  # type: ignore
        ppt_bytes = generate_pptx_from_template(
            outline=req.outline,
            topic=req.topic,
            customer=req.customer_name,
            template_id=req.template_id,
            style=req.style or {},
        )
    except Exception:
        # 2) Fallback: sehr einfache PPT aus Outline bauen
        ppt_bytes = _render_pptx_fallback(req)

    fname = req.filename or "deck.pptx"
    return StreamingResponse(
        content=BytesIO(ppt_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{fname}"'}
    )

