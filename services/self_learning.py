# -*- coding: utf-8 -*-
"""
STRATGEN Self-Learning System (Phase 8)
Lernt aus erfolgreichen Exports und User-Feedback.

FIX: get_embedding auf Modul-Level (war fälschlicherweise
     mitten in Methoden-Bodies definiert → IndentationError)
"""

from pathlib import Path
from datetime import datetime, timedelta
from typing import List, Dict, Optional
import json
import logging
import uuid
import threading
import time

log = logging.getLogger("stratgen.self_learning")

# Konfiguration
GENERATED_OUTPUTS_DIR = Path("/home/sodaen/stratgen/data/generated_outputs")
SESSION_COLLECTIONS_DIR = Path("/home/sodaen/stratgen/data/session_collections")
GENERATED_OUTPUTS_DIR.mkdir(parents=True, exist_ok=True)
SESSION_COLLECTIONS_DIR.mkdir(parents=True, exist_ok=True)

QUALITY_THRESHOLD = 0.7   # Minimum für Indexierung
SESSION_TTL_HOURS  = 24   # Session-Collections Lebenszeit
QDRANT_COLLECTION  = "generated_outputs"
EMBEDDING_MODEL    = "nomic-embed-text"


# ─────────────────────────────────────────────
# MODUL-LEVEL HILFSFUNKTIONEN
# ─────────────────────────────────────────────

def get_embedding(text: str) -> list:
    """Holt Embedding von Ollama (nomic-embed-text). Gibt [] zurück bei Fehler."""
    import os
    try:
        import httpx
        host = os.getenv("OLLAMA_HOST", "http://localhost:11434").rstrip("/")
        response = httpx.post(
            f"{host}/api/embeddings",
            json={"model": EMBEDDING_MODEL, "prompt": text},
            timeout=30.0,
        )
        if response.status_code == 200:
            return response.json().get("embedding", [])
    except Exception as e:
        log.debug("Embedding error: %s", e)
    return []


def _get_qdrant_client():
    """Gibt einen QdrantClient zurück oder None."""
    try:
        from qdrant_client import QdrantClient
        return QdrantClient(host="localhost", port=6333)
    except Exception as e:
        log.debug("Qdrant not available: %s", e)
        return None


def _ensure_collection(client, collection_name: str, vector_size: int = 768):
    """Stellt sicher dass eine Qdrant-Collection existiert."""
    try:
        from qdrant_client.models import VectorParams, Distance
        collections = [c.name for c in client.get_collections().collections]
        if collection_name not in collections:
            client.create_collection(
                collection_name=collection_name,
                vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
            )
            log.info("Created Qdrant collection: %s", collection_name)
    except Exception as e:
        log.warning("Could not ensure collection %s: %s", collection_name, e)


# ─────────────────────────────────────────────
# SELF-LEARNING SYSTEM
# ─────────────────────────────────────────────

class SelfLearningSystem:
    """Lernt aus erfolgreichen Exports und Feedback."""

    def __init__(self):
        self.feedback_scores: Dict[str, List[int]] = {}
        self._cleanup_thread = None
        self._start_cleanup_thread()

    def _start_cleanup_thread(self):
        def cleanup_loop():
            while True:
                try:
                    self._cleanup_old_sessions()
                except Exception as e:
                    log.error("Session cleanup error: %s", e)
                time.sleep(3600)

        self._cleanup_thread = threading.Thread(target=cleanup_loop, daemon=True)
        self._cleanup_thread.start()

    # ─────────────────────────────────────────
    # 8.1 EXPORT FEEDBACK LOOP
    # ─────────────────────────────────────────

    def on_export_complete(
        self,
        export_path: str,
        session_id: str,
        export_type: str = "pptx",
    ) -> Dict:
        """
        Wird aufgerufen wenn ein Export abgeschlossen ist.
        Analysiert Inhalt und indexiert bei guter Qualität in Qdrant.
        """
        result = {
            "session_id": session_id,
            "export_path": export_path,
            "export_type": export_type,
            "quality_score": 0.0,
            "indexed": False,
            "chunks_created": 0,
        }

        try:
            content_chunks = self._extract_export_content(export_path, export_type)

            if not content_chunks:
                result["error"] = "Keine Inhalte extrahiert"
                self._save_export_metadata(result)
                return result

            quality_score = self._assess_quality(content_chunks)
            result["quality_score"] = quality_score

            if quality_score >= QUALITY_THRESHOLD:
                client = _get_qdrant_client()
                if client:
                    _ensure_collection(client, QDRANT_COLLECTION)
                    try:
                        from qdrant_client.models import PointStruct
                        points = []
                        for i, chunk in enumerate(content_chunks):
                            embedding = get_embedding(chunk["text"])
                            if not embedding:
                                continue
                            points.append(PointStruct(
                                id=str(uuid.uuid4()),
                                vector=embedding,
                                payload={
                                    "text": chunk["text"],
                                    "source": f"export_{session_id}",
                                    "source_file": Path(export_path).name,
                                    "export_type": export_type,
                                    "quality_score": chunk.get("quality", quality_score),
                                    "session_id": session_id,
                                    "indexed_at": datetime.now().isoformat(),
                                    "chunk_index": i,
                                    "is_generated": True,
                                },
                            ))

                        if points:
                            client.upsert(
                                collection_name=QDRANT_COLLECTION,
                                points=points,
                            )
                            result["indexed"] = True
                            result["chunks_created"] = len(points)
                            log.info(
                                "Self-learning: indexed %d chunks from session %s (quality=%.2f)",
                                len(points), session_id, quality_score,
                            )
                    except Exception as e:
                        log.warning("Qdrant upsert failed: %s", e)
                        result["error"] = f"Qdrant: {e}"
                else:
                    result["error"] = "Qdrant nicht verfügbar"
            else:
                log.info(
                    "Self-learning: skipping indexing for session %s (quality=%.2f < %.2f)",
                    session_id, quality_score, QUALITY_THRESHOLD,
                )
                result["skipped_reason"] = "quality_below_threshold"

            self._log_learning_event(session_id, quality_score, result["chunks_created"])

        except Exception as e:
            log.error("on_export_complete error: %s", e)
            result["error"] = str(e)
        finally:
            self._save_export_metadata(result)

        return result

    def _extract_export_content(self, export_path: str, export_type: str) -> List[Dict]:
        chunks = []
        path = Path(export_path)
        if not path.exists():
            log.warning("Export file not found: %s", export_path)
            return chunks

        if export_type == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(path))
                for slide_num, slide in enumerate(prs.slides):
                    texts = [
                        shape.text.strip()
                        for shape in slide.shapes
                        if hasattr(shape, "text") and shape.text.strip()
                    ]
                    if texts:
                        combined = "\n".join(texts)
                        chunks.append({
                            "text": combined,
                            "slide_number": slide_num + 1,
                            "quality": self._assess_text_quality(combined),
                        })
            except ImportError:
                log.warning("python-pptx not installed")
            except Exception as e:
                log.warning("PPTX extraction error: %s", e)

        elif export_type == "json":
            try:
                data = json.loads(path.read_text(encoding="utf-8"))
                slides = data.get("slides") or data.get("result", {}).get("slides") or []
                for i, slide in enumerate(slides):
                    title = slide.get("title", "")
                    bullets = slide.get("bullets") or slide.get("content") or []
                    if isinstance(bullets, str):
                        bullets = [bullets]
                    text = title + "\n" + "\n".join(str(b) for b in bullets)
                    if text.strip():
                        chunks.append({
                            "text": text.strip(),
                            "slide_number": i + 1,
                            "quality": self._assess_text_quality(text),
                        })
            except Exception as e:
                log.warning("JSON extraction error: %s", e)

        elif export_type == "docx":
            try:
                from docx import Document
                doc = Document(str(path))
                section_text: List[str] = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        section_text.append(para.text.strip())
                        if len(" ".join(section_text).split()) > 500:
                            combined = "\n".join(section_text)
                            chunks.append({
                                "text": combined,
                                "quality": self._assess_text_quality(combined),
                            })
                            section_text = []
                if section_text:
                    combined = "\n".join(section_text)
                    chunks.append({
                        "text": combined,
                        "quality": self._assess_text_quality(combined),
                    })
            except ImportError:
                log.warning("python-docx not installed")
            except Exception as e:
                log.warning("DOCX extraction error: %s", e)

        return chunks

    def _assess_quality(self, chunks: List[Dict]) -> float:
        if not chunks:
            return 0.0
        scores = [chunk.get("quality", 0.5) for chunk in chunks]
        return sum(scores) / len(scores)

    def _assess_text_quality(self, text: str) -> float:
        score = 0.5
        words = text.split()
        word_count = len(words)
        if word_count >= 20:
            score += 0.05
        if word_count >= 50:
            score += 0.1
        if word_count >= 100:
            score += 0.05
        if any(
            line.strip().startswith(("-", "•", "*", "1.", "2.", "3."))
            for line in text.split("\n")
        ):
            score += 0.1
        keywords = [
            "marketing", "strategie", "strateg", "kunde", "kund", "zielgruppe",
            "kampagne", "content", "brand", "conversion", "roi", "kpi",
            "markt", "wettbewerb", "produkt", "lösung",
        ]
        hits = sum(1 for kw in keywords if kw.lower() in text.lower())
        score += min(0.15, hits * 0.03)
        sentences = text.split(".")
        if sentences:
            avg_len = word_count / max(len(sentences), 1)
            if avg_len > 6:
                score += 0.05
        return min(1.0, score)

    def _log_learning_event(self, session_id: str, quality_score: float, chunks: int):
        try:
            import requests as req
            req.post(
                "http://localhost:8011/knowledge/analytics/log/ingestion",
                params={
                    "source": f"self_learning_{session_id}",
                    "chunks_created": chunks,
                    "chunks_rejected": 0,
                    "duration_ms": 0,
                    "success": True,
                },
                json={"rejection_reasons": {}},
                timeout=5,
            )
        except Exception:
            pass

    def _save_export_metadata(self, result: Dict):
        try:
            meta_file = GENERATED_OUTPUTS_DIR / f"{result['session_id']}.json"
            with open(meta_file, "w", encoding="utf-8") as f:
                json.dump(result, f, indent=2, default=str)
        except Exception as e:
            log.warning("Could not save export metadata: %s", e)

    # ─────────────────────────────────────────
    # 8.2 USER FEEDBACK INTEGRATION
    # ─────────────────────────────────────────

    def record_feedback(
        self, chunk_id: str, score: int, session_id: Optional[str] = None
    ) -> Dict:
        # Normalisiere Score
        if score in (-1, 0):
            normalized = 2
        elif score == 1 and score != 1:  # thumbs up marker
            normalized = 4
        else:
            normalized = max(1, min(5, score))

        self.feedback_scores.setdefault(chunk_id, []).append(normalized)

        try:
            import requests as req
            req.post(
                "http://localhost:8011/knowledge/analytics/log/feedback",
                params={"chunk_id": chunk_id, "score": normalized, "session_id": session_id or ""},
                timeout=5,
            )
        except Exception:
            pass

        scores = self.feedback_scores[chunk_id]
        if len(scores) >= 3:
            avg = sum(scores) / len(scores)
            self._update_chunk_quality(chunk_id, avg / 5.0)

        return {
            "ok": True,
            "chunk_id": chunk_id,
            "feedback_count": len(scores),
            "avg_score": sum(scores) / len(scores),
        }

    def _update_chunk_quality(self, chunk_id: str, new_quality: float):
        client = _get_qdrant_client()
        if not client:
            return
        for collection in ("knowledge_base", "design_templates", QDRANT_COLLECTION):
            try:
                points = client.retrieve(
                    collection_name=collection,
                    ids=[chunk_id],
                    with_payload=True,
                )
                if points:
                    old = points[0].payload.get("quality_score", 0.5)
                    updated = round(0.6 * old + 0.4 * new_quality, 3)
                    client.set_payload(
                        collection_name=collection,
                        payload={"quality_score": updated, "feedback_adjusted": True},
                        points=[chunk_id],
                    )
                    log.info("Updated chunk %s quality: %.3f → %.3f", chunk_id, old, updated)
                    return
            except Exception:
                continue

    def get_low_quality_chunks(self, threshold: float = 0.4, limit: int = 20) -> List[Dict]:
        client = _get_qdrant_client()
        if not client:
            return []
        low: List[Dict] = []
        for collection in ("knowledge_base", QDRANT_COLLECTION):
            try:
                result = client.scroll(
                    collection_name=collection,
                    limit=500,
                    with_payload=True,
                    with_vectors=False,
                )[0]
                for point in result:
                    q = point.payload.get("quality_score", 0.5)
                    if q < threshold:
                        low.append({
                            "id": str(point.id),
                            "collection": collection,
                            "quality_score": q,
                            "source": point.payload.get("source_file", "unknown"),
                            "text_preview": point.payload.get("text", "")[:100],
                        })
            except Exception:
                continue
        low.sort(key=lambda x: x["quality_score"])
        return low[:limit]

    # ─────────────────────────────────────────
    # 8.3 SESSION ISOLATION
    # ─────────────────────────────────────────

    def create_session_collection(self, session_id: str) -> str:
        collection_name = f"session_{session_id}"
        client = _get_qdrant_client()
        if not client:
            return ""
        try:
            _ensure_collection(client, collection_name)
            self._register_session(session_id)
            return collection_name
        except Exception as e:
            log.warning("Session creation error: %s", e)
            return ""

    def add_to_session(self, session_id: str, text: str, metadata: Optional[Dict] = None) -> str:
        collection_name = f"session_{session_id}"
        chunk_id = str(uuid.uuid4())
        client = _get_qdrant_client()
        if not client:
            return ""
        try:
            from qdrant_client.models import PointStruct
            embedding = get_embedding(text)
            if not embedding:
                return ""
            client.upsert(
                collection_name=collection_name,
                points=[PointStruct(
                    id=chunk_id,
                    vector=embedding,
                    payload={
                        "text": text,
                        "session_id": session_id,
                        "indexed_at": datetime.now().isoformat(),
                        "quality_score": self._assess_text_quality(text),
                        **(metadata or {}),
                    },
                )],
            )
            return chunk_id
        except Exception as e:
            log.warning("Session add error: %s", e)
            return ""

    def search_session(self, session_id: str, query: str, limit: int = 5) -> List[Dict]:
        collection_name = f"session_{session_id}"
        client = _get_qdrant_client()
        if not client:
            return []
        try:
            query_vec = get_embedding(query)
            if not query_vec:
                return []
            results = client.search(
                collection_name=collection_name,
                query_vector=query_vec,
                limit=limit,
                with_payload=True,
            )
            return [
                {"id": str(r.id), "score": r.score, "text": r.payload.get("text", ""), "metadata": r.payload}
                for r in results
            ]
        except Exception as e:
            log.warning("Session search error: %s", e)
            return []

    def delete_session_collection(self, session_id: str) -> bool:
        client = _get_qdrant_client()
        if not client:
            return False
        try:
            client.delete_collection(f"session_{session_id}")
            self._unregister_session(session_id)
            return True
        except Exception as e:
            log.warning("Session delete error: %s", e)
            return False

    def _register_session(self, session_id: str):
        f = SESSION_COLLECTIONS_DIR / "active_sessions.json"
        try:
            sessions = json.loads(f.read_text()) if f.exists() else {}
            sessions[session_id] = {
                "created_at": datetime.now().isoformat(),
                "expires_at": (datetime.now() + timedelta(hours=SESSION_TTL_HOURS)).isoformat(),
            }
            f.write_text(json.dumps(sessions, indent=2))
        except Exception as e:
            log.warning("Session register error: %s", e)

    def _unregister_session(self, session_id: str):
        f = SESSION_COLLECTIONS_DIR / "active_sessions.json"
        try:
            if f.exists():
                sessions = json.loads(f.read_text())
                sessions.pop(session_id, None)
                f.write_text(json.dumps(sessions, indent=2))
        except Exception as e:
            log.warning("Session unregister error: %s", e)

    def _cleanup_old_sessions(self):
        f = SESSION_COLLECTIONS_DIR / "active_sessions.json"
        try:
            if not f.exists():
                return
            sessions = json.loads(f.read_text())
            now = datetime.now()
            expired = [
                sid for sid, info in sessions.items()
                if now > datetime.fromisoformat(info["expires_at"])
            ]
            for sid in expired:
                log.info("Cleaning up expired session: %s", sid)
                self.delete_session_collection(sid)
            if expired:
                log.info("Cleaned up %d expired sessions", len(expired))
        except Exception as e:
            log.warning("Session cleanup error: %s", e)


# ─────────────────────────────────────────────
# SINGLETON
# ─────────────────────────────────────────────

self_learning = SelfLearningSystem()


# ─────────────────────────────────────────────
# API ENDPOINTS
# ─────────────────────────────────────────────

def register_self_learning_routes(app):
    """Registriert Self-Learning API Endpoints."""
    from fastapi import APIRouter

    router = APIRouter(prefix="/learning", tags=["Self-Learning"])

    @router.post("/export-complete")
    async def on_export_complete(
        export_path: str, session_id: str, export_type: str = "pptx"
    ):
        """Callback wenn ein Export abgeschlossen ist. Indexiert gute Inhalte in Qdrant."""
        result = self_learning.on_export_complete(export_path, session_id, export_type)
        return result

    @router.post("/feedback")
    async def record_feedback(chunk_id: str, score: int, session_id: str = None):
        """Nimmt User-Feedback auf (1-5)."""
        return self_learning.record_feedback(chunk_id, score, session_id)

    @router.get("/low-quality")
    async def get_low_quality(threshold: float = 0.4, limit: int = 20):
        """Gibt niedrig bewertete Chunks zurück (Cleanup-Kandidaten)."""
        chunks = self_learning.get_low_quality_chunks(threshold, limit)
        return {"ok": True, "chunks": chunks, "count": len(chunks)}

    @router.post("/session/create")
    async def create_session(session_id: str):
        """Erstellt eine isolierte Session-Collection."""
        collection = self_learning.create_session_collection(session_id)
        return {"ok": bool(collection), "collection": collection}

    @router.post("/session/{session_id}/add")
    async def add_to_session(session_id: str, text: str, source: str = None):
        """Fügt Content zur Session hinzu."""
        chunk_id = self_learning.add_to_session(
            session_id, text, {"source": source} if source else None
        )
        return {"ok": bool(chunk_id), "chunk_id": chunk_id}

    @router.get("/session/{session_id}/search")
    async def search_session(session_id: str, query: str, limit: int = 5):
        """Sucht in der Session-Collection."""
        results = self_learning.search_session(session_id, query, limit)
        return {"ok": True, "results": results}

    @router.delete("/session/{session_id}")
    async def delete_session(session_id: str):
        """Löscht eine Session-Collection."""
        return {"ok": self_learning.delete_session_collection(session_id)}

    @router.get("/stats")
    async def learning_stats():
        """Gibt Lernstatistiken zurück."""
        meta_files = list(GENERATED_OUTPUTS_DIR.glob("*.json"))
        indexed = sum(
            1 for f in meta_files
            if json.loads(f.read_text()).get("indexed", False)
        )
        return {
            "ok": True,
            "total_exports_processed": len(meta_files),
            "total_indexed": indexed,
            "feedback_tracked": len(self_learning.feedback_scores),
        }

    return router
