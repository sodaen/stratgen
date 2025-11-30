# -*- coding: utf-8 -*-
"""
services/slide_dna_analyzer.py
==============================
Feature: Slide DNA Analyzer

Analysiert ALLE Templates und extrahiert:
1. Erfolgreiche Slide-Sequenzen
2. Optimale Bullet-Längen pro Slide-Typ
3. Wort-/Phrasen-Muster
4. Layout-Präferenzen pro Branche/Thema
5. DNA-Profile für Generierung

Author: StratGen Agent V3.6
"""
from __future__ import annotations
import os
import re
import json
import sqlite3
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field, asdict
from collections import Counter, defaultdict
from datetime import datetime

# ============================================
# CONFIGURATION
# ============================================

RAW_DIR = os.getenv("STRATGEN_RAW_DIR", "data/raw")
DB_PATH = os.getenv("STRATGEN_DNA_DB", "data/slide_dna.sqlite")

# ============================================
# DATA CLASSES
# ============================================

@dataclass
class SlidePattern:
    """Ein erkanntes Slide-Muster."""
    slide_type: str
    avg_bullets: float
    avg_bullet_length: float
    common_words: List[str]
    common_phrases: List[str]
    typical_titles: List[str]
    occurrence_count: int


@dataclass
class SequencePattern:
    """Eine erkannte Slide-Sequenz."""
    sequence: List[str]  # z.B. ["title", "problem", "solution", "benefits"]
    frequency: int
    success_score: float  # Basierend auf wie oft verwendet
    source_templates: List[str]


@dataclass
class SlideDNA:
    """Das DNA-Profil aller Templates."""
    total_templates: int
    total_slides: int
    slide_patterns: Dict[str, SlidePattern]
    common_sequences: List[SequencePattern]
    global_vocabulary: List[str]
    industry_patterns: Dict[str, Dict[str, Any]]
    analyzed_at: str


# ============================================
# DATABASE
# ============================================

def _connect():
    """Verbindet zur DNA-Datenbank."""
    os.makedirs("data", exist_ok=True)
    con = sqlite3.connect(DB_PATH, timeout=30)
    con.row_factory = sqlite3.Row
    return con


def ensure_tables():
    """Erstellt Tabellen falls nicht vorhanden."""
    con = _connect()
    cur = con.cursor()
    
    # Analysierte Templates
    cur.execute("""
    CREATE TABLE IF NOT EXISTS analyzed_templates (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        file_path TEXT NOT NULL UNIQUE,
        file_hash TEXT,
        slide_count INTEGER,
        analyzed_at TEXT DEFAULT CURRENT_TIMESTAMP,
        metadata TEXT  -- JSON
    )
    """)
    
    # Extrahierte Slides
    cur.execute("""
    CREATE TABLE IF NOT EXISTS extracted_slides (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        template_id INTEGER,
        slide_index INTEGER,
        slide_type TEXT,
        title TEXT,
        bullet_count INTEGER,
        avg_bullet_length REAL,
        word_count INTEGER,
        content_hash TEXT,
        raw_text TEXT,
        FOREIGN KEY (template_id) REFERENCES analyzed_templates(id)
    )
    """)
    
    # Slide-Sequenzen
    cur.execute("""
    CREATE TABLE IF NOT EXISTS slide_sequences (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        template_id INTEGER,
        sequence TEXT,  -- JSON Array
        sequence_length INTEGER,
        FOREIGN KEY (template_id) REFERENCES analyzed_templates(id)
    )
    """)
    
    # Vokabular
    cur.execute("""
    CREATE TABLE IF NOT EXISTS vocabulary (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        word TEXT NOT NULL,
        slide_type TEXT,
        frequency INTEGER DEFAULT 1,
        UNIQUE(word, slide_type)
    )
    """)
    
    # Phrasen
    cur.execute("""
    CREATE TABLE IF NOT EXISTS phrases (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        phrase TEXT NOT NULL,
        slide_type TEXT,
        frequency INTEGER DEFAULT 1,
        UNIQUE(phrase, slide_type)
    )
    """)
    
    con.commit()
    con.close()


# ============================================
# SLIDE TYPE DETECTION
# ============================================

SLIDE_TYPE_KEYWORDS = {
    "title": ["agenda", "inhalt", "gliederung", "titel", "willkommen"],
    "executive_summary": ["summary", "zusammenfassung", "überblick", "executive", "key findings"],
    "problem": ["problem", "herausforderung", "challenge", "pain", "issue", "schwierigkeit"],
    "solution": ["lösung", "solution", "ansatz", "approach", "antwort", "vorschlag"],
    "benefits": ["nutzen", "vorteile", "benefits", "mehrwert", "value", "why"],
    "roi": ["roi", "return", "investment", "business case", "kosten", "einsparung"],
    "roadmap": ["roadmap", "timeline", "zeitplan", "meilenstein", "milestone", "phase"],
    "team": ["team", "über uns", "about", "wer wir sind", "experten"],
    "case_study": ["case study", "referenz", "erfolgsgeschichte", "kunde", "projekt"],
    "competitive": ["wettbewerb", "vergleich", "konkurrenz", "markt", "positioning"],
    "next_steps": ["next steps", "nächste schritte", "empfehlung", "action", "todo"],
    "contact": ["kontakt", "contact", "fragen", "discussion", "vielen dank", "thank"],
    "appendix": ["appendix", "anhang", "backup", "zusatz", "detail"],
}


def detect_slide_type(title: str, content: str) -> str:
    """Erkennt den Slide-Typ basierend auf Titel und Content."""
    text = f"{title} {content}".lower()
    
    scores = {}
    for slide_type, keywords in SLIDE_TYPE_KEYWORDS.items():
        score = sum(1 for kw in keywords if kw in text)
        if score > 0:
            scores[slide_type] = score
    
    if scores:
        return max(scores, key=scores.get)
    
    return "content"  # Default


# ============================================
# TEMPLATE ANALYSIS
# ============================================

def analyze_template(file_path: str) -> Dict[str, Any]:
    """
    Analysiert ein einzelnes Template.
    
    Returns:
        Dictionary mit Analyse-Ergebnissen
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed"}
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        return {"ok": False, "error": str(e)}
    
    slides_data = []
    slide_types = []
    all_words = []
    all_phrases = []
    
    for idx, slide in enumerate(prs.slides):
        slide_text = []
        title = ""
        bullets = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
                        
                        # Erster Text ist oft Titel
                        if not title and len(text) < 100:
                            title = text
                        elif len(text) > 10:
                            bullets.append(text)
        
        full_text = " ".join(slide_text)
        slide_type = detect_slide_type(title, full_text)
        slide_types.append(slide_type)
        
        # Wörter extrahieren
        words = re.findall(r'\b[A-ZÄÖÜa-zäöüß]{4,}\b', full_text)
        all_words.extend(words)
        
        # Phrasen extrahieren (2-3 Wörter)
        word_list = full_text.split()
        for i in range(len(word_list) - 1):
            phrase = f"{word_list[i]} {word_list[i+1]}"
            if len(phrase) > 8:
                all_phrases.append(phrase.lower())
        
        slides_data.append({
            "index": idx,
            "type": slide_type,
            "title": title,
            "bullet_count": len(bullets),
            "avg_bullet_length": sum(len(b) for b in bullets) / max(1, len(bullets)),
            "word_count": len(words),
            "raw_text": full_text[:1000]
        })
    
    return {
        "ok": True,
        "file_path": file_path,
        "slide_count": len(prs.slides),
        "slides": slides_data,
        "sequence": slide_types,
        "words": all_words,
        "phrases": all_phrases
    }


def analyze_all_templates(force_reanalyze: bool = False) -> Dict[str, Any]:
    """
    Analysiert alle Templates im RAW_DIR.
    
    Args:
        force_reanalyze: Alle neu analysieren, auch wenn bereits in DB
    
    Returns:
        Zusammenfassung der Analyse
    """
    ensure_tables()
    
    raw_path = Path(RAW_DIR)
    if not raw_path.exists():
        return {"ok": False, "error": f"RAW_DIR not found: {RAW_DIR}"}
    
    pptx_files = list(raw_path.glob("**/*.pptx"))
    
    con = _connect()
    cur = con.cursor()
    
    analyzed = 0
    skipped = 0
    errors = []
    
    for pptx_file in pptx_files:
        file_path = str(pptx_file)
        
        # Bereits analysiert?
        if not force_reanalyze:
            cur.execute("SELECT id FROM analyzed_templates WHERE file_path = ?", (file_path,))
            if cur.fetchone():
                skipped += 1
                continue
        
        # Analysieren
        result = analyze_template(file_path)
        
        if not result.get("ok"):
            errors.append({"file": file_path, "error": result.get("error")})
            continue
        
        # In DB speichern
        cur.execute("""
        INSERT OR REPLACE INTO analyzed_templates (file_path, slide_count, metadata)
        VALUES (?, ?, ?)
        """, (file_path, result["slide_count"], json.dumps({"analyzed": True})))
        
        template_id = cur.lastrowid
        
        # Slides speichern
        for slide in result["slides"]:
            cur.execute("""
            INSERT INTO extracted_slides 
            (template_id, slide_index, slide_type, title, bullet_count, avg_bullet_length, word_count, raw_text)
            VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """, (template_id, slide["index"], slide["type"], slide["title"],
                  slide["bullet_count"], slide["avg_bullet_length"], slide["word_count"], slide["raw_text"]))
        
        # Sequenz speichern
        cur.execute("""
        INSERT INTO slide_sequences (template_id, sequence, sequence_length)
        VALUES (?, ?, ?)
        """, (template_id, json.dumps(result["sequence"]), len(result["sequence"])))
        
        # Vokabular updaten
        word_counts = Counter(w.lower() for w in result["words"])
        for word, count in word_counts.most_common(50):
            cur.execute("""
            INSERT INTO vocabulary (word, frequency) VALUES (?, ?)
            ON CONFLICT(word, slide_type) DO UPDATE SET frequency = frequency + ?
            """, (word, count, count))
        
        analyzed += 1
        
        if analyzed % 10 == 0:
            con.commit()
    
    con.commit()
    con.close()
    
    return {
        "ok": True,
        "total_files": len(pptx_files),
        "analyzed": analyzed,
        "skipped": skipped,
        "errors": len(errors),
        "error_details": errors[:5]
    }


# ============================================
# DNA EXTRACTION
# ============================================

def extract_slide_dna() -> SlideDNA:
    """
    Extrahiert das Slide-DNA-Profil aus allen analysierten Templates.
    
    Returns:
        SlideDNA-Objekt
    """
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    
    # Template-Statistiken
    cur.execute("SELECT COUNT(*) as cnt FROM analyzed_templates")
    total_templates = cur.fetchone()["cnt"]
    
    cur.execute("SELECT COUNT(*) as cnt FROM extracted_slides")
    total_slides = cur.fetchone()["cnt"]
    
    # Slide-Patterns pro Typ
    slide_patterns = {}
    cur.execute("""
    SELECT slide_type, 
           COUNT(*) as cnt,
           AVG(bullet_count) as avg_bullets,
           AVG(avg_bullet_length) as avg_bullet_len
    FROM extracted_slides
    GROUP BY slide_type
    """)
    
    for row in cur.fetchall():
        slide_type = row["slide_type"]
        
        # Häufige Titel für diesen Typ
        cur.execute("""
        SELECT title, COUNT(*) as cnt FROM extracted_slides
        WHERE slide_type = ? AND title != ''
        GROUP BY title ORDER BY cnt DESC LIMIT 10
        """, (slide_type,))
        titles = [r["title"] for r in cur.fetchall()]
        
        slide_patterns[slide_type] = SlidePattern(
            slide_type=slide_type,
            avg_bullets=round(row["avg_bullets"] or 0, 1),
            avg_bullet_length=round(row["avg_bullet_len"] or 0, 1),
            common_words=[],
            common_phrases=[],
            typical_titles=titles[:5],
            occurrence_count=row["cnt"]
        )
    
    # Häufige Sequenzen
    cur.execute("SELECT sequence, COUNT(*) as cnt FROM slide_sequences GROUP BY sequence ORDER BY cnt DESC LIMIT 20")
    sequences = []
    for row in cur.fetchall():
        seq = json.loads(row["sequence"])
        sequences.append(SequencePattern(
            sequence=seq,
            frequency=row["cnt"],
            success_score=min(1.0, row["cnt"] / 5),
            source_templates=[]
        ))
    
    # Globales Vokabular
    cur.execute("SELECT word, SUM(frequency) as freq FROM vocabulary GROUP BY word ORDER BY freq DESC LIMIT 100")
    vocabulary = [row["word"] for row in cur.fetchall()]
    
    con.close()
    
    return SlideDNA(
        total_templates=total_templates,
        total_slides=total_slides,
        slide_patterns=slide_patterns,
        common_sequences=sequences[:10],
        global_vocabulary=vocabulary[:50],
        industry_patterns={},
        analyzed_at=datetime.now().isoformat()
    )


# ============================================
# PATTERN MATCHING
# ============================================

def get_optimal_structure(
    topic: str,
    deck_size: str = "medium",
    industry: str = ""
) -> Dict[str, Any]:
    """
    Empfiehlt optimale Slide-Struktur basierend auf DNA.
    
    Args:
        topic: Hauptthema
        deck_size: short/medium/long
        industry: Branche
    
    Returns:
        Empfohlene Struktur mit Begründung
    """
    dna = extract_slide_dna()
    
    # Deck-Size zu Anzahl
    size_map = {"short": 5, "medium": 10, "long": 20}
    target_slides = size_map.get(deck_size, 10)
    
    # Beste Sequenz finden die zur Größe passt
    best_sequence = None
    for seq in dna.common_sequences:
        if abs(len(seq.sequence) - target_slides) < 3:
            best_sequence = seq
            break
    
    if not best_sequence and dna.common_sequences:
        best_sequence = dna.common_sequences[0]
    
    # Slide-Details ergänzen
    recommended_slides = []
    if best_sequence:
        for slide_type in best_sequence.sequence:
            pattern = dna.slide_patterns.get(slide_type)
            recommended_slides.append({
                "type": slide_type,
                "recommended_bullets": int(pattern.avg_bullets) if pattern else 4,
                "typical_titles": pattern.typical_titles[:3] if pattern else [],
                "avg_bullet_length": int(pattern.avg_bullet_length) if pattern else 50
            })
    
    return {
        "ok": True,
        "recommended_structure": recommended_slides,
        "sequence": best_sequence.sequence if best_sequence else [],
        "confidence": best_sequence.success_score if best_sequence else 0.5,
        "based_on_templates": dna.total_templates,
        "reasoning": f"Basierend auf {dna.total_templates} analysierten Templates"
    }


def get_slide_recommendations(slide_type: str) -> Dict[str, Any]:
    """
    Gibt Empfehlungen für einen spezifischen Slide-Typ.
    
    Args:
        slide_type: Der Slide-Typ
    
    Returns:
        Empfehlungen mit Beispielen
    """
    dna = extract_slide_dna()
    pattern = dna.slide_patterns.get(slide_type)
    
    if not pattern:
        return {"ok": False, "error": f"Unknown slide type: {slide_type}"}
    
    return {
        "ok": True,
        "slide_type": slide_type,
        "recommendations": {
            "bullet_count": int(pattern.avg_bullets),
            "bullet_length": int(pattern.avg_bullet_length),
            "typical_titles": pattern.typical_titles,
            "occurrence_in_templates": pattern.occurrence_count
        },
        "examples_available": pattern.occurrence_count
    }


# ============================================
# API FUNCTIONS
# ============================================

def check_status() -> Dict[str, Any]:
    """Gibt den Status des DNA Analyzers zurück."""
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    
    cur.execute("SELECT COUNT(*) as cnt FROM analyzed_templates")
    templates = cur.fetchone()["cnt"]
    
    cur.execute("SELECT COUNT(*) as cnt FROM extracted_slides")
    slides = cur.fetchone()["cnt"]
    
    cur.execute("SELECT COUNT(DISTINCT slide_type) as cnt FROM extracted_slides")
    types = cur.fetchone()["cnt"]
    
    con.close()
    
    return {
        "ok": True,
        "analyzed_templates": templates,
        "extracted_slides": slides,
        "slide_types_found": types,
        "db_path": DB_PATH
    }


def get_dna_summary() -> Dict[str, Any]:
    """Gibt eine Zusammenfassung der Slide-DNA zurück."""
    dna = extract_slide_dna()
    
    return {
        "ok": True,
        "total_templates": dna.total_templates,
        "total_slides": dna.total_slides,
        "slide_types": {
            st: {
                "count": p.occurrence_count,
                "avg_bullets": p.avg_bullets,
                "typical_titles": p.typical_titles[:3]
            }
            for st, p in dna.slide_patterns.items()
        },
        "top_sequences": [
            {"sequence": s.sequence, "frequency": s.frequency}
            for s in dna.common_sequences[:5]
        ],
        "analyzed_at": dna.analyzed_at
    }
