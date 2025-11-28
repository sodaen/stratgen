# -*- coding: utf-8 -*-
"""
services/template_learner.py
============================
Lernt Präsentationsstrukturen aus bestehenden PPTX in /raw.
Extrahiert: Layouts, Farbschemas, Slide-Strukturen, Best Practices.

Funktionen:
- scan_templates(): Scannt /raw nach PPTX und analysiert sie
- extract_structure(): Extrahiert Struktur einer PPTX
- learn_patterns(): Lernt wiederkehrende Muster
- suggest_structure(): Empfiehlt Struktur basierend auf gelernten Mustern
"""
from __future__ import annotations
import os
import re
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
from collections import Counter

# python-pptx für PPTX-Analyse
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import PP_PLACEHOLDER
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# ============================================
# KONFIGURATION
# ============================================

RAW_DIR = os.getenv("STRATGEN_RAW_DIR", "data/raw")
TEMPLATES_DB = os.getenv("STRATGEN_TEMPLATES_DB", "data/templates_learned.json")

# Bekannte Slide-Typen (für Pattern-Matching)
SLIDE_TYPE_PATTERNS = {
    "title": ["titel", "title", "cover", "start", "opening"],
    "agenda": ["agenda", "inhalt", "overview", "contents", "gliederung"],
    "executive_summary": ["executive", "summary", "zusammenfassung", "überblick", "key takeaways"],
    "problem": ["problem", "herausforderung", "challenge", "pain", "issue"],
    "solution": ["lösung", "solution", "approach", "ansatz", "our approach"],
    "use_case": ["use case", "use-case", "anwendung", "beispiel", "case study"],
    "benefits": ["benefits", "vorteile", "nutzen", "value", "advantages"],
    "roi": ["roi", "business case", "kosten", "savings", "investment", "payback"],
    "roadmap": ["roadmap", "timeline", "zeitplan", "phases", "milestones", "plan"],
    "team": ["team", "über uns", "about us", "wer wir sind", "organization"],
    "competitive": ["wettbewerb", "competition", "vergleich", "comparison", "vs"],
    "risks": ["risiken", "risks", "mitigation", "challenges"],
    "next_steps": ["next steps", "nächste schritte", "action", "call to action", "cta"],
    "appendix": ["appendix", "anhang", "backup", "additional"],
    "contact": ["kontakt", "contact", "questions", "fragen", "q&a"],
    "thank_you": ["thank", "danke", "ende", "end", "closing"],
}

# ============================================
# STORAGE
# ============================================

def _load_db() -> Dict[str, Any]:
    """Lädt die Template-Datenbank."""
    if os.path.exists(TEMPLATES_DB):
        try:
            with open(TEMPLATES_DB, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            pass
    return {
        "templates": {},
        "patterns": {
            "common_structures": [],
            "slide_type_frequency": {},
            "avg_bullets_per_slide": 0,
            "avg_slides_per_deck": 0,
            "color_schemes": [],
        },
        "updated_at": None
    }


def _save_db(db: Dict[str, Any]) -> None:
    """Speichert die Template-Datenbank."""
    os.makedirs(os.path.dirname(TEMPLATES_DB), exist_ok=True)
    db["updated_at"] = datetime.now().isoformat()
    with open(TEMPLATES_DB, "w", encoding="utf-8") as f:
        json.dump(db, f, indent=2, ensure_ascii=False)


def _file_hash(path: str) -> str:
    """Berechnet Hash für eindeutige ID."""
    h = hashlib.md5()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()[:12]


# ============================================
# SLIDE-TYP ERKENNUNG
# ============================================

def _detect_slide_type(title: str, content: str = "") -> str:
    """Erkennt den Slide-Typ basierend auf Titel und Content."""
    text = f"{title} {content}".lower()
    
    for slide_type, patterns in SLIDE_TYPE_PATTERNS.items():
        for pattern in patterns:
            if pattern in text:
                return slide_type
    
    return "content"  # Default


def _extract_colors_from_shape(shape) -> List[str]:
    """Extrahiert Farben aus einem Shape."""
    colors = []
    try:
        # Fill-Farbe
        if hasattr(shape, "fill") and shape.fill:
            if hasattr(shape.fill, "fore_color") and shape.fill.fore_color:
                try:
                    rgb = shape.fill.fore_color.rgb
                    if rgb:
                        colors.append(f"#{rgb}")
                except Exception:
                    pass
        
        # Text-Farbe
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.color and run.font.color.rgb:
                        colors.append(f"#{run.font.color.rgb}")
    except Exception:
        pass
    
    return colors


# ============================================
# PPTX ANALYSE
# ============================================

def extract_structure(pptx_path: str) -> Dict[str, Any]:
    """
    Extrahiert die Struktur einer PPTX-Datei.
    
    Args:
        pptx_path: Pfad zur PPTX-Datei
    
    Returns:
        Dict mit: slides, layouts, colors, metadata
    """
    if not HAS_PPTX:
        return {"ok": False, "error": "python-pptx not installed"}
    
    if not os.path.exists(pptx_path):
        return {"ok": False, "error": "File not found"}
    
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"ok": False, "error": f"Cannot open PPTX: {str(e)}"}
    
    # Metadaten
    metadata = {
        "filename": os.path.basename(pptx_path),
        "slide_count": len(prs.slides),
        "slide_width": prs.slide_width.inches if prs.slide_width else None,
        "slide_height": prs.slide_height.inches if prs.slide_height else None,
    }
    
    # Layouts sammeln
    layouts_used = []
    for layout in prs.slide_layouts:
        try:
            layouts_used.append({
                "name": layout.name,
                "placeholders": len(list(layout.placeholders))
            })
        except Exception:
            pass
    
    # Slides analysieren
    slides = []
    all_colors = []
    total_bullets = 0
    
    for i, slide in enumerate(prs.slides):
        slide_data = {
            "index": i,
            "title": "",
            "bullet_count": 0,
            "has_image": False,
            "has_chart": False,
            "has_table": False,
            "type": "content",
            "layout_name": None,
        }
        
        # Layout-Name
        try:
            if slide.slide_layout:
                slide_data["layout_name"] = slide.slide_layout.name
        except Exception:
            pass
        
        # Shapes analysieren
        content_text = []
        for shape in slide.shapes:
            # Titel
            if shape.is_placeholder:
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type == PP_PLACEHOLDER.TITLE:
                        slide_data["title"] = shape.text or ""
                except Exception:
                    pass
            
            # Text/Bullets
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        content_text.append(para.text.strip())
                        slide_data["bullet_count"] += 1
            
            # Bilder
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                slide_data["has_image"] = True
            
            # Charts
            if hasattr(shape, "has_chart") and shape.has_chart:
                slide_data["has_chart"] = True
            
            # Tabellen
            if hasattr(shape, "has_table") and shape.has_table:
                slide_data["has_table"] = True
            
            # Farben extrahieren
            colors = _extract_colors_from_shape(shape)
            all_colors.extend(colors)
        
        # Slide-Typ erkennen
        slide_data["type"] = _detect_slide_type(
            slide_data["title"], 
            " ".join(content_text)
        )
        
        total_bullets += slide_data["bullet_count"]
        slides.append(slide_data)
    
    # Farben aggregieren (Top 10)
    color_counts = Counter(all_colors)
    top_colors = [c for c, _ in color_counts.most_common(10)]
    
    return {
        "ok": True,
        "metadata": metadata,
        "slides": slides,
        "layouts": layouts_used,
        "colors": top_colors,
        "stats": {
            "total_slides": len(slides),
            "total_bullets": total_bullets,
            "avg_bullets_per_slide": round(total_bullets / len(slides), 1) if slides else 0,
            "slides_with_images": sum(1 for s in slides if s["has_image"]),
            "slides_with_charts": sum(1 for s in slides if s["has_chart"]),
        }
    }


# ============================================
# TEMPLATE SCANNING
# ============================================

def scan_templates(directory: str = None) -> Dict[str, Any]:
    """
    Scannt ein Verzeichnis nach PPTX und analysiert alle.
    
    Args:
        directory: Zu scannendes Verzeichnis (default: RAW_DIR)
    
    Returns:
        Dict mit: scanned, templates, patterns
    """
    directory = directory or RAW_DIR
    
    if not os.path.isdir(directory):
        return {"ok": False, "error": f"Directory not found: {directory}"}
    
    db = _load_db()
    results = {
        "ok": True,
        "scanned": 0,
        "new": 0,
        "updated": 0,
        "errors": []
    }
    
    # Alle PPTX finden
    pptx_files = []
    for root, _, files in os.walk(directory):
        for f in files:
            if f.lower().endswith(".pptx") and not f.startswith("~"):
                pptx_files.append(os.path.join(root, f))
    
    # Analysieren
    for pptx_path in pptx_files:
        try:
            file_hash = _file_hash(pptx_path)
            
            # Bereits analysiert?
            existing = db["templates"].get(file_hash)
            if existing and existing.get("path") == pptx_path:
                results["scanned"] += 1
                continue
            
            # Neue Analyse
            structure = extract_structure(pptx_path)
            if not structure.get("ok"):
                results["errors"].append(f"{pptx_path}: {structure.get('error')}")
                continue
            
            # In DB speichern
            db["templates"][file_hash] = {
                "id": file_hash,
                "path": pptx_path,
                "filename": os.path.basename(pptx_path),
                "structure": structure,
                "analyzed_at": datetime.now().isoformat()
            }
            
            if existing:
                results["updated"] += 1
            else:
                results["new"] += 1
            results["scanned"] += 1
            
        except Exception as e:
            results["errors"].append(f"{pptx_path}: {str(e)}")
    
    # Patterns aktualisieren
    _update_patterns(db)
    _save_db(db)
    
    results["total_templates"] = len(db["templates"])
    return results


def _update_patterns(db: Dict[str, Any]) -> None:
    """Aktualisiert die gelernten Patterns basierend auf allen Templates."""
    templates = list(db["templates"].values())
    if not templates:
        return
    
    # Slide-Typ Häufigkeit
    type_counter = Counter()
    for t in templates:
        slides = t.get("structure", {}).get("slides", [])
        for s in slides:
            type_counter[s.get("type", "content")] += 1
    
    db["patterns"]["slide_type_frequency"] = dict(type_counter)
    
    # Durchschnittliche Bullets pro Slide
    total_bullets = 0
    total_slides = 0
    for t in templates:
        stats = t.get("structure", {}).get("stats", {})
        total_bullets += stats.get("total_bullets", 0)
        total_slides += stats.get("total_slides", 0)
    
    db["patterns"]["avg_bullets_per_slide"] = round(total_bullets / total_slides, 1) if total_slides else 4
    db["patterns"]["avg_slides_per_deck"] = round(total_slides / len(templates), 0) if templates else 15
    
    # Häufigste Strukturen (Slide-Typ-Sequenzen)
    structures = []
    for t in templates:
        slides = t.get("structure", {}).get("slides", [])
        sequence = [s.get("type", "content") for s in slides[:10]]  # Erste 10 Slides
        structures.append(sequence)
    
    # Finde häufigste Anfangs-Sequenzen
    start_sequences = Counter()
    for seq in structures:
        if len(seq) >= 3:
            start_sequences[tuple(seq[:3])] += 1
    
    db["patterns"]["common_structures"] = [
        list(seq) for seq, _ in start_sequences.most_common(5)
    ]
    
    # Farbschemas
    all_colors = []
    for t in templates:
        colors = t.get("structure", {}).get("colors", [])
        if colors:
            all_colors.append(colors[:5])  # Top 5 Farben pro Template
    
    db["patterns"]["color_schemes"] = all_colors[:10]  # Max 10 Schemas


# ============================================
# STRUKTUR-VORSCHLÄGE
# ============================================

def suggest_structure(
    deck_size: str = "medium",
    topic: str = "",
    industry: str = ""
) -> Dict[str, Any]:
    """
    Empfiehlt eine Deck-Struktur basierend auf gelernten Mustern.
    
    Args:
        deck_size: "short" (5-10), "medium" (15-25), "large" (30-50)
        topic: Thema des Decks (für Anpassungen)
        industry: Branche (für Anpassungen)
    
    Returns:
        Dict mit: slides (Liste von Slide-Typen), recommendations
    """
    db = _load_db()
    patterns = db.get("patterns", {})
    
    # Basis-Struktur je nach Größe
    if deck_size == "short":
        base_structure = [
            {"type": "title", "title": "Titel"},
            {"type": "executive_summary", "title": "Executive Summary"},
            {"type": "problem", "title": "Herausforderung"},
            {"type": "solution", "title": "Unser Ansatz"},
            {"type": "benefits", "title": "Ihr Nutzen"},
            {"type": "next_steps", "title": "Nächste Schritte"},
            {"type": "contact", "title": "Kontakt"},
        ]
    elif deck_size == "large":
        base_structure = [
            {"type": "title", "title": "Titel"},
            {"type": "agenda", "title": "Agenda"},
            {"type": "executive_summary", "title": "Executive Summary"},
            {"type": "problem", "title": "Herausforderung"},
            {"type": "problem", "title": "Aktuelle Situation"},
            {"type": "solution", "title": "Unser Ansatz"},
            {"type": "solution", "title": "Lösungsarchitektur"},
            {"type": "use_case", "title": "Use Case 1"},
            {"type": "use_case", "title": "Use Case 2"},
            {"type": "use_case", "title": "Use Case 3"},
            {"type": "benefits", "title": "Ihr Nutzen"},
            {"type": "roi", "title": "ROI & Business Case"},
            {"type": "roi", "title": "Kostenübersicht"},
            {"type": "roadmap", "title": "Roadmap"},
            {"type": "roadmap", "title": "Meilensteine"},
            {"type": "team", "title": "Unser Team"},
            {"type": "competitive", "title": "Marktvergleich"},
            {"type": "risks", "title": "Risiken & Mitigation"},
            {"type": "next_steps", "title": "Nächste Schritte"},
            {"type": "appendix", "title": "Anhang"},
            {"type": "contact", "title": "Kontakt"},
            {"type": "thank_you", "title": "Vielen Dank"},
        ]
    else:  # medium
        base_structure = [
            {"type": "title", "title": "Titel"},
            {"type": "agenda", "title": "Agenda"},
            {"type": "executive_summary", "title": "Executive Summary"},
            {"type": "problem", "title": "Herausforderung"},
            {"type": "solution", "title": "Unser Ansatz"},
            {"type": "use_case", "title": "Use Case 1"},
            {"type": "use_case", "title": "Use Case 2"},
            {"type": "benefits", "title": "Ihr Nutzen"},
            {"type": "roi", "title": "ROI & Business Case"},
            {"type": "roadmap", "title": "Roadmap"},
            {"type": "team", "title": "Unser Team"},
            {"type": "risks", "title": "Risiken"},
            {"type": "next_steps", "title": "Nächste Schritte"},
            {"type": "contact", "title": "Kontakt"},
        ]
    
    # Empfehlungen aus Patterns
    recommendations = []
    
    avg_bullets = patterns.get("avg_bullets_per_slide", 4)
    recommendations.append(f"Durchschnittlich {avg_bullets} Bullets pro Slide empfohlen")
    
    if patterns.get("slide_type_frequency"):
        top_types = sorted(
            patterns["slide_type_frequency"].items(), 
            key=lambda x: -x[1]
        )[:5]
        recommendations.append(f"Häufigste Slide-Typen: {', '.join(t for t, _ in top_types)}")
    
    # Anpassungen basierend auf Topic
    topic_lower = topic.lower()
    if "roi" in topic_lower or "business case" in topic_lower:
        # ROI-fokussiertes Deck braucht mehr Zahlen-Slides
        for i, s in enumerate(base_structure):
            if s["type"] == "roi":
                base_structure.insert(i+1, {"type": "roi", "title": "Detaillierte Kalkulation"})
                break
        recommendations.append("ROI-Fokus erkannt: Zusätzliche Kalkulationsfolie eingefügt")
    
    if "technical" in topic_lower or "integration" in topic_lower:
        # Technisches Deck
        for i, s in enumerate(base_structure):
            if s["type"] == "solution":
                base_structure.insert(i+1, {"type": "content", "title": "Technische Architektur"})
                break
        recommendations.append("Technischer Fokus: Architektur-Folie eingefügt")
    
    return {
        "ok": True,
        "deck_size": deck_size,
        "slides": base_structure,
        "slide_count": len(base_structure),
        "recommendations": recommendations,
        "learned_from": len(db.get("templates", {})),
    }


# ============================================
# TEMPLATE STATISTICS
# ============================================

def get_statistics() -> Dict[str, Any]:
    """
    Gibt Statistiken über gelernte Templates zurück.
    """
    db = _load_db()
    templates = list(db["templates"].values())
    
    if not templates:
        return {
            "ok": True,
            "template_count": 0,
            "message": "Keine Templates gelernt. Führe scan_templates() aus."
        }
    
    return {
        "ok": True,
        "template_count": len(templates),
        "patterns": db.get("patterns", {}),
        "templates": [
            {
                "filename": t["filename"],
                "slides": t["structure"]["metadata"]["slide_count"],
                "analyzed": t["analyzed_at"]
            }
            for t in templates
        ],
        "updated_at": db.get("updated_at")
    }


# ============================================
# TEST
# ============================================

if __name__ == "__main__":
    print("=== Template Learner Test ===\n")
    
    if not HAS_PPTX:
        print("ERROR: python-pptx not installed!")
        print("Install with: pip install python-pptx")
        exit(1)
    
    # Scan
    print("--- Scanning /raw ---")
    result = scan_templates()
    print(f"Scanned: {result.get('scanned')}")
    print(f"New: {result.get('new')}")
    print(f"Total: {result.get('total_templates')}")
    
    if result.get("errors"):
        print(f"Errors: {result['errors'][:3]}")
    
    # Statistics
    print("\n--- Statistics ---")
    stats = get_statistics()
    print(f"Templates: {stats.get('template_count')}")
    if stats.get("patterns"):
        print(f"Avg bullets/slide: {stats['patterns'].get('avg_bullets_per_slide')}")
        print(f"Avg slides/deck: {stats['patterns'].get('avg_slides_per_deck')}")
    
    # Suggestion
    print("\n--- Structure Suggestion (medium) ---")
    suggestion = suggest_structure("medium", topic="Digital Transformation")
    print(f"Slides: {suggestion.get('slide_count')}")
    for s in suggestion.get("slides", [])[:5]:
        print(f"  - {s['type']}: {s['title']}")
    print("  ...")
