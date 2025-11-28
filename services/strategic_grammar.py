from __future__ import annotations
from typing import Any, Dict, List
from pathlib import Path
from datetime import datetime
import hashlib
from pptx import Presentation

def _hash_path(p: Path) -> str:
    return hashlib.sha1(str(p.resolve()).encode("utf-8")).hexdigest()[:12]

def _shape_text(sh) -> str:
    try:
        if hasattr(sh, "text"): return (sh.text or "").strip()
        if hasattr(sh, "title") and sh.title: return (sh.title.text or "").strip()
    except Exception: return ""
    return ""

def _guess_slide_kind(texts: List[str]) -> str:
    t = " ".join(texts).lower()
    if any(k in t for k in ["agenda", "inhalt", "overview"]): return "agenda"
    if "kpi" in t or "metr" in t: return "kpis"
    if any(k in t for k in ["ziel", "objective", "okr"]): return "goals"
    if any(k in t for k in ["strategie", "strategy", "roadmap"]): return "strategy"
    if any(k in t for k in ["insight", "hypoth", "research"]): return "insights"
    return "section"

def extract_pptx_features(pptx_path: str | Path) -> Dict[str, Any]:
    p = Path(pptx_path)
    prs = Presentation(str(p))
    slides_out = []
    for idx, slide in enumerate(prs.slides, start=1):
        titles, bullets, bodies = [], [], []
        img, tab, chart = 0, 0, 0
        for sh in slide.shapes:
            try:
                if getattr(sh, "shape_type", None) and str(sh.shape_type).endswith("PICTURE"): img += 1
            except Exception: pass
            if getattr(sh, "has_table", False): tab += 1
            if getattr(sh, "has_chart", False): chart += 1

            txt = _shape_text(sh)
            if not txt: continue
            if getattr(sh, "name", "").lower().startswith("title"):
                titles.append(txt)
            else:
                bodies.append(txt)
            try:
                tf = getattr(sh, "text_frame", None)
                if tf:
                    for para in tf.paragraphs:
                        bt = (para.text or "").strip()
                        if bt:
                            bullets.append({"text": bt, "level": int(getattr(para, "level", 0))})
            except Exception: pass

        flat = titles + [b["text"] for b in bullets] + bodies
        slides_out.append({
            "index": idx,
            "title": titles[0] if titles else None,
            "bullets": bullets,
            "body": bodies,
            "kind": _guess_slide_kind(flat),
            "images": img, "tables": tab, "charts": chart
        })

    deck_title = slides_out[0].get("title") if slides_out else None
    deck_sub = None
    for s in slides_out[:3]:
        if s["body"]:
            deck_sub = s["body"][0]; break

    return {
        "source": str(p),
        "source_id": _hash_path(p),
        "created_at": int(datetime.utcnow().timestamp()),
        "deck": {"title": deck_title, "subtitle": deck_sub, "slide_count": len(slides_out)},
        "slides": slides_out,
        "metrics": {
            "avg_bullets_per_slide": round(sum(len(s["bullets"]) for s in slides_out)/max(1,len(slides_out)), 2),
            "image_ratio": round(sum(s["images"] for s in slides_out)/max(1,len(slides_out)), 2),
        }
    }

def _recommend_flow(seq: List[str]) -> List[str]:
    base = ["title","agenda","goals","strategy","kpis","roadmap"]
    if "agenda" not in seq and "title" in base: base.insert(1, "agenda")
    return base

def derive_patterns(features: Dict[str, Any]) -> Dict[str, Any]:
    seq = [s["kind"] for s in features.get("slides", [])]
    kind_hist, lvl_hist = {}, {}
    for k in seq: kind_hist[k] = kind_hist.get(k, 0) + 1
    for s in features.get("slides", []):
        for b in s["bullets"]:
            lvl = int(b.get("level", 0))
            lvl_hist[lvl] = lvl_hist.get(lvl, 0) + 1
    return {
        "source_id": features.get("source_id"),
        "deck_title": (features.get("deck") or {}).get("title"),
        "sequence": seq,
        "kind_histogram": kind_hist,
        "bullet_level_histogram": lvl_hist,
        "recommended_flow": _recommend_flow(seq),
    }
