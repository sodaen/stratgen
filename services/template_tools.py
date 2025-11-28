from __future__ import annotations
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Tuple

STYLES_DIR = Path("styles")
MASTER_PATH = STYLES_DIR / "master.pptx"
BACKUP_DIR  = STYLES_DIR / "backups"
BACKUP_DIR.mkdir(parents=True, exist_ok=True)

TOKENS = ["#KUNDENLOGO","#TITEL","#THEMA","#DATUM",
          "#ZWISCHENFOLIE",
          "#TOPIC","#SUBTOPIC","#INHALT",
          "#ENDE"]

REQUIRED = {
    "start":   {"#KUNDENLOGO","#TITEL","#THEMA","#DATUM"},
    "section": {"#ZWISCHENFOLIE"},
    "content": {"#TOPIC","#SUBTOPIC","#INHALT"},
    "end":     {"#ENDE"},
}

def _layout_tokens(layout) -> List[str]:
    found = set()
    for shp in getattr(layout, "shapes", []):
        if hasattr(shp, "text_frame") and shp.text_frame and shp.text_frame.text:
            txt = shp.text_frame.text
            for t in TOKENS:
                if t in txt: found.add(t)
    return sorted(found)

def validate_master(path: Path = MASTER_PATH) -> Dict:
    assert path.exists(), f"{path} fehlt"
    prs = Presentation(path.as_posix())
    layouts = []
    for i, lo in enumerate(prs.slide_layouts):
        layouts.append({"index": i, "tokens": _layout_tokens(lo)})
    # Rollen erkennen
    indices = detect_layout_indices(prs)
    # Vollständigkeit je Rolle prüfen
    roles_ok = {}
    for role, idx in indices.items():
        tokens = set(_layout_tokens(prs.slide_layouts[idx]))
        roles_ok[role] = REQUIRED[role].issubset(tokens)
    return {"path": str(path), "layouts": layouts, "roles": indices, "roles_ok": roles_ok}

def detect_layout_indices(prs: Presentation) -> Dict[str,int]:
    def has(layout, needed): return set(_layout_tokens(layout)).issuperset(set(needed))
    idx = {"start": None, "section": None, "content": None, "end": None}
    for i, lo in enumerate(prs.slide_layouts):
        if idx["start"]   is None and has(lo, REQUIRED["start"]):   idx["start"] = i
        if idx["section"] is None and has(lo, REQUIRED["section"]): idx["section"] = i
        if idx["content"] is None and has(lo, REQUIRED["content"]): idx["content"] = i
        if idx["end"]     is None and has(lo, REQUIRED["end"]):     idx["end"] = i
    # Fallbacks auf typische Office-Layouts
    if idx["start"]   is None: idx["start"]   = 0
    if idx["section"] is None: idx["section"] = 2 if len(prs.slide_layouts)>2 else 0
    if idx["content"] is None: idx["content"] = 1 if len(prs.slide_layouts)>1 else 0
    if idx["end"]     is None: idx["end"]     = 5 if len(prs.slide_layouts)>5 else len(prs.slide_layouts)-1
    return idx

def _title_placeholder(layout):
    for shp in layout.placeholders:
        try:
            if shp.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                if hasattr(shp, "text_frame"):
                    return shp
        except Exception: continue
    for shp in layout.shapes:
        if hasattr(shp, "text_frame"):
            return shp
    return None

def _body_placeholder(layout, title_shape):
    # bevorzugt BODY
    for shp in layout.placeholders:
        try:
            if shp is title_shape: 
                continue
            if shp.placeholder_format.type == PP_PLACEHOLDER.BODY and hasattr(shp, "text_frame"):
                return shp
        except Exception: continue
    # sonst: erster beliebiger Text-Shape, der nicht der Titel ist
    for shp in layout.shapes:
        if shp is title_shape: 
            continue
        if hasattr(shp, "text_frame"):
            return shp
    return None

def propose_patch(path: Path = MASTER_PATH) -> Dict:
    prs = Presentation(path.as_posix())
    roles = detect_layout_indices(prs)
    plan = {}
    # Für jede Rolle prüfen, was fehlt
    for role, required in REQUIRED.items():
        idx = roles[role]
        layout = prs.slide_layouts[idx]
        existing = set(_layout_tokens(layout))
        missing  = list(required - existing)
        plan[role] = {"layout_index": idx, "existing": sorted(existing), "missing": sorted(missing)}
    return {"path": str(path), "plan": plan}

def apply_patch(path: Path = MASTER_PATH) -> Dict:
    prs = Presentation(path.as_posix())
    roles = detect_layout_indices(prs)
    changes = {}
    for role, req in REQUIRED.items():
        idx = roles[role]
        lo  = prs.slide_layouts[idx]
        have = set(_layout_tokens(lo))
        need = list(req - have)
        if not need:
            changes[role] = {"layout_index": idx, "status": "ok", "applied": []}
            continue
        title = _title_placeholder(lo)
        body  = _body_placeholder(lo, title)
        applied = []
        if role == "start":
            if title: title.text_frame.text = "#TITEL"; applied.append("#TITEL")
            if body:  body.text_frame.text = "#KUNDENLOGO\n#THEMA\n#DATUM"; applied += ["#KUNDENLOGO","#THEMA","#DATUM"]
        elif role == "section":
            if title: title.text_frame.text = "#ZWISCHENFOLIE"; applied.append("#ZWISCHENFOLIE")
        elif role == "content":
            if title: title.text_frame.text = "#TOPIC"; applied.append("#TOPIC")
            if body:  body.text_frame.text = "#SUBTOPIC\n#INHALT"; applied += ["#SUBTOPIC","#INHALT"]
        elif role == "end":
            if title: title.text_frame.text = "#ENDE"; applied.append("#ENDE")
        changes[role] = {"layout_index": idx, "status": "patched", "applied": applied, "was_missing": need}
    # Backup & Save
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    backup_path = BACKUP_DIR / f"master_{ts}.pptx"
    backup_path.write_bytes(Path(path).read_bytes())
    prs.save(path.as_posix())
    return {"path": str(path), "backup": str(backup_path), "changes": changes}
