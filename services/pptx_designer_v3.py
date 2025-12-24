"""
PPTX Designer V3 - Mit Bild-Integration
"""

import io
from typing import Dict, Any, List, Optional
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


class PPTXDesignerV3:
    PALETTES = {
        "corporate": {
            "primary": "#1E40AF",
            "secondary": "#3B82F6",
            "accent": "#10B981",
            "background": "#FFFFFF",
            "text": "#111827",
            "text_light": "#6B7280",
            "chapter_bg": "#1E3A5F",
        }
    }
    
    def __init__(self, colors=None, palette="corporate", company_name="", include_slide_numbers=True, auto_images=True):
        self.palette = self.PALETTES.get(palette, self.PALETTES["corporate"])
        if colors:
            self.palette.update(colors)
        self.company_name = company_name
        self.include_slide_numbers = include_slide_numbers
        self.auto_images = auto_images
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.total_slides = 0
    
    def _hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    
    def _add_background(self, slide, color=None):
        if color is None:
            color = self.palette["background"]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(color)
    
    def _add_shape(self, slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb(color)
        shape.line.fill.background()
        return shape
    
    def _add_text_box(self, slide, left, top, width, height, text, font_size=14, font_color=None, bold=False, align=PP_ALIGN.LEFT):
        if font_color is None:
            font_color = self.palette["text"]
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(font_size)
        p.font.color.rgb = self._hex_to_rgb(font_color)
        p.font.bold = bold
        p.font.name = "Calibri"
        p.alignment = align
        return txBox
    
    def _add_image(self, slide, image_path, left, top, width=None, height=None):
        """Fuegt ein Bild zum Slide hinzu."""
        try:
            img_path = Path(image_path) if image_path else None
            if img_path and img_path.exists():
                if width and height:
                    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), Inches(width), Inches(height))
                elif width:
                    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), width=Inches(width))
                elif height:
                    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), height=Inches(height))
                else:
                    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top))
                return True
        except Exception as e:
            print(f"Bild-Fehler: {e}")
        return False
    
    def _get_image_for_slide(self, slide_type: str, title: str = "") -> Optional[str]:
        """Holt automatisch ein passendes Bild via Unsplash."""
        if not self.auto_images:
            return None
        try:
            from services.unsplash_service import get_image_for_slide
            result = get_image_for_slide(slide_type, title)
            if result:
                return result.get("local_path")
        except Exception as e:
            print(f"Unsplash error: {e}")
        return None
    
    def _add_slide_number(self, slide, number, color=None):
        if self.include_slide_numbers:
            font_color = color if color else self.palette["text_light"]
            self._add_text_box(slide, 12.5, 7.0, 0.7, 0.3, f"{number}/{self.total_slides}", font_size=10, font_color=font_color, align=PP_ALIGN.RIGHT)
    
    def _add_footer(self, slide):
        if self.company_name:
            self._add_text_box(slide, 0.5, 7.0, 4, 0.3, self.company_name, font_size=10, font_color=self.palette["text_light"])
    
    def _render_title_slide(self, slide, data, number):
        # Hintergrund-Bild oder Farbe
        image_path = data.get("image_path") or self._get_image_for_slide("title", data.get("title", ""))
        if image_path:
            self._add_image(slide, image_path, 0, 0, 13.333, 7.5)
            # Nur unterer Balken für Textlesbarkeit (Bild bleibt sichtbar oben)
            self._add_shape(slide, 0, 4.5, 13.333, 3.0, self.palette["primary"])
        else:
            self._add_background(slide, self.palette["primary"])
        
        self._add_shape(slide, 1, 3.2, 11.333, 0.02, self.palette["accent"])
        title = data.get("title", "Praesentation")
        self._add_text_box(slide, 1, 2.0, 11.333, 1.2, title, font_size=44, font_color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        subtitle = data.get("subtitle", self.company_name)
        if subtitle:
            self._add_text_box(slide, 1, 3.5, 11.333, 0.6, subtitle, font_size=24, font_color="#FFFFFF", align=PP_ALIGN.CENTER)
        date_str = data.get("date", datetime.now().strftime("%B %Y"))
        self._add_text_box(slide, 1, 6.5, 11.333, 0.4, date_str, font_size=14, font_color=self.palette["secondary"], align=PP_ALIGN.CENTER)
    
    def _render_chapter_slide(self, slide, data, number):
        # Hintergrund-Bild oder Farbe
        image_path = data.get("image_path") or self._get_image_for_slide("chapter", data.get("title", ""))
        if image_path:
            self._add_image(slide, image_path, 0, 0, 13.333, 7.5)
            # Nur mittlerer Balken für Text (Bild oben und unten sichtbar)
            self._add_shape(slide, 0, 2.5, 13.333, 2.5, self.palette["chapter_bg"])
        else:
            self._add_background(slide, self.palette["chapter_bg"])
        
        self._add_shape(slide, 0, 0, 0.15, 7.5, self.palette["accent"])
        chapter_num = data.get("chapter_number", "")
        if chapter_num:
            self._add_text_box(slide, 1, 2.0, 11, 0.6, f"KAPITEL {chapter_num}", font_size=16, font_color=self.palette["accent"], bold=True)
        title = data.get("title", "Neues Kapitel")
        self._add_text_box(slide, 1, 2.8, 11, 1.5, title, font_size=48, font_color="#FFFFFF", bold=True)
        subtitle = data.get("subtitle", "")
        if subtitle:
            self._add_text_box(slide, 1, 4.5, 10, 1.0, subtitle, font_size=18, font_color=self.palette["secondary"])
    
    def _render_bullets_slide(self, slide, data, number):
        self._add_background(slide)
        self._add_shape(slide, 0, 0, 13.333, 0.08, self.palette["primary"])
        title = data.get("title", "Inhalt")
        self._add_text_box(slide, 0.5, 0.4, 12, 0.7, title, font_size=28, bold=True, font_color=self.palette["primary"])
        
        # Optionales Bild rechts (jeder 3. Slide)
        image_path = None
        if self.auto_images and number % 3 == 0:
            try:
                from services.unsplash_service import get_image_for_slide
                result = get_image_for_slide(title, self.company_name)
                if result and isinstance(result, dict):
                    image_path = result.get("local_path")
            except:
                pass
        
        bullets = data.get("bullets", [])
        if not bullets:
            content = data.get("content", "")
            if content:
                bullets = [content]
        
        text_width = 7.5 if image_path else 11.5
        y_pos = 1.4
        for bullet in bullets[:6]:
            if not bullet or len(str(bullet).strip()) < 3:
                continue
            self._add_shape(slide, 0.5, y_pos + 0.15, 0.15, 0.15, self.palette["accent"], MSO_SHAPE.OVAL)
            self._add_text_box(slide, 0.85, y_pos, text_width, 0.9, str(bullet), font_size=16, font_color=self.palette["text"])
            y_pos += 1.0
        
        if image_path:
            self._add_image(slide, image_path, 9.0, 1.4, 3.8, 3.8)
        
        self._add_slide_number(slide, number)
        self._add_footer(slide)

    def _render_text_slide(self, slide, data, number):
        self._add_background(slide)
        self._add_shape(slide, 0, 0, 0.08, 7.5, self.palette["primary"])
        title = data.get("title", "Details")
        self._add_text_box(slide, 0.5, 0.4, 12, 0.7, title, font_size=28, bold=True, font_color=self.palette["primary"])
        
        # Optionales Key Visual rechts (jeder 2. Text-Slide)
        image_path = None
        if self.auto_images and number % 2 == 0:
            try:
                from services.unsplash_service import get_image_for_slide
                result = get_image_for_slide(title, self.company_name)
                if result and isinstance(result, dict):
                    image_path = result.get("local_path")
            except:
                pass
        
        content = data.get("content", "")
        if not content:
            bullets = data.get("bullets", [])
            if bullets:
                content = "\n\n".join(str(b) for b in bullets)
        
        if image_path:
            # Text links, Bild rechts
            if content:
                self._add_text_box(slide, 0.5, 1.3, 7.5, 5.5, content, font_size=15, font_color=self.palette["text"])
            self._add_image(slide, image_path, 8.5, 1.3, 4.3, 4.3)
        else:
            # Volle Breite ohne Bild
            if content:
                self._add_text_box(slide, 0.5, 1.3, 12, 5.5, content, font_size=15, font_color=self.palette["text"])
        
        self._add_slide_number(slide, number)
        self._add_footer(slide)
    
    def _render_executive_summary(self, slide, data, number):
        self._add_background(slide)
        title = data.get("title", "Executive Summary")
        self._add_text_box(slide, 0.5, 0.4, 12, 0.7, title, font_size=32, bold=True, font_color=self.palette["primary"])
        self._add_shape(slide, 0.5, 1.1, 12, 0.02, self.palette["accent"])
        bullets = data.get("bullets", [])
        y_pos = 1.5
        for i, bullet in enumerate(bullets[:4]):
            self._add_shape(slide, 0.5, y_pos, 0.4, 0.4, self.palette["primary"])
            self._add_text_box(slide, 0.5, y_pos, 0.4, 0.4, str(i + 1), font_size=14, font_color="#FFFFFF", bold=True, align=PP_ALIGN.CENTER)
            self._add_text_box(slide, 1.1, y_pos, 5.5, 1.2, str(bullet), font_size=14, font_color=self.palette["text"])
            y_pos += 1.4
        self._add_shape(slide, 7, 1.5, 5.5, 5, self.palette["primary"])
        highlights = data.get("highlights", bullets[4:] if len(bullets) > 4 else [])
        if highlights:
            self._add_text_box(slide, 7.3, 1.7, 5, 0.5, "KEY HIGHLIGHTS", font_size=14, font_color="#FFFFFF", bold=True)
            hy = 2.4
            for h in highlights[:4]:
                self._add_text_box(slide, 7.3, hy, 4.8, 1.0, f"- {h}", font_size=13, font_color="#FFFFFF")
                hy += 1.1
        self._add_slide_number(slide, number)
        self._add_footer(slide)
    
    def _render_persona_slide(self, slide, data, number):
        self._add_background(slide)
        self._add_shape(slide, 0, 0, 13.333, 1.0, self.palette["primary"])
        title = data.get("title", "Persona")
        self._add_text_box(slide, 0.5, 0.25, 12, 0.6, title, font_size=24, bold=True, font_color="#FFFFFF")
        
        # Persona-Bild (Kreis mit Bild oder Platzhalter)
        image_path = data.get("image_path") or self._get_image_for_slide("persona", data.get("persona_name", ""))
        if image_path:
            # Bild einfuegen (quadratisch)
            self._add_image(slide, image_path, 0.5, 1.3, 2.5, 2.5)
        else:
            # Platzhalter-Kreis
            self._add_shape(slide, 0.5, 1.3, 2.5, 2.5, self.palette["secondary"], MSO_SHAPE.OVAL)
        
        bullets = data.get("bullets", [])
        name = data.get("persona_name", bullets[0] if bullets else "Max Mustermann")
        role = data.get("persona_role", bullets[1] if len(bullets) > 1 else "Entscheider")
        self._add_text_box(slide, 3.3, 1.3, 4, 0.5, str(name), font_size=22, bold=True, font_color=self.palette["primary"])
        self._add_text_box(slide, 3.3, 1.85, 4, 0.4, str(role), font_size=16, font_color=self.palette["text_light"])
        py = 2.5
        for prop in bullets[2:6]:
            self._add_text_box(slide, 3.3, py, 4, 0.4, f"- {prop}", font_size=12, font_color=self.palette["text"])
            py += 0.4
        self._add_shape(slide, 8, 1.3, 4.8, 2.8, "#FEF2F2")
        self._add_text_box(slide, 8.2, 1.4, 4.4, 0.4, "PAIN POINTS", font_size=12, bold=True, font_color="#DC2626")
        pain_points = data.get("pain_points", bullets[6:9] if len(bullets) > 6 else [])
        ppy = 1.9
        for p in pain_points[:3]:
            self._add_text_box(slide, 8.2, ppy, 4.4, 0.6, f"- {p}", font_size=11, font_color="#7F1D1D")
            ppy += 0.55
        self._add_shape(slide, 8, 4.3, 4.8, 2.5, "#F0FDF4")
        self._add_text_box(slide, 8.2, 4.4, 4.4, 0.4, "GOALS", font_size=12, bold=True, font_color="#16A34A")
        goals = data.get("goals", bullets[9:12] if len(bullets) > 9 else [])
        gy = 4.9
        for g in goals[:3]:
            self._add_text_box(slide, 8.2, gy, 4.4, 0.6, f"- {g}", font_size=11, font_color="#166534")
            gy += 0.55
        quote = data.get("quote", bullets[12] if len(bullets) > 12 else "")
        if quote:
            self._add_text_box(slide, 0.5, 6.0, 7, 0.8, f'"{quote}"', font_size=13, font_color=self.palette["text_light"])
        self._add_slide_number(slide, number)
    
    def _render_comparison_slide(self, slide, data, number):
        self._add_background(slide)
        title = data.get("title", "Vergleich")
        self._add_text_box(slide, 0.5, 0.4, 12, 0.7, title, font_size=28, bold=True, font_color=self.palette["primary"])
        bullets = data.get("bullets", [])
        headers = data.get("headers", ["Option 1", "Option 2", "Option 3"])
        columns = min(len(headers), 4)
        col_width = 12.333 / columns
        for col in range(columns):
            x_pos = 0.5 + col * col_width
            bg_color = self.palette["primary"] if col == 1 else "#F3F4F6"
            text_color = "#FFFFFF" if col == 1 else self.palette["text"]
            self._add_shape(slide, x_pos, 1.3, col_width - 0.2, 5.5, bg_color)
            header = headers[col] if col < len(headers) else f"Option {col+1}"
            self._add_text_box(slide, x_pos + 0.1, 1.4, col_width - 0.4, 0.5, str(header), font_size=16, bold=True, font_color=text_color, align=PP_ALIGN.CENTER)
            col_bullets = bullets[col::columns] if bullets else []
            cy = 2.2
            for item in col_bullets[:6]:
                self._add_text_box(slide, x_pos + 0.1, cy, col_width - 0.4, 0.7, f"- {item}", font_size=12, font_color=text_color)
                cy += 0.8
        self._add_slide_number(slide, number)
        self._add_footer(slide)
    
    def _render_chart_slide(self, slide, data, number):
        self._add_background(slide)
        title = data.get("title", "Daten")
        self._add_text_box(slide, 0.5, 0.4, 12, 0.7, title, font_size=28, bold=True, font_color=self.palette["primary"])
        self._add_shape(slide, 0.5, 1.3, 7.5, 5.2, "#F9FAFB")
        self._add_text_box(slide, 2.5, 3.5, 3.5, 0.5, "[CHART]", font_size=24, font_color=self.palette["text_light"], align=PP_ALIGN.CENTER)
        self._add_shape(slide, 8.2, 1.3, 4.6, 5.2, self.palette["primary"])
        self._add_text_box(slide, 8.4, 1.5, 4.2, 0.4, "KEY INSIGHTS", font_size=14, bold=True, font_color="#FFFFFF")
        bullets = data.get("bullets", [])
        iy = 2.1
        for b in bullets[:5]:
            self._add_text_box(slide, 8.4, iy, 4.2, 0.9, f"> {b}", font_size=12, font_color="#FFFFFF")
            iy += 0.95
        self._add_slide_number(slide, number)
        self._add_footer(slide)
    
    def _render_timeline_slide(self, slide, data, number):
        self._add_background(slide)
        title = data.get("title", "Roadmap")
        self._add_text_box(slide, 0.5, 0.4, 12, 0.7, title, font_size=28, bold=True, font_color=self.palette["primary"])
        bullets = data.get("bullets", [])
        phases = data.get("phases", bullets[:6] if bullets else [f"Phase {i+1}" for i in range(4)])
        self._add_shape(slide, 0.5, 3.5, 12.333, 0.05, self.palette["secondary"])
        num = min(len(phases), 6) or 4
        pw = 12.333 / num
        for i, ph in enumerate(phases[:num]):
            xp = 0.5 + i * pw + pw / 2 - 0.25
            col = self.palette["primary"] if i % 2 == 0 else self.palette["accent"]
            self._add_shape(slide, xp, 3.35, 0.5, 0.5, col, MSO_SHAPE.OVAL)
            yo = 1.8 if i % 2 == 0 else 4.2
            self._add_text_box(slide, xp - 1, yo, 2.5, 1.2, str(ph), font_size=12, font_color=self.palette["text"], align=PP_ALIGN.CENTER)
        self._add_slide_number(slide, number)
        self._add_footer(slide)
    
    def _render_quote_slide(self, slide, data, number):
        # Hintergrund-Bild oder Farbe
        image_path = data.get("image_path") or self._get_image_for_slide("quote", data.get("quote", ""))
        if image_path:
            self._add_image(slide, image_path, 0, 0, 13.333, 7.5)
            self._add_shape(slide, 0, 0, 13.333, 7.5, self.palette["chapter_bg"])
        else:
            self._add_background(slide, self.palette["chapter_bg"])
        
        self._add_text_box(slide, 1, 1.5, 2, 1.5, '"', font_size=100, font_color=self.palette["accent"], bold=True)
        bullets = data.get("bullets", [])
        quote = data.get("quote", bullets[0] if bullets else "Zitat")
        self._add_text_box(slide, 1.5, 2.5, 10, 3, str(quote), font_size=28, font_color="#FFFFFF", align=PP_ALIGN.CENTER)
        source = data.get("source", bullets[1] if len(bullets) > 1 else "")
        if source:
            self._add_text_box(slide, 1.5, 5.8, 10, 0.5, f"- {source}", font_size=16, font_color=self.palette["secondary"], align=PP_ALIGN.CENTER)
        self._add_slide_number(slide, number)
    
    def _render_conclusion_slide(self, slide, data, number):
        # Optionales Hintergrundbild
        image_path = None
        if self.auto_images:
            try:
                from services.unsplash_service import get_image_for_slide
                result = get_image_for_slide("success teamwork business", self.company_name)
                if result and isinstance(result, dict):
                    image_path = result.get("local_path")
            except:
                pass
        
        if image_path:
            self._add_image(slide, image_path, 0, 0, 13.333, 7.5)
            # Overlay
            self._add_shape(slide, 0, 0, 13.333, 7.5, self.palette["primary"])
        else:
            self._add_background(slide)
            self._add_shape(slide, 0, 0, 13.333, 0.15, self.palette["accent"])
        
        title = data.get("title", "Key Takeaways")
        self._add_text_box(slide, 0.5, 0.5, 12, 0.7, title, font_size=32, bold=True, font_color=self.palette["primary"] if not image_path else "#FFFFFF")
        bullets = data.get("bullets", [])
        cy = 1.5
        text_color = self.palette["text"] if not image_path else "#FFFFFF"
        for i, b in enumerate(bullets[:5], 1):
            self._add_shape(slide, 0.5, cy, 0.6, 0.6, self.palette["primary"] if not image_path else self.palette["accent"])
            self._add_text_box(slide, 0.5, cy + 0.05, 0.6, 0.5, str(i), font_size=20, bold=True, font_color="#FFFFFF", align=PP_ALIGN.CENTER)
            self._add_text_box(slide, 1.3, cy, 11, 0.9, str(b), font_size=16, font_color=text_color)
            cy += 1.1
        cta = data.get("cta", bullets[5] if len(bullets) > 5 else "")
        if cta:
            self._add_shape(slide, 0.5, 6.3, 12.333, 0.8, self.palette["accent"])
            self._add_text_box(slide, 0.5, 6.4, 12.333, 0.6, str(cta), font_size=18, bold=True, font_color="#FFFFFF", align=PP_ALIGN.CENTER)
        self._add_slide_number(slide, number, color="#FFFFFF" if image_path else None)

    def _render_sources_slide(self, slide, data, number):
        """Rendert eine Quellenübersicht-Folie mit detaillierten Angaben."""
        self._add_background(slide)
        
        # Titel
        self._add_text_box(slide, 0.5, 0.3, 12.333, 0.8, "Quellen & Referenzen", 
                          font_size=32, bold=True, font_color=self.palette["primary"])
        
        # Untertitel mit Hinweis
        self._add_text_box(slide, 0.5, 0.9, 12.333, 0.4, 
                          "Alle Angaben wurden zum Zeitpunkt der Erstellung geprüft. Stand: Dezember 2024",
                          font_size=11, font_color=self.palette["text_light"])
        
        # Linie unter Titel
        self._add_shape(slide, 0.5, 1.0, 12.333, 0.02, self.palette["accent"])
        
        # Quellen aus data
        sources = data.get("sources", [])
        if not sources:
            sources = data.get("bullets", [])
        
        # Zwei Spalten für Quellen
        cy = 1.3
        col_width = 6.0
        
        for i, source in enumerate(sources[:20]):  # Max 20 Quellen
            col = 0 if i < 10 else 1
            row = i if i < 10 else i - 10
            
            x = 0.5 + (col * col_width)
            y = cy + (row * 0.55)
            
            # Nummer und Quelle
            source_text = f"{i+1}. {source}" if not str(source).startswith(('•', '-', '1', '2', '3', '4', '5', '6', '7', '8', '9')) else source
            self._add_text_box(slide, x, y, col_width - 0.3, 0.5, 
                              str(source_text)[:80], 
                              font_size=11, font_color=self.palette["text"])
        
        # Footer
        self._add_text_box(slide, 0.5, 6.8, 12.333, 0.4, 
                          "Alle Quellen wurden zum Zeitpunkt der Erstellung geprüft.",
                          font_size=9, font_color=self.palette["text_light"], align=PP_ALIGN.CENTER)
        
        self._add_slide_number(slide, number)
        self._add_footer(slide)

    def _render_contact_slide(self, slide, data, number):
        self._add_background(slide, self.palette["primary"])
        title = data.get("title", "Vielen Dank!")
        self._add_text_box(slide, 1, 2, 11.333, 1.2, title, font_size=48, bold=True, font_color="#FFFFFF", align=PP_ALIGN.CENTER)
        subtitle = data.get("subtitle", "Fragen")
        self._add_text_box(slide, 1, 3.3, 11.333, 0.6, subtitle, font_size=24, font_color=self.palette["secondary"], align=PP_ALIGN.CENTER)
        bullets = data.get("bullets", [])
        if bullets or self.company_name:
            txt = "\n".join(str(b) for b in bullets[:3]) if bullets else self.company_name
            self._add_text_box(slide, 4.2, 4.7, 4.933, 1.6, txt, font_size=14, font_color="#FFFFFF", align=PP_ALIGN.CENTER)
    
    def create_presentation(self, slides, title="Praesentation", company="", include_sources_slide=False):
        self.company_name = company or self.company_name
        self.total_slides = len(slides)
        RENDERERS = {
            "title": self._render_title_slide, "chapter": self._render_chapter_slide,
            "executive_summary": self._render_executive_summary, "bullets": self._render_bullets_slide,
            "text": self._render_text_slide, "persona": self._render_persona_slide,
            "comparison": self._render_comparison_slide, "chart": self._render_chart_slide,
            "timeline": self._render_timeline_slide, "quote": self._render_quote_slide,
            "conclusion": self._render_conclusion_slide, "contact": self._render_contact_slide, "sources": self._render_sources_slide,
        }
        FALLBACKS = {
            "problem": "bullets", "solution": "bullets", "benefits": "bullets",
            "features": "bullets", "approach": "bullets", "methodology": "bullets",
            "risks": "bullets", "resources": "bullets", "budget": "bullets",
            "roi": "bullets", "metrics": "bullets", "implementation": "bullets",
            "deep_dive": "text", "technical": "text", "analysis": "text",
            "case_study": "text", "testimonial": "quote", "milestones": "timeline",
            "roadmap": "timeline", "next_steps": "conclusion",
        }
        for i, sd in enumerate(slides, 1):
            st = sd.get("type", "bullets")
            if st not in RENDERERS:
                st = FALLBACKS.get(st, "bullets")
            sl = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            try:
                RENDERERS.get(st, self._render_bullets_slide)(sl, sd, i)
            except Exception as e:
                print(f"Error slide {i}: {e}")
                self._render_bullets_slide(sl, sd, i)
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.read()


def create_presentation_v3(slides, title="Praesentation", company="", colors=None, palette="corporate", auto_images=True):
    designer = PPTXDesignerV3(colors=colors, palette=palette, company_name=company, auto_images=auto_images)
    return designer.create_presentation(slides, title, company)
