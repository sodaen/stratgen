"""
STRATGEN Knowledge Pipeline - Phase 2
Chunking Engine mit Quality Gates

Speichere als: /home/sodaen/stratgen/services/knowledge_pipeline.py
"""

import os
import re
import json
import hashlib
import logging
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
from dataclasses import dataclass, field, asdict
import tiktoken

logger = logging.getLogger(__name__)

# Konfiguration
CHUNK_CONFIG = {
    "knowledge": {"size": 600, "overlap": 100, "min_length": 50},
    "template": {"size": 800, "overlap": 150, "min_length": 50},
    "external": {"size": 400, "overlap": 50, "min_length": 80},
    "generated": {"size": 500, "overlap": 100, "min_length": 100}
}

DATA_ROOT = Path(os.getenv("STRATGEN_DATA", "/home/sodaen/stratgen/data"))


@dataclass
class ChunkMetadata:
    """Vollständige Metadata für jeden Chunk."""
    id: str
    source_file: str
    source_path: str
    source_type: str  # knowledge, template, external, generated
    collection: str
    
    chunk_index: int
    total_chunks: int
    chunk_size_tokens: int
    chunk_size_chars: int
    
    content_hash: str
    content_preview: str
    
    quality_score: float
    quality_checks: Dict[str, bool]
    
    indexed_at: str
    source_modified: str
    
    embedding_model: str = "nomic-embed-text"
    embedding_dimension: int = 768


@dataclass
class QualityReport:
    """Bericht über Quality-Checks."""
    passed: bool
    min_length_ok: bool
    no_spam: bool
    valid_utf8: bool
    language_ok: bool
    duplicate_of: Optional[str] = None
    issues: List[str] = field(default_factory=list)


class ChunkingEngine:
    """
    Intelligente Chunking Engine mit Quality Gates.
    """
    
    def __init__(self):
        self.tokenizer = tiktoken.get_encoding("cl100k_base")
        self.seen_hashes: Dict[str, str] = {}  # hash -> chunk_id
        self.stats = {
            "processed": 0,
            "chunks_created": 0,
            "rejected_min_length": 0,
            "rejected_spam": 0,
            "rejected_duplicate": 0
        }
    
    def count_tokens(self, text: str) -> int:
        """Zählt Tokens im Text."""
        return len(self.tokenizer.encode(text))
    
    def compute_hash(self, text: str) -> str:
        """Berechnet Content-Hash."""
        normalized = re.sub(r'\s+', ' ', text.strip().lower())
        return hashlib.sha256(normalized.encode()).hexdigest()[:16]
    
    def check_quality(self, text: str, min_length: int = 100) -> QualityReport:
        """
        Führt alle Quality-Checks durch.
        """
        issues = []
        
        # 1. Mindest-Länge
        min_length_ok = len(text.strip()) >= min_length
        if not min_length_ok:
            issues.append(f"Zu kurz: {len(text)} < {min_length} Zeichen")
        
        # 2. Spam/Junk-Erkennung
        spam_patterns = [
            r'^[\d\s\.\,]+$',  # Nur Zahlen
            r'^[\W\s]+$',      # Nur Sonderzeichen
            r'(.)\1{10,}',     # Wiederholte Zeichen
        ]
        no_spam = True
        for pattern in spam_patterns:
            if re.match(pattern, text):
                no_spam = False
                issues.append(f"Spam-Pattern erkannt: {pattern}")
                break
        
        # Sonderzeichen-Ratio
        special_ratio = len(re.findall(r'[^\w\s]', text)) / max(len(text), 1)
        if special_ratio > 0.3:
            no_spam = False
            issues.append(f"Zu viele Sonderzeichen: {special_ratio:.1%}")
        
        # 3. UTF-8 Validierung
        try:
            text.encode('utf-8').decode('utf-8')
            valid_utf8 = True
        except:
            valid_utf8 = False
            issues.append("Ungültiges UTF-8")
        
        # 4. Sprache (einfache Heuristik)
        german_words = ['und', 'der', 'die', 'das', 'ist', 'für', 'mit', 'auf', 'bei']
        english_words = ['the', 'and', 'for', 'with', 'this', 'that', 'from', 'have']
        
        words = text.lower().split()
        de_count = sum(1 for w in words if w in german_words)
        en_count = sum(1 for w in words if w in english_words)
        
        language_ok = (de_count + en_count) >= min(3, len(words) // 10)
        if not language_ok and len(words) > 20:
            issues.append("Sprache nicht erkannt (weder DE noch EN)")
        
        # 5. Duplikat-Check
        content_hash = self.compute_hash(text)
        duplicate_of = self.seen_hashes.get(content_hash)
        if duplicate_of:
            issues.append(f"Duplikat von: {duplicate_of}")
        
        passed = min_length_ok and no_spam and valid_utf8 and language_ok and not duplicate_of
        
        return QualityReport(
            passed=passed,
            min_length_ok=min_length_ok,
            no_spam=no_spam,
            valid_utf8=valid_utf8,
            language_ok=language_ok,
            duplicate_of=duplicate_of,
            issues=issues
        )
    
    def chunk_text(self, 
                   text: str, 
                   source_type: str = "knowledge",
                   source_file: str = "",
                   source_path: str = "") -> List[Tuple[str, ChunkMetadata]]:
        """
        Chunked Text mit semantischem Bewusstsein.
        
        Returns:
            List of (chunk_text, metadata) tuples
        """
        config = CHUNK_CONFIG.get(source_type, CHUNK_CONFIG["knowledge"])
        target_size = config["size"]
        overlap = config["overlap"]
        min_length = config["min_length"]
        
        # Collection bestimmen
        collection_map = {
            "knowledge": "knowledge_base",
            "template": "design_templates",
            "external": "external_sources",
            "generated": "generated_outputs"
        }
        collection = collection_map.get(source_type, "knowledge_base")
        
        self.stats["processed"] += 1
        
        # Absätze respektieren
        paragraphs = re.split(r'\n\s*\n', text)
        
        chunks = []
        current_chunk = ""
        current_tokens = 0
        
        for para in paragraphs:
            para = para.strip()
            if not para:
                continue
            
            para_tokens = self.count_tokens(para)
            
            # Absatz passt in aktuellen Chunk
            if current_tokens + para_tokens <= target_size:
                current_chunk += ("\n\n" if current_chunk else "") + para
                current_tokens += para_tokens
            else:
                # Aktuellen Chunk speichern
                if current_chunk and current_tokens >= min_length // 4:
                    chunks.append(current_chunk)
                
                # Großer Absatz: aufteilen
                if para_tokens > target_size:
                    sentences = re.split(r'(?<=[.!?])\s+', para)
                    current_chunk = ""
                    current_tokens = 0
                    
                    for sent in sentences:
                        sent_tokens = self.count_tokens(sent)
                        if current_tokens + sent_tokens <= target_size:
                            current_chunk += (" " if current_chunk else "") + sent
                            current_tokens += sent_tokens
                        else:
                            if current_chunk:
                                chunks.append(current_chunk)
                            current_chunk = sent
                            current_tokens = sent_tokens
                else:
                    current_chunk = para
                    current_tokens = para_tokens
        
        # Letzten Chunk
        if current_chunk:
            chunks.append(current_chunk)
        
        # Overlap hinzufügen
        if overlap > 0 and len(chunks) > 1:
            overlapped_chunks = []
            for i, chunk in enumerate(chunks):
                if i > 0:
                    # Overlap vom vorherigen Chunk
                    prev_words = chunks[i-1].split()
                    overlap_words = prev_words[-overlap//4:] if len(prev_words) > overlap//4 else []
                    chunk = " ".join(overlap_words) + " " + chunk if overlap_words else chunk
                overlapped_chunks.append(chunk)
            chunks = overlapped_chunks
        
        # Metadata erstellen und Quality-Check
        result = []
        total_chunks = len(chunks)
        source_modified = datetime.now().isoformat()
        
        if source_path:
            try:
                source_modified = datetime.fromtimestamp(
                    Path(source_path).stat().st_mtime
                ).isoformat()
            except:
                pass
        
        for i, chunk_text in enumerate(chunks):
            # Quality Check
            quality = self.check_quality(chunk_text, min_length)
            
            if not quality.passed:
                if not quality.min_length_ok:
                    self.stats["rejected_min_length"] += 1
                elif not quality.no_spam:
                    self.stats["rejected_spam"] += 1
                elif quality.duplicate_of:
                    self.stats["rejected_duplicate"] += 1
                logger.debug(f"Chunk rejected: {quality.issues}")
                continue
            
            # Hash registrieren
            content_hash = self.compute_hash(chunk_text)
            chunk_id = f"{content_hash}_{i}"
            self.seen_hashes[content_hash] = chunk_id
            
            metadata = ChunkMetadata(
                id=chunk_id,
                source_file=source_file or Path(source_path).name if source_path else "unknown",
                source_path=source_path,
                source_type=source_type,
                collection=collection,
                chunk_index=i,
                total_chunks=total_chunks,
                chunk_size_tokens=self.count_tokens(chunk_text),
                chunk_size_chars=len(chunk_text),
                content_hash=content_hash,
                content_preview=chunk_text[:200] + "..." if len(chunk_text) > 200 else chunk_text,
                quality_score=self._compute_quality_score(chunk_text, quality),
                quality_checks={
                    "min_length": quality.min_length_ok,
                    "no_spam": quality.no_spam,
                    "valid_utf8": quality.valid_utf8,
                    "language": quality.language_ok
                },
                indexed_at=datetime.now().isoformat(),
                source_modified=source_modified
            )
            
            result.append((chunk_text, metadata))
            self.stats["chunks_created"] += 1
        
        return result
    
    def _compute_quality_score(self, text: str, quality: QualityReport) -> float:
        """Berechnet Quality-Score (0-1)."""
        score = 0.5  # Basis
        
        # Länge (optimal 300-800 Zeichen)
        length = len(text)
        if 300 <= length <= 800:
            score += 0.2
        elif 200 <= length <= 1000:
            score += 0.1
        
        # Quality-Checks
        if quality.min_length_ok:
            score += 0.1
        if quality.no_spam:
            score += 0.1
        if quality.language_ok:
            score += 0.1
        
        return min(score, 1.0)
    
    def get_stats(self) -> Dict[str, int]:
        """Gibt Statistiken zurück."""
        return self.stats.copy()


class KnowledgePipeline:
    """
    Haupt-Pipeline für Knowledge Ingestion.
    """
    
    def __init__(self):
        self.chunker = ChunkingEngine()
        self.qdrant = None
        self._init_qdrant()
    
    def _init_qdrant(self):
        """Initialisiert Qdrant-Client."""
        try:
            from qdrant_client import QdrantClient
            self.qdrant = QdrantClient(host="localhost", port=6333)
        except Exception as e:
            logger.error(f"Qdrant init failed: {e}")
    
    def _get_embedding(self, text: str) -> List[float]:
        """Generiert Embedding via Ollama."""
        import httpx
        
        try:
            response = httpx.post(
                "http://localhost:11434/api/embeddings",
                json={"model": "nomic-embed-text", "prompt": text},
                timeout=30.0
            )
            if response.status_code == 200:
                return response.json().get("embedding", [])
        except Exception as e:
            logger.error(f"Embedding failed: {e}")
        
        return []
    
    def ingest_file(self, 
                    file_path: str,
                    source_type: str = "knowledge") -> Dict[str, Any]:
        """
        Indexiert eine einzelne Datei.
        """
        path = Path(file_path)
        
        if not path.exists():
            return {"ok": False, "error": f"File not found: {file_path}"}
        
        # Text extrahieren
        text = self._extract_text(path)
        if not text:
            return {"ok": False, "error": "Could not extract text"}
        
        # Chunken
        chunks = self.chunker.chunk_text(
            text=text,
            source_type=source_type,
            source_file=path.name,
            source_path=str(path)
        )
        
        if not chunks:
            return {"ok": False, "error": "No valid chunks created"}
        
        # In Qdrant indexieren
        if self.qdrant:
            from qdrant_client.http import models
            
            points = []
            for chunk_text, metadata in chunks:
                embedding = self._get_embedding(chunk_text)
                if not embedding:
                    continue
                
                points.append(models.PointStruct(
                    id=hash(metadata.id) & 0x7FFFFFFFFFFFFFFF,  # Positive int64
                    vector=embedding,
                    payload={
                        "text": chunk_text,
                        **asdict(metadata)
                    }
                ))
            
            if points:
                self.qdrant.upsert(
                    collection_name=metadata.collection,
                    points=points
                )
        
        return {
            "ok": True,
            "file": path.name,
            "chunks_created": len(chunks),
            "collection": chunks[0][1].collection if chunks else None
        }
    
    def _extract_text(self, path: Path) -> str:
        """Extrahiert Text aus verschiedenen Dateiformaten."""
        suffix = path.suffix.lower()
        
        try:
            if suffix in ['.txt', '.md']:
                return path.read_text(encoding='utf-8', errors='ignore')
            
            elif suffix == '.pdf':
                try:
                    import fitz  # PyMuPDF
                    doc = fitz.open(str(path))
                    text = ""
                    for page in doc:
                        text += page.get_text()
                    return text
                except:
                    return ""
            
            elif suffix == '.docx':
                try:
                    from docx import Document
                    doc = Document(str(path))
                    return "\n".join([p.text for p in doc.paragraphs])
                except:
                    return ""
            
            elif suffix == '.pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(str(path))
                    texts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                texts.append(shape.text)
                    return "\n\n".join(texts)
                except:
                    return ""
            
            else:
                # Versuche als Text zu lesen
                return path.read_text(encoding='utf-8', errors='ignore')
        
        except Exception as e:
            logger.error(f"Text extraction failed for {path}: {e}")
            return ""
    
    def ingest_directory(self, 
                         directory: str,
                         source_type: str = "knowledge",
                         extensions: List[str] = None) -> Dict[str, Any]:
        """
        Indexiert ein komplettes Verzeichnis.
        """
        dir_path = Path(directory)
        if not dir_path.exists():
            return {"ok": False, "error": f"Directory not found: {directory}"}
        
        extensions = extensions or ['.txt', '.md', '.pdf', '.docx', '.pptx']
        
        results = {
            "ok": True,
            "directory": str(dir_path),
            "files_processed": 0,
            "files_failed": 0,
            "chunks_total": 0,
            "details": []
        }
        
        for ext in extensions:
            for file_path in dir_path.glob(f"**/*{ext}"):
                result = self.ingest_file(str(file_path), source_type)
                
                if result.get("ok"):
                    results["files_processed"] += 1
                    results["chunks_total"] += result.get("chunks_created", 0)
                else:
                    results["files_failed"] += 1
                
                results["details"].append({
                    "file": file_path.name,
                    **result
                })
        
        # Statistiken
        results["chunking_stats"] = self.chunker.get_stats()
        
        return results


def run_full_ingestion():
    """
    Führt die komplette Ingestion durch.
    """
    print("=" * 50)
    print("  STRATGEN KNOWLEDGE INGESTION")
    print("=" * 50)
    
    pipeline = KnowledgePipeline()
    
    # 1. Knowledge-Ordner
    print("\n=== 1. Knowledge-Ordner indexieren ===")
    result = pipeline.ingest_directory(
        str(DATA_ROOT / "knowledge"),
        source_type="knowledge",
        extensions=['.txt', '.md']
    )
    print(f"Dateien: {result['files_processed']}, Chunks: {result['chunks_total']}")
    
    # 2. Raw-Ordner (PPTX Text)
    print("\n=== 2. Raw-Ordner (PPTX) indexieren ===")
    result = pipeline.ingest_directory(
        str(DATA_ROOT / "raw"),
        source_type="template",
        extensions=['.pptx']
    )
    print(f"Dateien: {result['files_processed']}, Chunks: {result['chunks_total']}")
    
    # Statistiken
    print("\n=== Chunking Statistiken ===")
    stats = pipeline.chunker.get_stats()
    for key, value in stats.items():
        print(f"  {key}: {value}")
    
    print("\n✅ Ingestion abgeschlossen!")


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    run_full_ingestion()
