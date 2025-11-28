from __future__ import annotations
from pathlib import Path
import sys, traceback

# optionale Parser – wenn Modul fehlt, wird Typ übersprungen
try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    from pypdf import PdfReader
except Exception:
    PdfReader = None

PROJECT_ROOT = Path(__file__).resolve().parents[1]
EXPORTS_DIR  = PROJECT_ROOT / "data" / "exports"
RAW_DIR      = PROJECT_ROOT / "data" / "raw"          # aktuell ungenutzt, aber vorbereitet
DERIVED_DIR  = PROJECT_ROOT / "data" / "knowledge" / "derived"
DERIVED_DIR.mkdir(parents=True, exist_ok=True)

def extract_pptx(p: Path) -> str:
    if Presentation is None:
        raise RuntimeError("python-pptx nicht installiert")
    prs = Presentation(str(p))
    lines = []
    for si, slide in enumerate(prs.slides, 1):
        lines.append(f"--- Slide {si} ---")
        for sh in slide.shapes:
            # Text
            if getattr(sh, "has_text_frame", False) and sh.text_frame:
                t = (sh.text_frame.text or "").strip()
                if t:
                    lines.append(t)
            # Tabellen
            if getattr(sh, "has_table", False) and sh.table:
                for r in sh.table.rows:
                    row_txt = [ (c.text if hasattr(c, "text") else "").strip() for c in r.cells ]
                    if any(row_txt):
                        lines.append(" | ".join(row_txt))
    return "\n".join(lines)

def extract_pdf(p: Path) -> str:
    if PdfReader is None:
        raise RuntimeError("pypdf nicht installiert")
    reader = PdfReader(str(p))
    lines = []
    for i, page in enumerate(reader.pages, 1):
        txt = (page.extract_text() or "").strip()
        lines.append(f"--- Page {i} ---")
        if txt:
            lines.append(txt)
    return "\n".join(lines)

def process_one(p: Path) -> tuple[bool, str, Path]:
    try:
        if p.suffix.lower() == ".pptx":
            text = extract_pptx(p)
        elif p.suffix.lower() == ".pdf":
            text = extract_pdf(p)
        else:
            return False, f"skip (unsupported): {p.name}", Path()
        out = DERIVED_DIR / (p.stem + ".txt")
        out.write_text(text, encoding="utf-8")
        return True, f"ok: {p.name} -> {out.relative_to(PROJECT_ROOT)}", out
    except Exception as e:
        return False, f"err: {p.name} -> {e}", Path()

def main() -> int:
    exts = {".pptx", ".pdf"}
    files = [p for p in EXPORTS_DIR.rglob("*") if p.is_file() and p.suffix.lower() in exts]
    done, skipped = 0, 0
    for f in sorted(files):
        ok, msg, _ = process_one(f)
        print(msg)
        done += (1 if ok else 0)
        skipped += (0 if ok else 1)
    print({"processed": len(files), "ok": done, "failed": skipped, "out_dir": str(DERIVED_DIR)})
    return 0

if __name__ == "__main__":
    sys.exit(main())
