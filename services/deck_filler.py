# -*- coding: utf-8 -*-
from __future__ import annotations
from typing import Any, Dict, List, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

from services.compat import ensure_outline_dict
from services.style_presets import get_style_options

TEMPLATE_DIR = Path("data/templates")

def _resolve_template_path(name: Optional[str]) -> Optional[str]:
    if not name:
        return None
    p = TEMPLATE_DIR / f"{name}.pptx"
    return str(p) if p.exists() else None

def _find_layout(prs: Presentation, names: List[str]) -> Any:
    names = [n.lower() for n in (names or [])]
    for lo in prs.slide_layouts:
        try:
            if lo.name and lo.name.lower() in names:
                return lo
        except Exception:
            pass
    # Fallbacks
    if len(prs.slide_layouts) > 1:
        return prs.slide_layouts[1]
    return prs.slide_layouts[0]

# ---------- Helpers ----------
def _first_placeholder(slide, type_name: str|None=None):
    for shp in getattr(slide, "placeholders", []):
        try:
            pf = shp.placeholder_format
            t = getattr(pf, "type", None)
            if type_name is None or (t and type_name.upper() in str(t).upper()):
                return shp
        except Exception:
            continue
    return None

def _body_placeholder(slide):
    # versuche Textkörper/Content-Placeholder
    for shp in getattr(slide, "placeholders", []):
        try:
            if hasattr(shp, "text_frame") and shp.text_frame is not None:
                # keine Titel-Placeholder
                pf = getattr(shp, "placeholder_format", None)
                if pf and "TITLE" in str(getattr(pf, "type", "")).upper():
                    continue
                return shp
        except Exception:
            continue
    # Fallback: erstes shape mit Text
    for shp in getattr(slide, "shapes", []):
        if hasattr(shp, "text_frame") and shp.text_frame is not None:
            return shp
    return None

def _insert_bullets(placeholder_or_shape, bullets: List[str]):
    if not placeholder_or_shape or not bullets:
        return
    tf = placeholder_or_shape.text_frame
    if tf is None:
        return
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = str(b)
        p.level = 0

def _add_notes(slide, notes: str):
    try:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = str(notes)
    except Exception:
        pass

def _insert_logo(slide, logo_path: Optional[str]):
    if not logo_path:
        return
    try:
        # unten rechts ca. 2.5 cm
        width = Inches(1.2)
        pic = slide.shapes.add_picture(logo_path, slide.width - width - Inches(0.6), slide.height - Inches(1.2), width=width)
        return pic
    except Exception:
        return None

# ---------- Public API ----------
def build_deck(project: Dict[str, Any],
               slide_plan: Optional[List[Dict[str, Any]]],
               outfile: str,
               template: Optional[str] = None) -> str:
    """Baut eine PPTX aus Outline/Slide-Plan. Nutzt Template (falls vorhanden)."""
    outline = ensure_outline_dict(project.get("outline") or {})
    style_name = (project.get("style") or "brand")
    style = get_style_options(style_name)
    max_bullets = int(style.get("max_bullets", 5) or 5)

    tpl_path = _resolve_template_path(template)
    prs = Presentation(tpl_path) if tpl_path else Presentation()

    # Layouts
    lo_title   = _find_layout(prs, ["Title Slide", "Titelfolie"])
    lo_content = _find_layout(prs, ["Title and Content", "Titel und Inhalt", "Title and Body", "Inhalt"])
    lo_2col    = _find_layout(prs, ["Two Content", "Two Column", "Zwei Inhalte"])

    # Titelfolie
    slide = prs.slides.add_slide(lo_title)
    title_text = outline.get("title") or (project.get("topic") or "Deck")
    subtitle_text = project.get("customer_name") or project.get("brand") or None
    try:
        if slide.shapes.title:
            slide.shapes.title.text = title_text
    except Exception:
        pass
    subph = _first_placeholder(slide, type_name="SUBTITLE")
    if subph and subtitle_text:
        _insert_bullets(subph, [subtitle_text])

    logo_path = None
    lp = project.get("logo") or "data/assets/logo.png"
    if lp and Path(lp).exists():
        logo_path = lp
    _insert_logo(slide, logo_path)

    # Sections
    sections = outline.get("sections") or []
    for sec in sections:
        layout_hint = (sec.get("layout_hint") or "").lower()
        lo = lo_2col if ("two" in layout_hint or "2" in layout_hint) else lo_content
        slide = prs.slides.add_slide(lo)
        try:
            if slide.shapes.title:
                slide.shapes.title.text = sec.get("title") or ""
        except Exception:
            pass
        body = _body_placeholder(slide)
        bullets = [str(x) for x in (sec.get("bullets") or [])][:max_bullets]
        if body and bullets:
            _insert_bullets(body, bullets)
        if sec.get("notes"):
            _add_notes(slide, sec["notes"])
        _insert_logo(slide, logo_path)

    out = Path(outfile)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)

# --- ADD/UPDATE in services/deck_filler.py ---

from pptx import Presentation
from pptx.util import Pt, Inches

def _add_notes(slide, text: str):
    try:
        if not text:
            return
        notes = slide.notes_slide or slide.notes_slide  # ensure created
        notes.notes_text_frame.text = text
    except Exception:
        pass

def _insert_bullets(shape, bullets: List[str], max_items: int = 7):
    tf = getattr(shape, "text_frame", None)
    if not tf:
        return
    tf.clear()
    for i, b in enumerate((bullets or [])[:max_items]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = str(b)
        p.level = 0

def _body_placeholder(slide):
    # einfache Suche: erster Platzhalter mit TextFrame außer dem Titel
    for shp in getattr(slide, "placeholders", []):
        try:
            if shp.placeholder_format.type is not None and getattr(shp, "text_frame", None):
                # Titel meist der erste – wir nehmen den ersten, der nicht der Title ist
                if "TITLE" in str(shp.placeholder_format.type):
                    continue
                return shp
        except Exception:
            continue
    # Fallback: erstbeste Form mit text_frame
    for shp in getattr(slide, "shapes", []):
        if getattr(shp, "text_frame", None):
            return shp
    return None

def _add_kpi_table(slide, metrics: Dict[str, Any]):
    if not metrics:
        return
    rows = len(metrics) + 1
    cols = 2
    left = Inches(1.0); top = Inches(2.0); width = Inches(8.0); height = Inches(0.8 + 0.3*rows)
    tbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    tbl.cell(0,0).text = "KPI"; tbl.cell(0,1).text = "Wert"
    r = 1
    for k,v in metrics.items():
        tbl.cell(r,0).text = str(k).upper()
        tbl.cell(r,1).text = str(v)
        r += 1

def build_deck(project: Dict[str, Any],
               slide_plan: Optional[List[Dict[str, Any]]],
               outfile: str,
               template: Optional[str] = None) -> str:
    """Baut PPTX aus Slide-Plan. Notes/Agenda/KPIs/Quellen werden unterstützt."""
    prs = Presentation(_resolve_template_path(template) or None)
    style = get_style_options(project.get("style") or "brand")
    outline = ensure_outline_dict(project.get("outline") or {})
    facts = project.get("facts") or {}

    plan = slide_plan or [{"layout_hint":"Title Slide","title": outline.get("title") or project.get("topic") or "Deck", "kind":"title"}]

    for item in plan:
        layout = _find_layout(prs, [item.get("layout_hint") or "Title and Content"])
        slide = prs.slides.add_slide(layout)
        # Titel
        if slide.shapes.title:
            slide.shapes.title.text = item.get("title") or ""
        # Content
        kind = item.get("kind")
        if kind == "agenda":
            body = _body_placeholder(slide)
            _insert_bullets(body, item.get("bullets") or [])
        elif str(kind).lower() in ("kpis","auto:kpis"):
            _add_kpi_table(slide, (facts or {}).get("metrics") or {})
        elif kind == "sources":
            body = _body_placeholder(slide)
            _insert_bullets(body, item.get("bullets") or [])
        else:
            body = _body_placeholder(slide)
            _insert_bullets(body, item.get("bullets") or [])
            _add_notes(slide, item.get("notes") or "")

    Path(outfile).parent.mkdir(parents=True, exist_ok=True)
    prs.save(outfile)
    return str(outfile)



def _add_table(slide, table_dict):
    # erwartet: {"columns": [...], "rows": [[...], ...]}
    from pptx.util import Inches, Pt
    cols = table_dict.get("columns") or []
    rows = table_dict.get("rows") or []
    nrows = max(1, len(rows)+1)
    ncols = max(1, len(cols))
    x, y, cx, cy = Inches(1.0), Inches(2.0), Inches(8.0), Inches(3.0)
    shape = slide.shapes.add_table(nrows, ncols, x, y, cx, cy)
    table = shape.table
    # Header
    for j, c in enumerate(cols):
        table.cell(0, j).text = str(c)
    # Body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row[:ncols]):
            table.cell(i, j).text = str(val)
    return table
