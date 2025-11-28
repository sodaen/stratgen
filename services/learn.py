# -*- coding: utf-8 -*-
from __future__ import annotations
from typing import Any, Dict, List
from pathlib import Path
from pptx import Presentation

CORPUS_DIR = Path("data/corpus")
MODEL_PATH = Path("data/learned_model.json")

def ingest_export(pptx_path: str|Path) -> Dict[str, Any]:
    CORPUS_DIR.mkdir(parents=True, exist_ok=True)
    p = Path(pptx_path)
    if not p.exists():
        return {"ok": False, "error": "not_found"}
    # Kopie in corpus
    dst = CORPUS_DIR / p.name
    if str(p) != str(dst):
        dst.write_bytes(p.read_bytes())
    # very light stats (für spätere RL-Signale ausbaubar)
    prs = Presentation(str(dst))
    n = len(prs.slides)
    avg_shapes = sum(len(s.shapes) for s in prs.slides) / max(n,1)
    return {"ok": True, "slides": n, "avg_shapes": avg_shapes}

def get_learned_prefs() -> Dict[str, Any]:
    # Placeholder: später Scores aus Telemetrie/Feedback mergen
    return {
        "title_layout_bias": ["Title and Content", "Title Slide"],
        "max_bullets": 5,
        "prefer_body_on_left": True
    }
