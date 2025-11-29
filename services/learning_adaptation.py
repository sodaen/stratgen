# -*- coding: utf-8 -*-
"""
services/learning_adaptation.py
===============================
Stufe 4: Learning & Adaptation

Features:
1. Feedback Integration - Lernen aus User-Feedback
2. Style Learning - Stil aus bestehenden PPTX lernen
3. Quality Prediction - Qualität vorhersagen vor Generierung
4. Preference Learning - User-Präferenzen lernen
5. Performance Analytics - Tracking und Optimierung

Author: StratGen Agent V3.3
"""
from __future__ import annotations
import os
import re
import json
import hashlib
import sqlite3
import time
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field, asdict
from datetime import datetime, timedelta
from collections import defaultdict

# ============================================
# CONFIGURATION
# ============================================

DB_PATH = os.getenv("STRATGEN_DB_PATH", "data/projects.sqlite")
LEARNING_DB = os.getenv("STRATGEN_LEARNING_DB", "data/learning.sqlite")
RAW_DIR = os.getenv("STRATGEN_RAW_DIR", "data/raw")

# Learning Parameters
MIN_FEEDBACK_FOR_LEARNING = 5
PREFERENCE_DECAY_DAYS = 30
QUALITY_PREDICTION_THRESHOLD = 0.7

# ============================================
# DATABASE SETUP
# ============================================

def _connect_learning_db():
    """Verbindet zur Learning-Datenbank."""
    os.makedirs("data", exist_ok=True)
    con = sqlite3.connect(LEARNING_DB, timeout=30)
    con.row_factory = sqlite3.Row
    return con


def ensure_learning_tables():
    """Erstellt Learning-Tabellen falls nicht vorhanden."""
    con = _connect_learning_db()
    cur = con.cursor()
    
    # Feedback-Tabelle
    cur.execute("""
    CREATE TABLE IF NOT EXISTS feedback (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        project_id TEXT NOT NULL,
        slide_index INTEGER,
        feedback_type TEXT NOT NULL,  -- positive, negative, edit, rating
        feedback_value TEXT,
        slide_type TEXT,
        content_hash TEXT,
        created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
    )
    """)
    
    # Learned Patterns
    cur.execute("""
    CREATE TABLE IF NOT EXISTS learned_patterns (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        pattern_type TEXT NOT NULL,  -- structure, style, content, visual
        pattern_key TEXT NOT NULL,
        pattern_value TEXT NOT NULL,
        confidence REAL DEFAULT 0.5,
        usage_count INTEGER DEFAULT 0,
        success_rate REAL DEFAULT 0.5,
        updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
        UNIQUE(pattern_type, pattern_key)
    )
    """)
    
    # User Preferences
    cur.execute("""
    CREATE TABLE IF NOT EXISTS user_preferences (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        preference_key TEXT NOT NULL UNIQUE,
        preference_value TEXT NOT NULL,
        weight REAL DEFAULT 1.0,
        learned_from_count INTEGER DEFAULT 0,
        updated_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
    )
    """)
    
    # Quality History
    cur.execute("""
    CREATE TABLE IF NOT EXISTS quality_history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        project_id TEXT NOT NULL,
        predicted_score REAL,
        actual_score REAL,
        slide_count INTEGER,
        deck_size TEXT,
        industry TEXT,
        topic_hash TEXT,
        created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
    )
    """)
    
    # Style Profiles
    cur.execute("""
    CREATE TABLE IF NOT EXISTS style_profiles (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        profile_name TEXT NOT NULL UNIQUE,
        source_file TEXT,
        bullet_style TEXT,  -- JSON: avg_length, avg_count, capitalization
        layout_preferences TEXT,  -- JSON: preferred layouts
        color_scheme TEXT,  -- JSON: primary, secondary colors
        font_preferences TEXT,  -- JSON: title_size, body_size
        created_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP
    )
    """)
    
    con.commit()
    con.close()


# ============================================
# DATA CLASSES
# ============================================

@dataclass
class FeedbackEntry:
    """Ein Feedback-Eintrag."""
    project_id: str
    feedback_type: str  # positive, negative, edit, rating
    feedback_value: str = ""
    slide_index: int = -1
    slide_type: str = ""
    content_hash: str = ""


@dataclass
class LearnedPattern:
    """Ein gelerntes Muster."""
    pattern_type: str  # structure, style, content, visual
    pattern_key: str
    pattern_value: Any
    confidence: float = 0.5
    usage_count: int = 0
    success_rate: float = 0.5


@dataclass
class StyleProfile:
    """Ein Style-Profil aus gelernten Templates."""
    profile_name: str
    bullet_style: Dict[str, Any] = field(default_factory=dict)
    layout_preferences: Dict[str, Any] = field(default_factory=dict)
    color_scheme: Dict[str, Any] = field(default_factory=dict)
    font_preferences: Dict[str, Any] = field(default_factory=dict)
    source_files: List[str] = field(default_factory=list)


@dataclass
class QualityPrediction:
    """Eine Qualitätsvorhersage."""
    predicted_score: float
    confidence: float
    factors: Dict[str, float] = field(default_factory=dict)
    recommendations: List[str] = field(default_factory=list)


# ============================================
# FEEDBACK INTEGRATION
# ============================================

def record_feedback(
    project_id: str,
    feedback_type: str,
    feedback_value: str = "",
    slide_index: int = -1,
    slide_type: str = "",
    content: str = ""
) -> Dict[str, Any]:
    """
    Speichert User-Feedback.
    
    Args:
        project_id: Projekt-ID
        feedback_type: positive, negative, edit, rating
        feedback_value: z.B. "5" für Rating, Kommentar, etc.
        slide_index: Index des betroffenen Slides (-1 für ganzes Deck)
        slide_type: Typ des Slides
        content: Content des Slides (für Hash)
    
    Returns:
        Status-Dictionary
    """
    ensure_learning_tables()
    
    content_hash = hashlib.md5(content.encode()).hexdigest()[:16] if content else ""
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    cur.execute("""
    INSERT INTO feedback (project_id, feedback_type, feedback_value, slide_index, slide_type, content_hash)
    VALUES (?, ?, ?, ?, ?, ?)
    """, (project_id, feedback_type, feedback_value, slide_index, slide_type, content_hash))
    
    con.commit()
    feedback_id = cur.lastrowid
    con.close()
    
    # Automatisch lernen wenn genug Feedback
    _auto_learn_from_feedback()
    
    return {"ok": True, "feedback_id": feedback_id}


def get_feedback_stats() -> Dict[str, Any]:
    """Gibt Feedback-Statistiken zurück."""
    ensure_learning_tables()
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    stats = {}
    
    # Gesamt-Feedback
    cur.execute("SELECT COUNT(*) FROM feedback")
    stats["total_feedback"] = cur.fetchone()[0]
    
    # Nach Typ
    cur.execute("""
    SELECT feedback_type, COUNT(*) as cnt
    FROM feedback
    GROUP BY feedback_type
    """)
    stats["by_type"] = {row["feedback_type"]: row["cnt"] for row in cur.fetchall()}
    
    # Nach Slide-Typ
    cur.execute("""
    SELECT slide_type, feedback_type, COUNT(*) as cnt
    FROM feedback
    WHERE slide_type != ''
    GROUP BY slide_type, feedback_type
    """)
    by_slide_type = defaultdict(lambda: {"positive": 0, "negative": 0})
    for row in cur.fetchall():
        by_slide_type[row["slide_type"]][row["feedback_type"]] = row["cnt"]
    stats["by_slide_type"] = dict(by_slide_type)
    
    # Durchschnittliches Rating
    cur.execute("""
    SELECT AVG(CAST(feedback_value AS REAL)) as avg_rating
    FROM feedback
    WHERE feedback_type = 'rating' AND feedback_value != ''
    """)
    row = cur.fetchone()
    stats["avg_rating"] = round(row["avg_rating"], 2) if row["avg_rating"] else None
    
    con.close()
    return stats


def _auto_learn_from_feedback():
    """Lernt automatisch aus gesammeltem Feedback."""
    ensure_learning_tables()
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    # Check ob genug Feedback vorhanden
    cur.execute("SELECT COUNT(*) FROM feedback")
    total = cur.fetchone()[0]
    
    if total < MIN_FEEDBACK_FOR_LEARNING:
        con.close()
        return
    
    # Lerne aus Slide-Typ Feedback
    cur.execute("""
    SELECT slide_type,
           SUM(CASE WHEN feedback_type = 'positive' THEN 1 ELSE 0 END) as positive,
           SUM(CASE WHEN feedback_type = 'negative' THEN 1 ELSE 0 END) as negative
    FROM feedback
    WHERE slide_type != ''
    GROUP BY slide_type
    HAVING positive + negative >= 3
    """)
    
    for row in cur.fetchall():
        slide_type = row["slide_type"]
        positive = row["positive"]
        negative = row["negative"]
        total_type = positive + negative
        success_rate = positive / total_type if total_type > 0 else 0.5
        
        # Pattern speichern
        cur.execute("""
        INSERT OR REPLACE INTO learned_patterns (pattern_type, pattern_key, pattern_value, confidence, success_rate, updated_at)
        VALUES ('slide_type', ?, ?, ?, ?, CURRENT_TIMESTAMP)
        """, (slide_type, json.dumps({"positive": positive, "negative": negative}), 
              min(1.0, total_type / 10), success_rate))
    
    con.commit()
    con.close()


# ============================================
# STYLE LEARNING
# ============================================

def learn_style_from_pptx(file_path: str, profile_name: str = None) -> Dict[str, Any]:
    """
    Lernt Stil aus einer PPTX-Datei.
    
    Args:
        file_path: Pfad zur PPTX
        profile_name: Name für das Profil (default: Dateiname)
    
    Returns:
        StyleProfile als Dictionary
    """
    from pathlib import Path
    
    if not Path(file_path).exists():
        return {"ok": False, "error": "File not found"}
    
    if not profile_name:
        profile_name = Path(file_path).stem
    
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed"}
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        return {"ok": False, "error": str(e)}
    
    # Analysiere Präsentation
    bullet_lengths = []
    bullet_counts = []
    layouts_used = defaultdict(int)
    font_sizes = {"title": [], "body": []}
    
    for slide in prs.slides:
        layout_name = slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown"
        layouts_used[layout_name] += 1
        
        bullet_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        bullet_lengths.append(len(text))
                        bullet_count += 1
                        
                        # Font-Größe
                        if para.runs:
                            try:
                                size = para.runs[0].font.size
                                if size:
                                    pt = size.pt if hasattr(size, 'pt') else int(size) / 12700
                                    if pt > 24:
                                        font_sizes["title"].append(pt)
                                    else:
                                        font_sizes["body"].append(pt)
                            except:
                                pass
        
        if bullet_count > 0:
            bullet_counts.append(bullet_count)
    
    # Erstelle StyleProfile
    profile = StyleProfile(
        profile_name=profile_name,
        bullet_style={
            "avg_length": round(sum(bullet_lengths) / max(1, len(bullet_lengths)), 1),
            "avg_count": round(sum(bullet_counts) / max(1, len(bullet_counts)), 1),
            "max_length": max(bullet_lengths) if bullet_lengths else 100,
        },
        layout_preferences={
            "most_used": max(layouts_used, key=layouts_used.get) if layouts_used else "Title and Content",
            "usage": dict(layouts_used)
        },
        font_preferences={
            "title_size": round(sum(font_sizes["title"]) / max(1, len(font_sizes["title"])), 1) if font_sizes["title"] else 32,
            "body_size": round(sum(font_sizes["body"]) / max(1, len(font_sizes["body"])), 1) if font_sizes["body"] else 18,
        },
        source_files=[file_path]
    )
    
    # In DB speichern
    _save_style_profile(profile)
    
    return {"ok": True, "profile": asdict(profile)}


def _save_style_profile(profile: StyleProfile):
    """Speichert ein StyleProfile in der DB."""
    ensure_learning_tables()
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    cur.execute("""
    INSERT OR REPLACE INTO style_profiles (profile_name, source_file, bullet_style, layout_preferences, font_preferences)
    VALUES (?, ?, ?, ?, ?)
    """, (
        profile.profile_name,
        profile.source_files[0] if profile.source_files else "",
        json.dumps(profile.bullet_style),
        json.dumps(profile.layout_preferences),
        json.dumps(profile.font_preferences)
    ))
    
    con.commit()
    con.close()


def get_style_profile(profile_name: str) -> Optional[StyleProfile]:
    """Lädt ein StyleProfile aus der DB."""
    ensure_learning_tables()
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    cur.execute("""
    SELECT * FROM style_profiles WHERE profile_name = ?
    """, (profile_name,))
    
    row = cur.fetchone()
    con.close()
    
    if not row:
        return None
    
    return StyleProfile(
        profile_name=row["profile_name"],
        bullet_style=json.loads(row["bullet_style"] or "{}"),
        layout_preferences=json.loads(row["layout_preferences"] or "{}"),
        font_preferences=json.loads(row["font_preferences"] or "{}"),
        source_files=[row["source_file"]] if row["source_file"] else []
    )


def learn_from_all_templates() -> Dict[str, Any]:
    """Lernt aus allen Templates in /raw."""
    raw_path = Path(RAW_DIR)
    if not raw_path.exists():
        return {"ok": False, "error": "RAW_DIR not found"}
    
    learned = 0
    errors = []
    
    for pptx_file in raw_path.glob("*.pptx"):
        result = learn_style_from_pptx(str(pptx_file))
        if result.get("ok"):
            learned += 1
        else:
            errors.append(f"{pptx_file.name}: {result.get('error')}")
    
    return {"ok": True, "learned": learned, "errors": errors}


def get_merged_style() -> Dict[str, Any]:
    """Gibt einen gemittelten Stil aus allen Profilen zurück."""
    ensure_learning_tables()
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    cur.execute("SELECT * FROM style_profiles")
    rows = cur.fetchall()
    con.close()
    
    if not rows:
        return {
            "bullet_style": {"avg_length": 60, "avg_count": 4, "max_length": 120},
            "font_preferences": {"title_size": 32, "body_size": 18}
        }
    
    # Mittelwerte berechnen
    bullet_lengths = []
    bullet_counts = []
    title_sizes = []
    body_sizes = []
    
    for row in rows:
        bs = json.loads(row["bullet_style"] or "{}")
        fp = json.loads(row["font_preferences"] or "{}")
        
        if bs.get("avg_length"):
            bullet_lengths.append(bs["avg_length"])
        if bs.get("avg_count"):
            bullet_counts.append(bs["avg_count"])
        if fp.get("title_size"):
            title_sizes.append(fp["title_size"])
        if fp.get("body_size"):
            body_sizes.append(fp["body_size"])
    
    return {
        "bullet_style": {
            "avg_length": round(sum(bullet_lengths) / max(1, len(bullet_lengths)), 1),
            "avg_count": round(sum(bullet_counts) / max(1, len(bullet_counts)), 1),
            "max_length": 120
        },
        "font_preferences": {
            "title_size": round(sum(title_sizes) / max(1, len(title_sizes)), 1),
            "body_size": round(sum(body_sizes) / max(1, len(body_sizes)), 1)
        },
        "profile_count": len(rows)
    }


# ============================================
# QUALITY PREDICTION
# ============================================

def predict_quality(
    topic: str,
    brief: str,
    deck_size: str,
    industry: str = "",
    slide_types: List[str] = None
) -> QualityPrediction:
    """
    Sagt die erwartete Qualität vorher.
    
    Args:
        topic: Thema der Präsentation
        brief: Briefing-Text
        deck_size: short, medium, large
        industry: Branche
        slide_types: Geplante Slide-Typen
    
    Returns:
        QualityPrediction
    """
    ensure_learning_tables()
    
    factors = {}
    recommendations = []
    base_score = 6.5
    
    # Faktor 1: Briefing-Länge
    brief_length = len(brief)
    if brief_length < 50:
        factors["brief_length"] = -0.5
        recommendations.append("Detaillierteres Briefing würde bessere Ergebnisse liefern")
    elif brief_length > 200:
        factors["brief_length"] = 0.5
        recommendations.append("Ausführliches Briefing - gute Basis")
    else:
        factors["brief_length"] = 0.0
    
    # Faktor 2: Historische Qualität für ähnliche Topics
    topic_hash = hashlib.md5(topic.lower().encode()).hexdigest()[:8]
    con = _connect_learning_db()
    cur = con.cursor()
    
    cur.execute("""
    SELECT AVG(actual_score) as avg_score, COUNT(*) as cnt
    FROM quality_history
    WHERE topic_hash = ? OR industry = ?
    """, (topic_hash, industry))
    
    row = cur.fetchone()
    if row["cnt"] and row["cnt"] >= 2:
        historical_avg = row["avg_score"] or 6.5
        factors["historical"] = (historical_avg - 6.5) / 2
    else:
        factors["historical"] = 0.0
    
    # Faktor 3: Slide-Typ Success Rates
    if slide_types:
        cur.execute("""
        SELECT pattern_key, success_rate
        FROM learned_patterns
        WHERE pattern_type = 'slide_type' AND pattern_key IN ({})
        """.format(",".join("?" * len(slide_types))), slide_types)
        
        success_rates = [row["success_rate"] for row in cur.fetchall()]
        if success_rates:
            avg_success = sum(success_rates) / len(success_rates)
            factors["slide_types"] = (avg_success - 0.5) * 2
    
    con.close()
    
    # Faktor 4: Deck-Size Angemessenheit
    words_in_brief = len(brief.split())
    if deck_size == "short" and words_in_brief > 100:
        factors["deck_size_fit"] = -0.3
        recommendations.append("Briefing ist umfangreich - 'medium' könnte besser passen")
    elif deck_size == "large" and words_in_brief < 30:
        factors["deck_size_fit"] = -0.3
        recommendations.append("Kurzes Briefing - 'short' könnte ausreichen")
    else:
        factors["deck_size_fit"] = 0.2
    
    # Finalen Score berechnen
    adjustment = sum(factors.values())
    predicted_score = max(4.0, min(10.0, base_score + adjustment))
    
    # Confidence basierend auf verfügbaren Daten
    data_points = sum(1 for v in factors.values() if v != 0)
    confidence = min(0.9, 0.5 + data_points * 0.1)
    
    return QualityPrediction(
        predicted_score=round(predicted_score, 1),
        confidence=round(confidence, 2),
        factors=factors,
        recommendations=recommendations
    )


def record_quality_result(
    project_id: str,
    predicted_score: float,
    actual_score: float,
    slide_count: int,
    deck_size: str,
    industry: str,
    topic: str
):
    """Speichert ein Qualitäts-Ergebnis für zukünftiges Lernen."""
    ensure_learning_tables()
    
    topic_hash = hashlib.md5(topic.lower().encode()).hexdigest()[:8]
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    cur.execute("""
    INSERT INTO quality_history (project_id, predicted_score, actual_score, slide_count, deck_size, industry, topic_hash)
    VALUES (?, ?, ?, ?, ?, ?, ?)
    """, (project_id, predicted_score, actual_score, slide_count, deck_size, industry, topic_hash))
    
    con.commit()
    con.close()


# ============================================
# PREFERENCE LEARNING
# ============================================

def update_preference(key: str, value: str, weight: float = 1.0):
    """Aktualisiert eine User-Präferenz."""
    ensure_learning_tables()
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    cur.execute("""
    INSERT INTO user_preferences (preference_key, preference_value, weight, learned_from_count, updated_at)
    VALUES (?, ?, ?, 1, CURRENT_TIMESTAMP)
    ON CONFLICT(preference_key) DO UPDATE SET
        preference_value = ?,
        weight = (weight + ?) / 2,
        learned_from_count = learned_from_count + 1,
        updated_at = CURRENT_TIMESTAMP
    """, (key, value, weight, value, weight))
    
    con.commit()
    con.close()


def get_preferences() -> Dict[str, Any]:
    """Gibt alle gelernten Präferenzen zurück."""
    ensure_learning_tables()
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    cur.execute("SELECT * FROM user_preferences ORDER BY weight DESC")
    rows = cur.fetchall()
    con.close()
    
    return {
        row["preference_key"]: {
            "value": row["preference_value"],
            "weight": row["weight"],
            "learned_from": row["learned_from_count"]
        }
        for row in rows
    }


def apply_preferences_to_request(request: Dict[str, Any]) -> Dict[str, Any]:
    """Wendet gelernte Präferenzen auf einen Request an."""
    preferences = get_preferences()
    
    # Beispiel-Anwendungen
    if "deck_size" in preferences and not request.get("deck_size"):
        request["deck_size"] = preferences["deck_size"]["value"]
    
    if "preferred_language" in preferences:
        request["language"] = preferences["preferred_language"]["value"]
    
    return request


# ============================================
# IMPROVEMENT SUGGESTIONS
# ============================================

def get_improvement_suggestions(
    content: str,
    slide_types: List[str]
) -> Dict[str, Any]:
    """
    Gibt Verbesserungsvorschläge basierend auf gelernten Mustern.
    
    Args:
        content: Generierter Content
        slide_types: Verwendete Slide-Typen
    
    Returns:
        Dictionary mit suggestions
    """
    ensure_learning_tables()
    
    suggestions = []
    
    # Stil-Check gegen gelernten Stil
    merged_style = get_merged_style()
    
    bullet_style = merged_style.get("bullet_style", {})
    optimal_length = bullet_style.get("avg_length", 60)
    optimal_count = bullet_style.get("avg_count", 4)
    
    # Analysiere Content
    lines = [l.strip() for l in content.split("\n") if l.strip()]
    avg_line_length = sum(len(l) for l in lines) / max(1, len(lines))
    
    if avg_line_length > optimal_length * 1.5:
        suggestions.append({
            "type": "style",
            "issue": "Bullets sind länger als gewöhnlich",
            "recommendation": f"Ziel: ~{int(optimal_length)} Zeichen pro Bullet",
            "severity": "low"
        })
    
    # Slide-Typ basierte Vorschläge
    con = _connect_learning_db()
    cur = con.cursor()
    
    for slide_type in slide_types:
        cur.execute("""
        SELECT success_rate FROM learned_patterns
        WHERE pattern_type = 'slide_type' AND pattern_key = ?
        """, (slide_type,))
        
        row = cur.fetchone()
        if row and row["success_rate"] < 0.4:
            suggestions.append({
                "type": "content",
                "issue": f"Slide-Typ '{slide_type}' hat niedrige Erfolgsrate",
                "recommendation": "Besondere Aufmerksamkeit auf diesen Slide-Typ",
                "severity": "medium"
            })
    
    con.close()
    
    return {"ok": True, "suggestions": suggestions}


# ============================================
# API FUNCTIONS
# ============================================

def check_status() -> Dict[str, Any]:
    """Gibt den Status des Learning-Systems zurück."""
    ensure_learning_tables()
    
    con = _connect_learning_db()
    cur = con.cursor()
    
    # Feedback-Count
    cur.execute("SELECT COUNT(*) FROM feedback")
    feedback_count = cur.fetchone()[0]
    
    # Patterns
    cur.execute("SELECT COUNT(*) FROM learned_patterns")
    pattern_count = cur.fetchone()[0]
    
    # Style Profiles
    cur.execute("SELECT COUNT(*) FROM style_profiles")
    profile_count = cur.fetchone()[0]
    
    # Quality History
    cur.execute("SELECT COUNT(*) FROM quality_history")
    quality_count = cur.fetchone()[0]
    
    con.close()
    
    return {
        "ok": True,
        "learning_active": feedback_count >= MIN_FEEDBACK_FOR_LEARNING,
        "stats": {
            "feedback_entries": feedback_count,
            "learned_patterns": pattern_count,
            "style_profiles": profile_count,
            "quality_history": quality_count,
            "min_feedback_threshold": MIN_FEEDBACK_FOR_LEARNING
        }
    }
