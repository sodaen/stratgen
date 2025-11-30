# -*- coding: utf-8 -*-
"""
services/semantic_slide_matcher.py
==================================
Feature: Semantic Slide Matcher

Findet semantisch ähnliche Slides aus deinen Templates:
1. Embeddings für alle Template-Slides
2. Similarity Search bei neuem Briefing
3. Konkrete Formulierungsvorschläge
4. "Best Practice" Slides vorschlagen

Author: StratGen Agent V3.6
"""
from __future__ import annotations
import os
import re
import json
import sqlite3
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field, asdict
from datetime import datetime

# ============================================
# CONFIGURATION
# ============================================

DB_PATH = os.getenv("STRATGEN_MATCHER_DB", "data/slide_matcher.sqlite")
EMBEDDING_MODEL = "nomic-embed-text"
SIMILARITY_THRESHOLD = 0.7

# ============================================
# EMBEDDING SERVICE
# ============================================

try:
    import requests
    OLLAMA_HOST = os.getenv("OLLAMA_HOST", "http://127.0.0.1:11434")
    
    def get_embedding(text: str) -> List[float]:
        """Holt Embedding von Ollama."""
        try:
            response = requests.post(
                f"{OLLAMA_HOST}/api/embeddings",
                json={"model": EMBEDDING_MODEL, "prompt": text[:2000]},
                timeout=30
            )
            if response.ok:
                return response.json().get("embedding", [])
        except Exception:
            pass
        return []
    
    HAS_EMBEDDINGS = True
except ImportError:
    HAS_EMBEDDINGS = False
    def get_embedding(text: str) -> List[float]:
        return []


# ============================================
# DATA CLASSES
# ============================================

@dataclass
class SlideMatch:
    """Ein gefundener Slide-Match."""
    slide_id: int
    template_name: str
    slide_index: int
    slide_type: str
    title: str
    content_preview: str
    similarity_score: float
    source_file: str


@dataclass 
class FormulationSuggestion:
    """Ein Formulierungsvorschlag."""
    original_text: str
    suggested_text: str
    source_template: str
    confidence: float


# ============================================
# DATABASE
# ============================================

def _connect():
    """Verbindet zur Matcher-Datenbank."""
    os.makedirs("data", exist_ok=True)
    con = sqlite3.connect(DB_PATH, timeout=30)
    con.row_factory = sqlite3.Row
    return con


def ensure_tables():
    """Erstellt Tabellen falls nicht vorhanden."""
    con = _connect()
    cur = con.cursor()
    
    # Slide-Index mit Embeddings
    cur.execute("""
    CREATE TABLE IF NOT EXISTS slide_index (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        template_path TEXT NOT NULL,
        template_name TEXT,
        slide_index INTEGER,
        slide_type TEXT,
        title TEXT,
        content TEXT,
        content_hash TEXT UNIQUE,
        embedding BLOB,
        created_at TEXT DEFAULT CURRENT_TIMESTAMP
    )
    """)
    
    # Bullet-Index für Formulierungen
    cur.execute("""
    CREATE TABLE IF NOT EXISTS bullet_index (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        slide_id INTEGER,
        bullet_text TEXT,
        bullet_hash TEXT UNIQUE,
        embedding BLOB,
        FOREIGN KEY (slide_id) REFERENCES slide_index(id)
    )
    """)
    
    con.commit()
    con.close()


# ============================================
# INDEXING
# ============================================

def index_template(file_path: str) -> Dict[str, Any]:
    """
    Indiziert ein Template für Semantic Search.
    
    Args:
        file_path: Pfad zur PPTX-Datei
    
    Returns:
        Indexierungs-Ergebnis
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed"}
    
    ensure_tables()
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        return {"ok": False, "error": str(e)}
    
    template_name = Path(file_path).stem
    con = _connect()
    cur = con.cursor()
    
    indexed_slides = 0
    indexed_bullets = 0
    
    for idx, slide in enumerate(prs.slides):
        slide_texts = []
        title = ""
        bullets = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
                        if not title and len(text) < 100:
                            title = text
                        elif len(text) > 15:
                            bullets.append(text)
        
        if not slide_texts:
            continue
        
        content = " ".join(slide_texts)
        content_hash = hashlib.md5(content.encode()).hexdigest()
        
        # Bereits indiziert?
        cur.execute("SELECT id FROM slide_index WHERE content_hash = ?", (content_hash,))
        if cur.fetchone():
            continue
        
        # Slide-Typ erkennen
        from services.slide_dna_analyzer import detect_slide_type
        slide_type = detect_slide_type(title, content)
        
        # Embedding generieren
        embedding = get_embedding(content)
        embedding_blob = json.dumps(embedding).encode() if embedding else None
        
        # Speichern
        cur.execute("""
        INSERT INTO slide_index (template_path, template_name, slide_index, slide_type, title, content, content_hash, embedding)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
        """, (file_path, template_name, idx, slide_type, title, content[:2000], content_hash, embedding_blob))
        
        slide_id = cur.lastrowid
        indexed_slides += 1
        
        # Bullets indizieren
        for bullet in bullets[:10]:
            bullet_hash = hashlib.md5(bullet.encode()).hexdigest()
            
            cur.execute("SELECT id FROM bullet_index WHERE bullet_hash = ?", (bullet_hash,))
            if cur.fetchone():
                continue
            
            bullet_embedding = get_embedding(bullet)
            bullet_blob = json.dumps(bullet_embedding).encode() if bullet_embedding else None
            
            cur.execute("""
            INSERT INTO bullet_index (slide_id, bullet_text, bullet_hash, embedding)
            VALUES (?, ?, ?, ?)
            """, (slide_id, bullet, bullet_hash, bullet_blob))
            
            indexed_bullets += 1
    
    con.commit()
    con.close()
    
    return {
        "ok": True,
        "template": template_name,
        "indexed_slides": indexed_slides,
        "indexed_bullets": indexed_bullets
    }


def index_all_templates() -> Dict[str, Any]:
    """Indiziert alle Templates."""
    from pathlib import Path
    
    raw_dir = Path(os.getenv("STRATGEN_RAW_DIR", "data/raw"))
    if not raw_dir.exists():
        return {"ok": False, "error": "RAW_DIR not found"}
    
    results = {
        "ok": True,
        "total_files": 0,
        "indexed_slides": 0,
        "indexed_bullets": 0,
        "errors": []
    }
    
    for pptx_file in raw_dir.glob("**/*.pptx"):
        results["total_files"] += 1
        
        result = index_template(str(pptx_file))
        if result.get("ok"):
            results["indexed_slides"] += result.get("indexed_slides", 0)
            results["indexed_bullets"] += result.get("indexed_bullets", 0)
        else:
            results["errors"].append({"file": str(pptx_file), "error": result.get("error")})
    
    return results


# ============================================
# SIMILARITY SEARCH
# ============================================

def cosine_similarity(vec1: List[float], vec2: List[float]) -> float:
    """Berechnet Cosine Similarity zwischen zwei Vektoren."""
    if not vec1 or not vec2 or len(vec1) != len(vec2):
        return 0.0
    
    dot_product = sum(a * b for a, b in zip(vec1, vec2))
    norm1 = sum(a * a for a in vec1) ** 0.5
    norm2 = sum(b * b for b in vec2) ** 0.5
    
    if norm1 == 0 or norm2 == 0:
        return 0.0
    
    return dot_product / (norm1 * norm2)


def find_similar_slides(
    query: str,
    slide_type: str = None,
    top_k: int = 5,
    min_similarity: float = 0.5
) -> List[SlideMatch]:
    """
    Findet ähnliche Slides zu einer Query.
    
    Args:
        query: Suchtext (Briefing, Thema, etc.)
        slide_type: Optional: Nur bestimmten Slide-Typ
        top_k: Anzahl Ergebnisse
        min_similarity: Minimale Ähnlichkeit
    
    Returns:
        Liste von SlideMatch-Objekten
    """
    ensure_tables()
    
    # Query Embedding
    query_embedding = get_embedding(query)
    if not query_embedding:
        return []
    
    con = _connect()
    cur = con.cursor()
    
    # Alle Slides mit Embeddings laden
    if slide_type:
        cur.execute("""
        SELECT id, template_path, template_name, slide_index, slide_type, title, content, embedding
        FROM slide_index WHERE slide_type = ? AND embedding IS NOT NULL
        """, (slide_type,))
    else:
        cur.execute("""
        SELECT id, template_path, template_name, slide_index, slide_type, title, content, embedding
        FROM slide_index WHERE embedding IS NOT NULL
        """)
    
    matches = []
    for row in cur.fetchall():
        if not row["embedding"]:
            continue
        
        slide_embedding = json.loads(row["embedding"])
        similarity = cosine_similarity(query_embedding, slide_embedding)
        
        if similarity >= min_similarity:
            matches.append(SlideMatch(
                slide_id=row["id"],
                template_name=row["template_name"],
                slide_index=row["slide_index"],
                slide_type=row["slide_type"],
                title=row["title"] or "",
                content_preview=row["content"][:200] if row["content"] else "",
                similarity_score=round(similarity, 3),
                source_file=row["template_path"]
            ))
    
    con.close()
    
    # Nach Similarity sortieren
    matches.sort(key=lambda x: x.similarity_score, reverse=True)
    
    return matches[:top_k]


def find_similar_bullets(
    query: str,
    top_k: int = 10,
    min_similarity: float = 0.6
) -> List[FormulationSuggestion]:
    """
    Findet ähnliche Bullet-Formulierungen.
    
    Args:
        query: Der zu matchende Text
        top_k: Anzahl Ergebnisse
        min_similarity: Minimale Ähnlichkeit
    
    Returns:
        Liste von FormulationSuggestion-Objekten
    """
    ensure_tables()
    
    query_embedding = get_embedding(query)
    if not query_embedding:
        return []
    
    con = _connect()
    cur = con.cursor()
    
    cur.execute("""
    SELECT b.id, b.bullet_text, b.embedding, s.template_name
    FROM bullet_index b
    JOIN slide_index s ON b.slide_id = s.id
    WHERE b.embedding IS NOT NULL
    """)
    
    suggestions = []
    for row in cur.fetchall():
        if not row["embedding"]:
            continue
        
        bullet_embedding = json.loads(row["embedding"])
        similarity = cosine_similarity(query_embedding, bullet_embedding)
        
        if similarity >= min_similarity:
            suggestions.append(FormulationSuggestion(
                original_text=query,
                suggested_text=row["bullet_text"],
                source_template=row["template_name"],
                confidence=round(similarity, 3)
            ))
    
    con.close()
    
    suggestions.sort(key=lambda x: x.confidence, reverse=True)
    
    return suggestions[:top_k]


# ============================================
# HIGH-LEVEL API
# ============================================

def get_slide_suggestions(
    briefing: str,
    slide_type: str = None,
    industry: str = ""
) -> Dict[str, Any]:
    """
    Gibt Slide-Vorschläge basierend auf Briefing.
    
    Args:
        briefing: Das Projekt-Briefing
        slide_type: Optional: Spezifischer Slide-Typ
        industry: Optional: Branche
    
    Returns:
        Dictionary mit Vorschlägen
    """
    query = f"{briefing} {industry}".strip()
    
    similar_slides = find_similar_slides(query, slide_type, top_k=5)
    
    if not similar_slides:
        return {
            "ok": True,
            "suggestions": [],
            "message": "Keine ähnlichen Slides gefunden. Möglicherweise müssen Templates erst indiziert werden."
        }
    
    return {
        "ok": True,
        "query": query[:100],
        "suggestions": [
            {
                "template": m.template_name,
                "slide_index": m.slide_index,
                "type": m.slide_type,
                "title": m.title,
                "preview": m.content_preview,
                "similarity": m.similarity_score,
                "source": m.source_file
            }
            for m in similar_slides
        ],
        "best_match": {
            "template": similar_slides[0].template_name,
            "similarity": similar_slides[0].similarity_score
        } if similar_slides else None
    }


def get_formulation_suggestions(
    text: str,
    context: str = ""
) -> Dict[str, Any]:
    """
    Gibt Formulierungsvorschläge für einen Text.
    
    Args:
        text: Der zu verbessernde Text
        context: Zusätzlicher Kontext
    
    Returns:
        Dictionary mit Vorschlägen
    """
    query = f"{text} {context}".strip()
    
    suggestions = find_similar_bullets(query, top_k=5)
    
    return {
        "ok": True,
        "original": text,
        "suggestions": [
            {
                "text": s.suggested_text,
                "source": s.source_template,
                "confidence": s.confidence
            }
            for s in suggestions
        ]
    }


# ============================================
# STATUS
# ============================================

def check_status() -> Dict[str, Any]:
    """Gibt den Status des Semantic Matchers zurück."""
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    
    cur.execute("SELECT COUNT(*) as cnt FROM slide_index")
    slides = cur.fetchone()["cnt"]
    
    cur.execute("SELECT COUNT(*) as cnt FROM slide_index WHERE embedding IS NOT NULL")
    slides_with_embedding = cur.fetchone()["cnt"]
    
    cur.execute("SELECT COUNT(*) as cnt FROM bullet_index")
    bullets = cur.fetchone()["cnt"]
    
    cur.execute("SELECT COUNT(DISTINCT template_name) as cnt FROM slide_index")
    templates = cur.fetchone()["cnt"]
    
    con.close()
    
    return {
        "ok": True,
        "indexed_templates": templates,
        "indexed_slides": slides,
        "slides_with_embeddings": slides_with_embedding,
        "indexed_bullets": bullets,
        "embedding_model": EMBEDDING_MODEL,
        "has_embeddings": HAS_EMBEDDINGS
    }
