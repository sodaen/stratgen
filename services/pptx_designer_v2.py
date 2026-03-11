"""
PPTX Designer Service v2 für Stratgen.
- Nutzt Wizard-Farbauswahl
- Fügt Quellenangaben hinzu
- Professionelle Layouts
"""

import os
import logging
from pathlib import Path
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json

logger = logging.getLogger(__name__)

DATA_ROOT = Path(os.getenv("STRATGEN_DATA", "/home/sodaen/stratgen/data"))


class PPTXDesignerV2:
    """Erstellt professionelle PPTX-Präsentationen mit Wizard-Farben und Quellen."""
    
    # Fallback-Paletten
    PALETTES = {
        "corporate": {
            "primary": "#1E40AF",
            "secondary": "#3B82F6", 
            "accent": "#10B981",
            "background": "#FFFFFF",
            "text": "#111827"
        },
        "modern": {
            "primary": "#7C3AED",
            "secondary": "#A78BFA",
            "accent": "#F59E0B",
            "background": "#FAFAFA",
            "text": "#27272A"
        },
        "minimal": {
            "primary": "#000000",
            "secondary": "#525252",
            "accent": "#3B82F6",
            "background": "#FFFFFF",
            "text": "#262626"
        },
        "vibrant": {
            "primary": "#EC4899",
            "secondary": "#8B5CF6",
            "accent": "#06B6D4",
            "background": "#FDF4FF",
            "text": "#581C87"
        }
    }
    
    def __init__(self, colors: Dict[str, str] = None, palette: str = "corporate"):
        """
        Args:
            colors: Wizard-Farben dict mit primary, secondary, accent, background, text
            palette: Fallback-Palette wenn keine Farben übergeben
        """
        if colors:
            self.colors = {
                "primary": colors.get("primary", "#1E40AF"),
                "secondary": colors.get("secondary", "#3B82F6"),
                "accent": colors.get("accent", "#10B981"),
                "background": colors.get("background", "#FFFFFF"),
                "text": colors.get("text", "#111827")
            }
        else:
            self.colors = self.PALETTES.get(palette, self.PALETTES["corporate"])
        
        self.prs = None
        self.slide_sources = []  # Sammelt Quellen pro Slide
        # Sprint 2: Image + Layout tracking
        self._customer_name = ""
        self._topic = ""
        self._use_images = True  # Auto Images aktiviert
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Konvertiert Hex zu RGB."""
        hex_color = hex_color.lstrip('#')
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    
    def _add_source_footer(self, slide, sources: List[str]):
        """Fügt Quellenangaben als Footer hinzu."""
        if not sources:
            return
        
        # Formatiere Quellen
        source_text = "Quellen: " + " | ".join(sources[:3])  # Max 3 Quellen
        if len(sources) > 3:
            source_text += f" (+{len(sources)-3} weitere)"
        
        # Footer Box
        footer = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.1), Inches(12.333), Inches(0.3)
        )
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = source_text
        p.font.size = Pt(8)
        p.font.color.rgb = self._hex_to_rgb(self.colors.get("secondary", "#666666"))
        p.alignment = PP_ALIGN.LEFT
    
    def _add_slide_number(self, slide, number: int, total: int):
        """Fügt Seitenzahl hinzu."""
        num_box = slide.shapes.add_textbox(
            Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{number}/{total}"
        p.font.size = Pt(9)
        p.font.color.rgb = self._hex_to_rgb(self.colors.get("secondary", "#666666"))
        p.alignment = PP_ALIGN.RIGHT
    
    def create_presentation(self, slides: List[Dict],
                           title: str = "Präsentation",
                           company: str = "",
                           include_sources_slide: bool = True,
                           customer_name: str = "",
                           use_images: bool = True) -> bytes:
        """
        Erstellt eine komplette PPTX-Präsentation.
        
        Args:
            slides: Liste von Slide-Dicts mit type, title, content, sources
            title: Titel der Präsentation
            company: Firmenname für Footer
            include_sources_slide: Fügt am Ende eine Quellenübersicht hinzu
        
        Returns:
            PPTX als Bytes
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)
        self._customer_name = customer_name or company
        self._topic = title
        self._use_images = use_images
        
        all_sources = []
        total_slides = len(slides) + (1 if include_sources_slide else 0)
        
        for i, slide_data in enumerate(slides):
            slide_type = slide_data.get("type", "content")
            sources = slide_data.get("sources", [])
            all_sources.extend(sources)
            
            # Erstelle Slide basierend auf Typ – Sprint 2: vollständiges Routing
            if slide_type == "title" or i == 0:
                slide = self._add_title_slide(slide_data, company)
            elif slide_type == "section":
                slide = self._add_section_slide(slide_data)
            elif slide_type == "two_column" or slide_type == "comparison":
                slide = self._add_two_column_slide(slide_data)
            elif slide_type == "bullets" or slide_type == "list":
                slide = self._add_bullet_slide(slide_data)
            elif slide_type == "quote" or slide_type == "statement":
                slide = self._add_quote_slide(slide_data)
            elif slide_type in ("data", "chart", "bar", "line", "pie"):
                slide = self._add_data_slide(slide_data)
            elif slide_type == "image" or slide_type == "visual":
                slide = self._add_image_slide(slide_data)
            elif slide_type == "agenda":
                slide = self._add_agenda_slide(slide_data)
            elif slide_type == "cta" or slide_type == "next_steps":
                slide = self._add_cta_slide(slide_data)
            elif slide_type == "kpi" or slide_type == "metrics":
                slide = self._add_kpi_slide(slide_data)
            elif slide_type == "swot":
                slide = self._add_swot_slide(slide_data)
            elif slide_type == "timeline" or slide_type == "roadmap":
                slide = self._add_timeline_slide(slide_data)
            else:
                slide = self._add_content_slide(slide_data)
            
            # Auto Images: Bild einfügen wenn vorhanden oder gefunden
            if self._use_images and slide_type not in ("title", "section"):
                self._try_add_image(slide, slide_data)
            
            # Füge Quellen-Footer hinzu (außer bei Title-Slide)
            if sources and slide_type != "title":
                self._add_source_footer(slide, sources)
            
            # Seitenzahl (außer bei Title-Slide)
            if slide_type != "title":
                self._add_slide_number(slide, i + 1, total_slides)
        
        # Quellenübersicht am Ende
        if include_sources_slide and all_sources:
            self._add_sources_overview_slide(list(set(all_sources)))
        
        # Speichere als Bytes
        from io import BytesIO
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.read()
    
    def _add_background(self, slide, use_primary: bool = False):
        """Fügt Hintergrundfarbe hinzu."""
        background = slide.background
        fill = background.fill
        fill.solid()
        color_key = "primary" if use_primary else "background"
        fill.fore_color.rgb = self._hex_to_rgb(self.colors[color_key])
    
    def _add_title_slide(self, data: Dict, company: str = ""):
        """Erstellt Title-Slide."""
        layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide, use_primary=True)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb("#FFFFFF")
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        if data.get("subtitle"):
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(1)
            )
            tf = sub_box.text_frame
            p = tf.paragraphs[0]
            p.text = data.get("subtitle", "")
            p.font.size = Pt(24)
            p.font.color.rgb = self._hex_to_rgb("#FFFFFF")
            p.alignment = PP_ALIGN.CENTER
        
        # Accent line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5.5), Inches(4), Inches(2.333), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.colors["accent"])
        line.line.fill.background()
        
        # Company name
        if company:
            comp_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.5)
            )
            tf = comp_box.text_frame
            p = tf.paragraphs[0]
            p.text = company
            p.font.size = Pt(14)
            p.font.color.rgb = self._hex_to_rgb("#FFFFFF")
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _add_section_slide(self, data: Dict):
        """Erstellt Section-Divider-Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # Sekundäre Farbe als Hintergrund
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(self.colors["secondary"])
        
        # Section Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb("#FFFFFF")
        p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _add_content_slide(self, data: Dict):
        """Erstellt Content-Slide mit Titel und Text."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide)
        
        # Accent bar oben
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(self.colors["primary"])
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
        
        # Content
        content = data.get("content", "") or data.get("body", "")
        if content:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.2)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            
            lines = content.split('\n')
            for i, line in enumerate(lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                line = line.strip()
                if line.startswith('- ') or line.startswith('* '):
                    p.text = "• " + line[2:]
                    p.level = 0
                else:
                    p.text = line
                
                p.font.size = Pt(18)
                p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
                p.space_after = Pt(12)
        
        return slide
    
    def _add_bullet_slide(self, data: Dict):
        """Erstellt Bullet-Point-Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide)
        
        # Accent bar links
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(0.15), Inches(7.5)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(self.colors["primary"])
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
        
        # Bullets
        bullets = data.get("bullets", [])
        if not bullets and data.get("content"):
            bullets = [l.strip().lstrip('- *') for l in data["content"].split('\n') if l.strip()]
        
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.2)
            )
            tf = content_box.text_frame
            
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                
                p.text = "• " + bullet
                p.font.size = Pt(20)
                p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
                p.space_after = Pt(16)
        
        return slide
    
    def _add_two_column_slide(self, data: Dict):
        """Erstellt Two-Column-Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
        
        # Left Column
        left_content = data.get("left_content", "")
        if not left_content and data.get("content"):
            left_content = data["content"][:len(data["content"])//2]
        
        left_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(5.5), Inches(5.2)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = left_content
        p.font.size = Pt(16)
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
        
        # Right Column
        right_box = slide.shapes.add_textbox(
            Inches(7.083), Inches(1.5), Inches(5.5), Inches(5.2)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("right_content", "")
        p.font.size = Pt(16)
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
        
        # Divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(6.583), Inches(1.5), Inches(0.02), Inches(5)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self._hex_to_rgb(self.colors["secondary"])
        line.line.fill.background()
        
        return slide
    
    def _add_quote_slide(self, data: Dict):
        """Erstellt Quote-Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        # Dunkler Hintergrund
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(self.colors["text"])
        
        # Quote marks
        quote_mark = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(2), Inches(1.5)
        )
        tf = quote_mark.text_frame
        p = tf.paragraphs[0]
        p.text = '"'
        p.font.size = Pt(120)
        p.font.color.rgb = self._hex_to_rgb(self.colors["accent"])
        
        # Quote text
        quote_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(11.333), Inches(3)
        )
        tf = quote_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = data.get("quote", data.get("content", ""))
        p.font.size = Pt(28)
        p.font.italic = True
        p.font.color.rgb = self._hex_to_rgb("#FFFFFF")
        p.alignment = PP_ALIGN.CENTER
        
        # Attribution
        if data.get("author"):
            attr_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5), Inches(11.333), Inches(0.5)
            )
            tf = attr_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"— {data['author']}"
            p.font.size = Pt(18)
            p.font.color.rgb = self._hex_to_rgb(self.colors["secondary"])
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def _add_data_slide(self, data: Dict):
        """Erstellt Daten/Chart-Slide mit Platzhalter."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = data.get("title", "Daten")
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
        
        # Data points
        data_points = data.get("data_points", [])
        if data_points:
            y_pos = 1.8
            for dp in data_points[:4]:  # Max 4 Data Points
                # Value box
                val_box = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(1), Inches(y_pos), Inches(3), Inches(1.2)
                )
                val_box.fill.solid()
                val_box.fill.fore_color.rgb = self._hex_to_rgb(self.colors["primary"])
                
                # Value text
                tf = val_box.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.text = str(dp.get("value", ""))
                p.font.size = Pt(36)
                p.font.bold = True
                p.font.color.rgb = self._hex_to_rgb("#FFFFFF")
                p.alignment = PP_ALIGN.CENTER
                
                # Label
                label_box = slide.shapes.add_textbox(
                    Inches(4.5), Inches(y_pos + 0.3), Inches(8), Inches(0.6)
                )
                tf = label_box.text_frame
                p = tf.paragraphs[0]
                p.text = dp.get("label", "")
                p.font.size = Pt(20)
                p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
                
                y_pos += 1.4
        else:
            # Fallback: Zeige Content als Text
            content = data.get("content", "")
            if content:
                content_box = slide.shapes.add_textbox(
                    Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.2)
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(18)
                p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
        
        return slide
    

    # ─────────────────────────────────────────
    # SPRINT 2: AUTO IMAGES
    # ─────────────────────────────────────────

    def _try_add_image(self, slide, data: Dict):
        """
        Versucht ein passendes Bild einzufügen.
        Priorität: 1. slide["image"] explizit gesetzt
                   2. resolve_for_slide() aus image_store
        Bild wird rechts oben eingefügt (2.5" x 2.5") wenn Platz ist.
        """
        # _SPRINT2_PATCH_
        img_path = None

        # 1. Explizit im Slide-Dict
        explicit = data.get("image") or data.get("image_path")
        if explicit:
            p = Path(explicit)
            if p.exists():
                img_path = p

        # 2. image_store resolve
        if not img_path and self._customer_name:
            try:
                from services.image_store import resolve_for_slide
                tokens = []
                title = data.get("title", "")
                bullets = data.get("bullets", [])
                # Tokens aus Titel + ersten Bullets
                words = (title + " " + " ".join(str(b) for b in bullets[:2])).split()
                tokens = [w.lower().strip(".,!?:;") for w in words if len(w) > 4][:8]
                topic = data.get("topic") or self._topic
                subtopic = data.get("subtopic") or data.get("type", "")
                resolved = resolve_for_slide(
                    self._customer_name, topic, subtopic, tokens
                )
                if resolved and resolved.exists():
                    img_path = resolved
            except Exception as e:
                logger.debug("Auto image resolve failed: %s", e)

        if not img_path:
            return

        # Bild rechts oben einfügen (overlay, klein)
        try:
            from pptx.util import Inches as _In
            W, H = _In(2.8), _In(2.8)
            LEFT = self.prs.slide_width - W - _In(0.3)
            TOP = _In(0.5)
            slide.shapes.add_picture(str(img_path), LEFT, TOP, W, H)
            logger.debug("Auto image inserted: %s", img_path.name)
        except Exception as e:
            logger.debug("Image insert failed: %s", e)

    # ─────────────────────────────────────────
    # SPRINT 2: NEUE LAYOUT-TYPEN
    # ─────────────────────────────────────────

    def _add_agenda_slide(self, data: Dict):
        """Agenda-Slide mit nummerierten Punkten."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._add_background(slide)

        # Linke Farbfläche
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.5), Inches(7.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = self._hex_to_rgb(self.colors["primary"]); bar.line.fill.background()

        # Titel
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.8), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = data.get("title", "Agenda"); p.font.size = Pt(32); p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["primary"])

        # Trennlinie
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.25), Inches(11.5), Inches(0.04))
        line.fill.solid(); line.fill.fore_color.rgb = self._hex_to_rgb(self.colors["accent"]); line.line.fill.background()

        # Nummerierte Punkte
        bullets = data.get("bullets", [])
        y = 1.5
        for i, bullet in enumerate(bullets[:8]):
            # Nummer-Badge
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), Inches(y), Inches(0.4), Inches(0.4))
            badge.fill.solid(); badge.fill.fore_color.rgb = self._hex_to_rgb(self.colors["accent"]); badge.line.fill.background()
            bt = badge.text_frame.paragraphs[0]
            bt.text = str(i + 1); bt.font.size = Pt(12); bt.font.bold = True
            bt.font.color.rgb = self._hex_to_rgb("#FFFFFF"); bt.alignment = PP_ALIGN.CENTER

            # Text
            tb2 = slide.shapes.add_textbox(Inches(1.3), Inches(y - 0.05), Inches(11.0), Inches(0.5))
            p2 = tb2.text_frame.paragraphs[0]
            p2.text = str(bullet); p2.font.size = Pt(20)
            p2.font.color.rgb = self._hex_to_rgb(self.colors["text"])
            y += 0.65

        return slide

    def _add_cta_slide(self, data: Dict):
        """CTA / Nächste Schritte Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._add_background(slide, use_primary=True)

        # Titel
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.333), Inches(1))
        p = tb.text_frame.paragraphs[0]
        p.text = data.get("title", "Nächste Schritte")
        p.font.size = Pt(40); p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb("#FFFFFF"); p.alignment = PP_ALIGN.CENTER

        # Accent-Linie
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.9), Inches(3.333), Inches(0.06))
        line.fill.solid(); line.fill.fore_color.rgb = self._hex_to_rgb(self.colors["accent"]); line.line.fill.background()

        # Steps als Kästen
        bullets = data.get("bullets", [])
        n = min(len(bullets), 4)
        if n > 0:
            box_w = 11.5 / n
            for i, step in enumerate(bullets[:4]):
                x = 0.75 + i * box_w
                box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.5), Inches(box_w - 0.2), Inches(3.5))
                box.fill.solid(); box.fill.fore_color.rgb = self._hex_to_rgb("#FFFFFF"); box.line.fill.background()

                num_tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(2.7), Inches(box_w - 0.4), Inches(0.5))
                np_ = num_tb.text_frame.paragraphs[0]
                np_.text = str(i + 1); np_.font.size = Pt(28); np_.font.bold = True
                np_.font.color.rgb = self._hex_to_rgb(self.colors["accent"]); np_.alignment = PP_ALIGN.CENTER

                step_tb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(3.4), Inches(box_w - 0.4), Inches(2.3))
                step_tb.text_frame.word_wrap = True
                sp = step_tb.text_frame.paragraphs[0]
                sp.text = str(step); sp.font.size = Pt(16)
                sp.font.color.rgb = self._hex_to_rgb(self.colors["primary"]); sp.alignment = PP_ALIGN.CENTER

        return slide

    def _add_kpi_slide(self, data: Dict):
        """KPI / Metriken Slide mit großen Zahlen."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._add_background(slide)

        # Oben: Balken
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = self._hex_to_rgb(self.colors["primary"]); bar.line.fill.background()

        # Titel
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.35), Inches(11.8), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = data.get("title", "KPIs & Metriken"); p.font.size = Pt(32); p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])

        # KPI-Kästen (aus bullets: "Label: Wert")
        bullets = data.get("bullets", [])
        kpis = []
        for b in bullets:
            b = str(b)
            if ":" in b:
                parts = b.split(":", 1)
                kpis.append({"label": parts[0].strip(), "value": parts[1].strip()})
            else:
                kpis.append({"label": b, "value": "—"})

        n = min(len(kpis), 4)
        if n > 0:
            box_w = 11.5 / n
            for i, kpi in enumerate(kpis[:4]):
                x = 0.75 + i * box_w
                box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.8), Inches(box_w - 0.2), Inches(4.5))
                box.fill.solid(); box.fill.fore_color.rgb = self._hex_to_rgb(self.colors["primary"] + "18"); box.line.fill.background()

                # Wert groß
                val_tb = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(box_w - 0.2), Inches(2))
                vp = val_tb.text_frame.paragraphs[0]
                vp.text = kpi["value"]; vp.font.size = Pt(44); vp.font.bold = True
                vp.font.color.rgb = self._hex_to_rgb(self.colors["accent"]); vp.alignment = PP_ALIGN.CENTER

                # Label klein
                lbl_tb = slide.shapes.add_textbox(Inches(x), Inches(4.4), Inches(box_w - 0.2), Inches(1))
                lp = lbl_tb.text_frame.paragraphs[0]
                lp.text = kpi["label"]; lp.font.size = Pt(16)
                lp.font.color.rgb = self._hex_to_rgb(self.colors["text"]); lp.alignment = PP_ALIGN.CENTER

        return slide

    def _add_swot_slide(self, data: Dict):
        """2x2 SWOT-Matrix Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._add_background(slide)

        # Titel
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.7))
        p = tb.text_frame.paragraphs[0]
        p.text = data.get("title", "SWOT-Analyse"); p.font.size = Pt(28); p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])

        # 4 Quadranten
        bullets = data.get("bullets", [])
        quadrants = [
            ("S – Stärken", self.colors["accent"], bullets[0] if len(bullets) > 0 else ""),
            ("W – Schwächen", "#EF4444", bullets[1] if len(bullets) > 1 else ""),
            ("O – Chancen", "#3B82F6", bullets[2] if len(bullets) > 2 else ""),
            ("T – Risiken", "#F59E0B", bullets[3] if len(bullets) > 3 else ""),
        ]
        positions = [(0.4, 1.0), (6.8, 1.0), (0.4, 4.2), (6.8, 4.2)]

        for (label, color, text), (x, y) in zip(quadrants, positions):
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(6.2), Inches(3.0))
            box.fill.solid(); box.fill.fore_color.rgb = self._hex_to_rgb(color + "22"); box.line.color.rgb = self._hex_to_rgb(color)

            hdr = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.1), Inches(6.0), Inches(0.5))
            hp = hdr.text_frame.paragraphs[0]
            hp.text = label; hp.font.size = Pt(16); hp.font.bold = True
            hp.font.color.rgb = self._hex_to_rgb(color)

            if text:
                tb2 = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.65), Inches(6.0), Inches(2.1))
                tb2.text_frame.word_wrap = True
                tp = tb2.text_frame.paragraphs[0]
                tp.text = str(text); tp.font.size = Pt(14)
                tp.font.color.rgb = self._hex_to_rgb(self.colors["text"])

        return slide

    def _add_timeline_slide(self, data: Dict):
        """Timeline / Roadmap Slide."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._add_background(slide)

        # Titel
        tb = slide.shapes.add_textbox(Inches(0.75), Inches(0.3), Inches(11.8), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = data.get("title", "Roadmap"); p.font.size = Pt(32); p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])

        # Horizontale Zeitlinie
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(3.6), Inches(11.8), Inches(0.08))
        line.fill.solid(); line.fill.fore_color.rgb = self._hex_to_rgb(self.colors["primary"]); line.line.fill.background()

        bullets = data.get("bullets", [])
        n = min(len(bullets), 6)
        if n > 0:
            spacing = 11.5 / n
            for i, step in enumerate(bullets[:6]):
                x = 0.75 + i * spacing

                # Punkt auf der Linie
                dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + spacing/2 - 0.15), Inches(3.44), Inches(0.3), Inches(0.3))
                dot.fill.solid(); dot.fill.fore_color.rgb = self._hex_to_rgb(self.colors["accent"]); dot.line.fill.background()

                # Text abwechselnd oben/unten
                if i % 2 == 0:
                    ty = 1.5
                else:
                    ty = 4.2

                step_tb = slide.shapes.add_textbox(Inches(x), Inches(ty), Inches(spacing - 0.1), Inches(1.8))
                step_tb.text_frame.word_wrap = True
                sp = step_tb.text_frame.paragraphs[0]
                sp.text = str(step); sp.font.size = Pt(14)
                sp.font.color.rgb = self._hex_to_rgb(self.colors["text"]); sp.alignment = PP_ALIGN.CENTER

        return slide

    def _add_image_slide(self, data: Dict):
        """Bild-dominierter Slide mit Titel und Caption."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        self._add_background(slide)

        # Titel
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.8))
        p = tb.text_frame.paragraphs[0]
        p.text = data.get("title", ""); p.font.size = Pt(28); p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])

        # Bild laden
        img_path = None
        explicit = data.get("image") or data.get("image_path")
        if explicit:
            p2 = Path(explicit)
            if p2.exists():
                img_path = p2
        if not img_path and self._customer_name:
            try:
                from services.image_store import resolve_for_slide
                title = data.get("title", self._topic)
                resolved = resolve_for_slide(self._customer_name, title, "", [])
                if resolved and resolved.exists():
                    img_path = resolved
            except Exception:
                pass

        if img_path:
            try:
                slide.shapes.add_picture(str(img_path), Inches(1.5), Inches(1.2), Inches(10.333), Inches(5.5))
            except Exception as e:
                logger.debug("Image slide picture failed: %s", e)

        # Caption aus Bullets
        bullets = data.get("bullets", [])
        if bullets:
            cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.4))
            cp = cap.text_frame.paragraphs[0]
            cp.text = " | ".join(str(b) for b in bullets[:3]); cp.font.size = Pt(12)
            cp.font.color.rgb = self._hex_to_rgb(self.colors["secondary"]); cp.alignment = PP_ALIGN.CENTER

        return slide

    def _add_sources_overview_slide(self, sources: List[str]):
        """Erstellt Quellenübersicht-Slide am Ende."""
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        
        self._add_background(slide)
        
        # Accent bar
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(0.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._hex_to_rgb(self.colors["primary"])
        bar.line.fill.background()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Quellen & Referenzen"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
        
        # Sources list
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.5), Inches(11.833), Inches(5.5)
        )
        tf = content_box.text_frame
        
        for i, source in enumerate(sources[:15]):  # Max 15 Quellen
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {source}"
            p.font.size = Pt(14)
            p.font.color.rgb = self._hex_to_rgb(self.colors["text"])
            p.space_after = Pt(8)
        
        if len(sources) > 15:
            p = tf.add_paragraph()
            p.text = f"... und {len(sources) - 15} weitere Quellen"
            p.font.size = Pt(12)
            p.font.color.rgb = self._hex_to_rgb(self.colors["secondary"])
        
        return slide


def create_presentation_v2(slides: List[Dict], 
                          title: str = "Präsentation",
                          company: str = "",
                          colors: Dict[str, str] = None,
                          palette: str = "corporate",
                          include_sources: bool = True) -> bytes:
    """
    Erstellt PPTX mit Wizard-Farben und Quellenangaben.
    
    Args:
        slides: Liste von Slide-Dicts mit type, title, content, sources
        title: Präsentationstitel
        company: Firmenname
        colors: Wizard-Farben (primary, secondary, accent, background, text)
        palette: Fallback-Palette
        include_sources: Quellenübersicht am Ende
    
    Returns:
        PPTX als Bytes
    """
    designer = PPTXDesignerV2(colors=colors, palette=palette)
    return designer.create_presentation(slides, title, company, include_sources)
