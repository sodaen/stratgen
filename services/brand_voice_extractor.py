# -*- coding: utf-8 -*-
"""
services/brand_voice_extractor.py
=================================
Feature: Brand Voice Extractor

Analysiert Texte aus Templates und lernt:
1. Tonalität (formal/informell)
2. Typische Verben, Adjektive
3. Satzlängen, Komplexität
4. Branchenspezifische Formulierungen
5. Generiert dann in DEINEM Stil

Author: StratGen Agent V3.6
"""
from __future__ import annotations
import os
import re
import json
import sqlite3
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field, asdict
from collections import Counter
from datetime import datetime

# ============================================
# CONFIGURATION
# ============================================

DB_PATH = os.getenv("STRATGEN_VOICE_DB", "data/brand_voice.sqlite")

# ============================================
# DATA CLASSES
# ============================================

@dataclass
class VoiceProfile:
    """Ein Brand Voice Profil."""
    name: str
    
    # Tonalität
    formality_score: float  # 0=informal, 1=formal
    confidence_score: float  # 0=vorsichtig, 1=bestimmt
    emotion_score: float  # 0=sachlich, 1=emotional
    
    # Satzstruktur
    avg_sentence_length: float
    avg_word_length: float
    sentence_complexity: float  # Anteil komplexer Sätze
    
    # Vokabular
    top_verbs: List[str]
    top_adjectives: List[str]
    top_nouns: List[str]
    power_words: List[str]
    avoided_words: List[str]
    
    # Muster
    typical_openings: List[str]
    typical_closings: List[str]
    transition_phrases: List[str]
    
    # Statistiken
    templates_analyzed: int
    total_words_analyzed: int


@dataclass
class TextAnalysis:
    """Analyse eines einzelnen Textes."""
    word_count: int
    sentence_count: int
    avg_sentence_length: float
    formality_indicators: int
    informal_indicators: int
    verbs: List[str]
    adjectives: List[str]
    nouns: List[str]


# ============================================
# LINGUISTIC PATTERNS
# ============================================

# Formale Indikatoren
FORMAL_INDICATORS = [
    "gemäß", "hinsichtlich", "bezüglich", "diesbezüglich", "entsprechend",
    "daher", "folglich", "somit", "demzufolge", "infolgedessen",
    "gewährleisten", "sicherstellen", "ermöglichen", "realisieren",
    "Implementierung", "Optimierung", "Evaluierung", "Strategie"
]

# Informale Indikatoren
INFORMAL_INDICATORS = [
    "super", "toll", "cool", "easy", "okay", "klar", "echt",
    "halt", "mal", "einfach", "irgendwie", "quasi"
]

# Power Words
POWER_WORDS = [
    "revolutionär", "bahnbrechend", "führend", "exklusiv", "garantiert",
    "sofort", "kostenlos", "neu", "bewährt", "erfolgreich", "professionell",
    "innovativ", "effizient", "nachhaltig", "strategisch", "intelligent"
]

# Confidence Indicators
CONFIDENCE_HIGH = [
    "werden", "ist", "garantiert", "sicher", "definitiv", "nachweislich",
    "bewährt", "erfolgreich", "führend"
]

CONFIDENCE_LOW = [
    "könnte", "möglicherweise", "eventuell", "vielleicht", "unter Umständen",
    "gegebenenfalls", "tendenziell"
]

# Typische deutsche Verben in Business-Kontexten
BUSINESS_VERBS = [
    "optimieren", "steigern", "verbessern", "entwickeln", "implementieren",
    "analysieren", "evaluieren", "transformieren", "digitalisieren",
    "automatisieren", "skalieren", "integrieren", "realisieren"
]

# ============================================
# DATABASE
# ============================================

def _connect():
    """Verbindet zur Voice-Datenbank."""
    os.makedirs("data", exist_ok=True)
    con = sqlite3.connect(DB_PATH, timeout=30)
    con.row_factory = sqlite3.Row
    return con


def ensure_tables():
    """Erstellt Tabellen falls nicht vorhanden."""
    con = _connect()
    cur = con.cursor()
    
    # Voice Profile
    cur.execute("""
    CREATE TABLE IF NOT EXISTS voice_profiles (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT NOT NULL UNIQUE,
        profile_data TEXT,  -- JSON
        created_at TEXT DEFAULT CURRENT_TIMESTAMP,
        updated_at TEXT DEFAULT CURRENT_TIMESTAMP
    )
    """)
    
    # Vokabular
    cur.execute("""
    CREATE TABLE IF NOT EXISTS vocabulary_stats (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        profile_name TEXT,
        word TEXT,
        word_type TEXT,  -- verb, adjective, noun, other
        frequency INTEGER DEFAULT 1,
        UNIQUE(profile_name, word)
    )
    """)
    
    # Phrasen
    cur.execute("""
    CREATE TABLE IF NOT EXISTS phrase_stats (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        profile_name TEXT,
        phrase TEXT,
        phrase_type TEXT,  -- opening, closing, transition
        frequency INTEGER DEFAULT 1,
        UNIQUE(profile_name, phrase)
    )
    """)
    
    con.commit()
    con.close()


# ============================================
# TEXT ANALYSIS
# ============================================

def analyze_text(text: str) -> TextAnalysis:
    """
    Analysiert einen einzelnen Text.
    
    Args:
        text: Der zu analysierende Text
    
    Returns:
        TextAnalysis-Objekt
    """
    # Sätze trennen
    sentences = re.split(r'[.!?]+', text)
    sentences = [s.strip() for s in sentences if s.strip()]
    
    # Wörter
    words = re.findall(r'\b[A-ZÄÖÜa-zäöüß]+\b', text)
    
    # Formality Score
    formal_count = sum(1 for w in FORMAL_INDICATORS if w.lower() in text.lower())
    informal_count = sum(1 for w in INFORMAL_INDICATORS if w.lower() in text.lower())
    
    # Wortarten erkennen (vereinfacht)
    verbs = [w for w in words if w.lower().endswith(('en', 'ern', 'eln', 'ieren'))]
    adjectives = [w for w in words if w.lower().endswith(('ig', 'lich', 'isch', 'bar', 'sam', 'haft'))]
    
    # Nomen (Großgeschrieben, nicht am Satzanfang)
    nouns = []
    for i, w in enumerate(words):
        if w[0].isupper() and len(w) > 3:
            # Prüfe ob am Satzanfang
            if i > 0 and words[i-1][-1] not in '.!?':
                nouns.append(w)
    
    avg_sentence_length = len(words) / max(1, len(sentences))
    
    return TextAnalysis(
        word_count=len(words),
        sentence_count=len(sentences),
        avg_sentence_length=avg_sentence_length,
        formality_indicators=formal_count,
        informal_indicators=informal_count,
        verbs=verbs,
        adjectives=adjectives,
        nouns=nouns
    )


def extract_phrases(text: str) -> Dict[str, List[str]]:
    """
    Extrahiert typische Phrasen aus Text.
    
    Returns:
        Dictionary mit openings, closings, transitions
    """
    sentences = re.split(r'[.!?]+', text)
    sentences = [s.strip() for s in sentences if s.strip() and len(s) > 20]
    
    result = {
        "openings": [],
        "closings": [],
        "transitions": []
    }
    
    # Openings (erste Wörter von Sätzen)
    for sent in sentences[:10]:
        words = sent.split()[:4]
        if len(words) >= 3:
            result["openings"].append(" ".join(words))
    
    # Transition Patterns
    transition_patterns = [
        r'(Darüber hinaus|Zusätzlich|Weiterhin|Außerdem|Des Weiteren)',
        r'(Im Gegensatz dazu|Andererseits|Jedoch|Allerdings)',
        r'(Zusammenfassend|Abschließend|Insgesamt|Letztlich)',
        r'(Das bedeutet|Das heißt|Mit anderen Worten)',
    ]
    
    for pattern in transition_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        result["transitions"].extend(matches)
    
    return result


# ============================================
# PROFILE BUILDING
# ============================================

def analyze_template_voice(file_path: str) -> Dict[str, Any]:
    """
    Analysiert die Voice eines Templates.
    
    Args:
        file_path: Pfad zur PPTX-Datei
    
    Returns:
        Voice-Analyse
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed"}
    
    try:
        prs = Presentation(file_path)
    except Exception as e:
        return {"ok": False, "error": str(e)}
    
    all_text = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text and len(text) > 10:
                        all_text.append(text)
    
    full_text = " ".join(all_text)
    
    if len(full_text) < 100:
        return {"ok": False, "error": "Not enough text to analyze"}
    
    # Analyse
    analysis = analyze_text(full_text)
    phrases = extract_phrases(full_text)
    
    # Formality Score berechnen
    total_indicators = analysis.formality_indicators + analysis.informal_indicators
    formality = analysis.formality_indicators / max(1, total_indicators)
    
    # Confidence Score
    confidence_high = sum(1 for w in CONFIDENCE_HIGH if w in full_text.lower())
    confidence_low = sum(1 for w in CONFIDENCE_LOW if w in full_text.lower())
    confidence = confidence_high / max(1, confidence_high + confidence_low)
    
    # Power Words
    found_power_words = [w for w in POWER_WORDS if w.lower() in full_text.lower()]
    
    return {
        "ok": True,
        "file": file_path,
        "word_count": analysis.word_count,
        "sentence_count": analysis.sentence_count,
        "avg_sentence_length": round(analysis.avg_sentence_length, 1),
        "formality_score": round(formality, 2),
        "confidence_score": round(confidence, 2),
        "top_verbs": Counter(v.lower() for v in analysis.verbs).most_common(10),
        "top_adjectives": Counter(a.lower() for a in analysis.adjectives).most_common(10),
        "top_nouns": Counter(n for n in analysis.nouns).most_common(10),
        "power_words_used": found_power_words,
        "phrases": phrases
    }


def build_voice_profile(profile_name: str = "default") -> VoiceProfile:
    """
    Baut ein Voice Profile aus allen analysierten Templates.
    
    Args:
        profile_name: Name des Profils
    
    Returns:
        VoiceProfile-Objekt
    """
    from pathlib import Path
    
    raw_dir = Path(os.getenv("STRATGEN_RAW_DIR", "data/raw"))
    
    all_analyses = []
    all_verbs = []
    all_adjectives = []
    all_nouns = []
    all_power_words = []
    all_openings = []
    all_transitions = []
    
    for pptx_file in raw_dir.glob("**/*.pptx"):
        result = analyze_template_voice(str(pptx_file))
        if result.get("ok"):
            all_analyses.append(result)
            all_verbs.extend([v[0] for v in result.get("top_verbs", [])])
            all_adjectives.extend([a[0] for a in result.get("top_adjectives", [])])
            all_nouns.extend([n[0] for n in result.get("top_nouns", [])])
            all_power_words.extend(result.get("power_words_used", []))
            all_openings.extend(result.get("phrases", {}).get("openings", []))
            all_transitions.extend(result.get("phrases", {}).get("transitions", []))
    
    if not all_analyses:
        return VoiceProfile(
            name=profile_name,
            formality_score=0.7,
            confidence_score=0.7,
            emotion_score=0.3,
            avg_sentence_length=15,
            avg_word_length=6,
            sentence_complexity=0.3,
            top_verbs=BUSINESS_VERBS[:10],
            top_adjectives=["effizient", "innovativ", "strategisch"],
            top_nouns=["Unternehmen", "Lösung", "Strategie"],
            power_words=POWER_WORDS[:10],
            avoided_words=INFORMAL_INDICATORS[:5],
            typical_openings=["Wir bieten", "Unsere Lösung"],
            typical_closings=["Kontaktieren Sie uns"],
            transition_phrases=["Darüber hinaus", "Zusätzlich"],
            templates_analyzed=0,
            total_words_analyzed=0
        )
    
    # Aggregieren
    avg_formality = sum(a["formality_score"] for a in all_analyses) / len(all_analyses)
    avg_confidence = sum(a["confidence_score"] for a in all_analyses) / len(all_analyses)
    avg_sentence_len = sum(a["avg_sentence_length"] for a in all_analyses) / len(all_analyses)
    total_words = sum(a["word_count"] for a in all_analyses)
    
    # Top Items
    verb_counter = Counter(all_verbs)
    adj_counter = Counter(all_adjectives)
    noun_counter = Counter(all_nouns)
    
    profile = VoiceProfile(
        name=profile_name,
        formality_score=round(avg_formality, 2),
        confidence_score=round(avg_confidence, 2),
        emotion_score=0.3,  # TODO: besser berechnen
        avg_sentence_length=round(avg_sentence_len, 1),
        avg_word_length=6.0,  # TODO: berechnen
        sentence_complexity=0.4,  # TODO: berechnen
        top_verbs=[v[0] for v in verb_counter.most_common(15)],
        top_adjectives=[a[0] for a in adj_counter.most_common(15)],
        top_nouns=[n[0] for n in noun_counter.most_common(15)],
        power_words=list(set(all_power_words))[:15],
        avoided_words=INFORMAL_INDICATORS[:10],
        typical_openings=list(set(all_openings))[:10],
        typical_closings=["Kontaktieren Sie uns", "Wir freuen uns auf Sie"],
        transition_phrases=list(set(all_transitions))[:10],
        templates_analyzed=len(all_analyses),
        total_words_analyzed=total_words
    )
    
    # In DB speichern
    save_profile(profile)
    
    return profile


def save_profile(profile: VoiceProfile):
    """Speichert ein Voice Profile in der Datenbank."""
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    
    cur.execute("""
    INSERT OR REPLACE INTO voice_profiles (name, profile_data, updated_at)
    VALUES (?, ?, CURRENT_TIMESTAMP)
    """, (profile.name, json.dumps(asdict(profile))))
    
    con.commit()
    con.close()


def load_profile(name: str = "default") -> Optional[VoiceProfile]:
    """Lädt ein Voice Profile aus der Datenbank."""
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    
    cur.execute("SELECT profile_data FROM voice_profiles WHERE name = ?", (name,))
    row = cur.fetchone()
    con.close()
    
    if row:
        data = json.loads(row["profile_data"])
        return VoiceProfile(**data)
    
    return None


# ============================================
# TEXT GENERATION GUIDANCE
# ============================================

def get_writing_guidelines(profile_name: str = "default") -> Dict[str, Any]:
    """
    Gibt Schreibrichtlinien basierend auf dem Voice Profile.
    
    Returns:
        Dictionary mit Guidelines für LLM-Prompts
    """
    profile = load_profile(profile_name)
    
    if not profile:
        profile = build_voice_profile(profile_name)
    
    # Tonalitäts-Beschreibung
    if profile.formality_score > 0.7:
        tone = "formal und professionell"
    elif profile.formality_score > 0.4:
        tone = "professionell aber zugänglich"
    else:
        tone = "locker und direkt"
    
    # Confidence-Beschreibung
    if profile.confidence_score > 0.7:
        confidence = "bestimmt und überzeugend"
    elif profile.confidence_score > 0.4:
        confidence = "ausgewogen"
    else:
        confidence = "vorsichtig und abwägend"
    
    return {
        "ok": True,
        "profile_name": profile.name,
        "guidelines": {
            "tone": tone,
            "confidence": confidence,
            "sentence_length": f"Durchschnittlich {int(profile.avg_sentence_length)} Wörter pro Satz",
            "preferred_verbs": profile.top_verbs[:10],
            "preferred_adjectives": profile.top_adjectives[:10],
            "power_words_to_use": profile.power_words[:10],
            "words_to_avoid": profile.avoided_words,
            "typical_openings": profile.typical_openings[:5],
            "transition_phrases": profile.transition_phrases[:5]
        },
        "prompt_instruction": f"""
Schreibe im folgenden Stil:
- Tonalität: {tone}
- Formulierung: {confidence}
- Satzlänge: ca. {int(profile.avg_sentence_length)} Wörter
- Verwende bevorzugt diese Verben: {', '.join(profile.top_verbs[:5])}
- Verwende diese Power-Words: {', '.join(profile.power_words[:5])}
- Vermeide: {', '.join(profile.avoided_words[:5])}
""",
        "statistics": {
            "templates_analyzed": profile.templates_analyzed,
            "words_analyzed": profile.total_words_analyzed
        }
    }


def adapt_text_to_voice(text: str, profile_name: str = "default") -> Dict[str, Any]:
    """
    Schlägt Anpassungen vor um Text an Voice anzupassen.
    
    Args:
        text: Der anzupassende Text
        profile_name: Name des Voice Profiles
    
    Returns:
        Dictionary mit Anpassungsvorschlägen
    """
    profile = load_profile(profile_name)
    if not profile:
        return {"ok": False, "error": "Profile not found"}
    
    suggestions = []
    
    # Prüfe auf informelle Wörter
    for word in profile.avoided_words:
        if word.lower() in text.lower():
            suggestions.append({
                "type": "avoid_word",
                "word": word,
                "suggestion": f"'{word}' vermeiden - zu informell"
            })
    
    # Prüfe Satzlänge
    sentences = re.split(r'[.!?]+', text)
    for sent in sentences:
        words = sent.split()
        if len(words) > profile.avg_sentence_length * 1.5:
            suggestions.append({
                "type": "sentence_length",
                "sentence": sent[:50] + "...",
                "suggestion": f"Satz aufteilen (aktuell {len(words)} Wörter, Ziel: ~{int(profile.avg_sentence_length)})"
            })
    
    # Schlage Power Words vor
    missing_power = [pw for pw in profile.power_words[:5] if pw.lower() not in text.lower()]
    if missing_power:
        suggestions.append({
            "type": "power_words",
            "suggestion": f"Erwäge Power-Words: {', '.join(missing_power[:3])}"
        })
    
    return {
        "ok": True,
        "original_text": text[:200],
        "suggestions": suggestions,
        "voice_match_score": max(0, 1 - len(suggestions) * 0.1)
    }


# ============================================
# API FUNCTIONS
# ============================================

def check_status() -> Dict[str, Any]:
    """Gibt den Status des Voice Extractors zurück."""
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    
    cur.execute("SELECT COUNT(*) as cnt FROM voice_profiles")
    profiles = cur.fetchone()["cnt"]
    
    con.close()
    
    profile = load_profile("default")
    
    return {
        "ok": True,
        "profiles_stored": profiles,
        "default_profile_exists": profile is not None,
        "templates_analyzed": profile.templates_analyzed if profile else 0
    }
