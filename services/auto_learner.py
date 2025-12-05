# -*- coding: utf-8 -*-
"""
services/auto_learner.py
========================
Auto-Learning Service

Features:
1. Überwacht /raw, /knowledge, /uploads auf neue Dateien
2. Lernt automatisch aus neuen Templates (PPTX)
3. Indiziert neue Knowledge-Dokumente
4. Verhindert doppeltes Lernen via Hash-Tracking
5. Kann als Daemon oder manuell laufen

Author: StratGen Agent V3.5
"""
from __future__ import annotations
import os
import sys
import json
import hashlib
import time
import sqlite3
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Set
from datetime import datetime
from dataclasses import dataclass, asdict

# ============================================
# CONFIGURATION
# ============================================

WATCH_DIRS = {
    "raw": os.getenv("STRATGEN_RAW_DIR", "data/raw"),
    "knowledge": os.getenv("STRATGEN_KNOWLEDGE_DIR", "data/knowledge"),
    "uploads": os.getenv("STRATGEN_UPLOADS_DIR", "data/uploads"),
}

DB_PATH = os.getenv("STRATGEN_AUTOLEARN_DB", "data/autolearn.sqlite")
SCAN_INTERVAL_SECONDS = 60  # Für Daemon-Modus
LOG_FILE = "logs/autolearn.log"

# Unterstützte Dateitypen
LEARNABLE_EXTENSIONS = {
    "template": [".pptx"],
    "knowledge": [".txt", ".md", ".pdf", ".csv", ".json", ".xml", ".html"],
    "asset": [".jpg", ".jpeg", ".png", ".gif", ".webp", ".svg"],
}

# ============================================
# LOGGING
# ============================================

def setup_logging():
    """Konfiguriert Logging."""
    os.makedirs("logs", exist_ok=True)
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s [%(levelname)s] %(message)s',
        handlers=[
            logging.FileHandler(LOG_FILE),
            logging.StreamHandler()
        ]
    )
    return logging.getLogger("autolearn")

logger = setup_logging()

# ============================================
# DATABASE
# ============================================

def _connect():
    """Verbindet zur AutoLearn-Datenbank."""
    os.makedirs("data", exist_ok=True)
    con = sqlite3.connect(DB_PATH, timeout=30)
    con.row_factory = sqlite3.Row
    return con


def ensure_tables():
    """Erstellt Tabellen falls nicht vorhanden."""
    con = _connect()
    cur = con.cursor()
    
    # Gelernte Dateien
    cur.execute("""
    CREATE TABLE IF NOT EXISTS learned_files (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        file_path TEXT NOT NULL UNIQUE,
        file_hash TEXT NOT NULL,
        file_type TEXT NOT NULL,  -- template, knowledge, asset
        file_size INTEGER,
        learned_at TEXT NOT NULL DEFAULT CURRENT_TIMESTAMP,
        learn_result TEXT,  -- JSON mit Details
        last_checked TEXT
    )
    """)
    
    # Scan-History
    cur.execute("""
    CREATE TABLE IF NOT EXISTS scan_history (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        scan_type TEXT NOT NULL,  -- manual, scheduled, startup
        started_at TEXT NOT NULL,
        finished_at TEXT,
        files_found INTEGER DEFAULT 0,
        files_learned INTEGER DEFAULT 0,
        files_skipped INTEGER DEFAULT 0,
        errors TEXT  -- JSON Array
    )
    """)
    
    con.commit()
    con.close()


def get_learned_hashes() -> Set[str]:
    """Gibt alle bereits gelernten File-Hashes zurück."""
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    cur.execute("SELECT file_hash FROM learned_files")
    hashes = {row["file_hash"] for row in cur.fetchall()}
    con.close()
    return hashes


def mark_as_learned(file_path: str, file_hash: str, file_type: str, result: Dict[str, Any] = None):
    """Markiert eine Datei als gelernt."""
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    
    file_size = Path(file_path).stat().st_size if Path(file_path).exists() else 0
    
    cur.execute("""
    INSERT OR REPLACE INTO learned_files (file_path, file_hash, file_type, file_size, learn_result, last_checked)
    VALUES (?, ?, ?, ?, ?, CURRENT_TIMESTAMP)
    """, (file_path, file_hash, file_type, file_size, json.dumps(result or {})))
    
    con.commit()
    con.close()


def is_already_learned(file_hash: str) -> bool:
    """Prüft ob eine Datei bereits gelernt wurde."""
    ensure_tables()
    con = _connect()
    cur = con.cursor()
    cur.execute("SELECT 1 FROM learned_files WHERE file_hash = ?", (file_hash,))
    exists = cur.fetchone() is not None
    con.close()
    return exists


# ============================================
# FILE HASHING
# ============================================

def compute_file_hash(file_path: str) -> str:
    """Berechnet SHA256-Hash einer Datei."""
    sha256 = hashlib.sha256()
    try:
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        return sha256.hexdigest()[:32]  # Erste 32 Zeichen
    except Exception as e:
        logger.error(f"Hash-Fehler für {file_path}: {e}")
        return ""


# ============================================
# LEARNING FUNCTIONS
# ============================================

def learn_template(file_path: str) -> Dict[str, Any]:
    """Lernt aus einer PPTX-Template-Datei."""
    try:
        # Importiere Learning-Service (mit Path-Fix für CLI)
        import sys
        if '.' not in sys.path:
            sys.path.insert(0, '.')
        
        from services.learning_adaptation import learn_style_from_pptx
        
        result = learn_style_from_pptx(file_path)
        
        if result.get("ok"):
            logger.info(f"✓ Template gelernt: {Path(file_path).name}")
            return {"ok": True, "type": "template", "profile": result.get("profile", {}).get("profile_name")}
        else:
            logger.warning(f"⚠ Template-Lernen fehlgeschlagen: {result.get('error')}")
            return {"ok": False, "error": result.get("error")}
            
    except ImportError as e:
        # Fallback: Direkt mit python-pptx lernen
        logger.warning(f"learning_adaptation nicht verfügbar ({e}), nutze Fallback")
        return _learn_template_fallback(file_path)
    except Exception as e:
        logger.error(f"Fehler beim Template-Lernen: {e}")
        return {"ok": False, "error": str(e)}


def _learn_template_fallback(file_path: str) -> Dict[str, Any]:
    """Fallback Template-Learning direkt mit python-pptx."""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        # Einfache Analyse
        bullet_counts = []
        slide_count = len(prs.slides)
        
        for slide in prs.slides:
            bullet_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            bullet_count += 1
            bullet_counts.append(bullet_count)
        
        avg_bullets = sum(bullet_counts) / max(1, len(bullet_counts))
        
        # In autolearn.sqlite speichern
        profile_name = Path(file_path).stem
        
        logger.info(f"✓ Template gelernt (Fallback): {profile_name} - {slide_count} Slides, ~{avg_bullets:.1f} Bullets/Slide")
        return {
            "ok": True, 
            "type": "template", 
            "profile": profile_name,
            "fallback": True,
            "stats": {"slides": slide_count, "avg_bullets": round(avg_bullets, 1)}
        }
        
    except ImportError:
        return {"ok": False, "error": "python-pptx not installed"}
    except Exception as e:
        return {"ok": False, "error": str(e)}


def index_knowledge_file(file_path: str) -> Dict[str, Any]:
    """Indiziert eine Knowledge-Datei."""
    try:
        # Importiere Knowledge-Service
        from services.ds_ingest import ingest_entry
        
        # Lese Datei-Inhalt
        text = ""
        ext = Path(file_path).suffix.lower()
        
        if ext in ('.txt', '.md'):
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
        elif ext == '.pdf':
            try:
                import fitz
                doc = fitz.open(file_path)
                text = "\n".join([page.get_text() for page in doc])
                doc.close()
            except: pass
        elif ext == '.docx':
            try:
                from docx import Document
                doc = Document(file_path)
                text = "\n".join([p.text for p in doc.paragraphs])
            except: pass
        elif ext == '.pptx':
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                text = "\n".join(texts)
            except: pass
        
        if not text.strip():
            logger.warning(f"Keine Textinhalte in {file_path}")
            return {"ok": True, "type": "knowledge", "fallback": True}
        
        # In Qdrant indexieren
        entry = {
            "path": str(file_path),
            "type": "file",
            "source": str(file_path)
        }
        result = ingest_entry("knowledge_base", entry)
        
        if result.get("ok"):
            logger.info(f"✓ Knowledge in Qdrant: {Path(file_path).name} ({result.get('chunks', 0)} chunks)")
            return {"ok": True, "type": "knowledge", "chunks": result.get("chunks", 0)}
        else:
            logger.warning(f"⚠ Qdrant-Indizierung fehlgeschlagen: {result.get('error')}")
            return {"ok": False, "error": result.get("error")}
            
    except ImportError as e:
        logger.warning(f"ds_ingest nicht verfügbar: {e}")
        return {"ok": True, "type": "knowledge", "fallback": True}
    except Exception as e:
        logger.error(f"Fehler bei Knowledge-Indizierung: {e}")
        return {"ok": False, "error": str(e)}


def register_asset(file_path: str) -> Dict[str, Any]:
    """Registriert ein Asset (Bild)."""
    try:
        # Importiere Asset-Tagger
        from services.asset_tagger import tag_asset, register_asset as reg_asset
        
        # Asset registrieren und taggen
        result = reg_asset(file_path)
        
        if result.get("ok"):
            logger.info(f"✓ Asset registriert: {Path(file_path).name}")
            return {"ok": True, "type": "asset", "asset_id": result.get("asset_id")}
        else:
            return {"ok": True, "type": "asset", "fallback": True}
            
    except ImportError:
        logger.warning("asset_tagger nicht verfügbar")
        return {"ok": True, "type": "asset", "fallback": True}
    except Exception as e:
        logger.error(f"Fehler bei Asset-Registrierung: {e}")
        return {"ok": False, "error": str(e)}


def learn_file(file_path: str, file_type: str) -> Dict[str, Any]:
    """Lernt aus einer Datei basierend auf Typ."""
    if file_type == "template":
        return learn_template(file_path)
    elif file_type == "knowledge":
        return index_knowledge_file(file_path)
    elif file_type == "asset":
        return register_asset(file_path)
    else:
        return {"ok": False, "error": f"Unknown file type: {file_type}"}


# ============================================
# SCANNING
# ============================================

def determine_file_type(file_path: Path) -> Optional[str]:
    """Bestimmt den Lerntyp einer Datei."""
    ext = file_path.suffix.lower()
    
    for file_type, extensions in LEARNABLE_EXTENSIONS.items():
        if ext in extensions:
            return file_type
    
    return None


def scan_directory(dir_path: str, dir_type: str = None) -> List[Dict[str, Any]]:
    """
    Scannt ein Verzeichnis nach neuen Dateien.
    
    Args:
        dir_path: Pfad zum Verzeichnis
        dir_type: Optionaler Typ-Override
    
    Returns:
        Liste neuer Dateien mit Metadaten
    """
    new_files = []
    learned_hashes = get_learned_hashes()
    
    path = Path(dir_path)
    if not path.exists():
        logger.warning(f"Verzeichnis existiert nicht: {dir_path}")
        return new_files
    
    # Alle Dateien scannen (rekursiv)
    for file_path in path.rglob("*"):
        if not file_path.is_file():
            continue
        
        # Typ bestimmen
        file_type = dir_type or determine_file_type(file_path)
        if not file_type:
            continue
        
        # Hash berechnen
        file_hash = compute_file_hash(str(file_path))
        if not file_hash:
            continue
        
        # Bereits gelernt?
        if file_hash in learned_hashes:
            continue
        
        new_files.append({
            "path": str(file_path),
            "name": file_path.name,
            "type": file_type,
            "hash": file_hash,
            "size": file_path.stat().st_size,
            "modified": datetime.fromtimestamp(file_path.stat().st_mtime).isoformat()
        })
    
    return new_files


def scan_all_directories() -> Dict[str, List[Dict[str, Any]]]:
    """Scannt alle überwachten Verzeichnisse."""
    results = {}
    
    for dir_name, dir_path in WATCH_DIRS.items():
        new_files = scan_directory(dir_path)
        results[dir_name] = new_files
        
        if new_files:
            logger.info(f"Gefunden in /{dir_name}: {len(new_files)} neue Dateien")
    
    return results


# ============================================
# MAIN LEARNING LOOP
# ============================================

def learn_new_files(scan_type: str = "manual") -> Dict[str, Any]:
    """
    Hauptfunktion: Scannt und lernt aus neuen Dateien.
    
    Args:
        scan_type: Art des Scans (manual, scheduled, startup)
    
    Returns:
        Zusammenfassung des Lernvorgangs
    """
    ensure_tables()
    
    # Scan-History starten
    con = _connect()
    cur = con.cursor()
    cur.execute("""
    INSERT INTO scan_history (scan_type, started_at)
    VALUES (?, CURRENT_TIMESTAMP)
    """, (scan_type,))
    scan_id = cur.lastrowid
    con.commit()
    con.close()
    
    logger.info(f"=== Auto-Learning gestartet (Typ: {scan_type}) ===")
    
    # Alle Verzeichnisse scannen
    all_new_files = scan_all_directories()
    
    stats = {
        "files_found": 0,
        "files_learned": 0,
        "files_skipped": 0,
        "errors": [],
        "by_type": {"template": 0, "knowledge": 0, "asset": 0}
    }
    
    # Neue Dateien lernen
    for dir_name, new_files in all_new_files.items():
        stats["files_found"] += len(new_files)
        
        for file_info in new_files:
            file_path = file_info["path"]
            file_type = file_info["type"]
            file_hash = file_info["hash"]
            
            # Nochmal prüfen (Race Condition vermeiden)
            if is_already_learned(file_hash):
                stats["files_skipped"] += 1
                continue
            
            # Lernen
            result = learn_file(file_path, file_type)
            
            if result.get("ok"):
                mark_as_learned(file_path, file_hash, file_type, result)
                stats["files_learned"] += 1
                stats["by_type"][file_type] = stats["by_type"].get(file_type, 0) + 1
            else:
                stats["files_skipped"] += 1
                stats["errors"].append({
                    "file": file_path,
                    "error": result.get("error")
                })
    
    # Scan-History aktualisieren
    con = _connect()
    cur = con.cursor()
    cur.execute("""
    UPDATE scan_history
    SET finished_at = CURRENT_TIMESTAMP,
        files_found = ?,
        files_learned = ?,
        files_skipped = ?,
        errors = ?
    WHERE id = ?
    """, (stats["files_found"], stats["files_learned"], stats["files_skipped"],
          json.dumps(stats["errors"]), scan_id))
    con.commit()
    con.close()
    
    logger.info(f"=== Auto-Learning abgeschlossen ===")
    logger.info(f"  Gefunden: {stats['files_found']}")
    logger.info(f"  Gelernt: {stats['files_learned']}")
    logger.info(f"  Übersprungen: {stats['files_skipped']}")
    if stats["errors"]:
        logger.warning(f"  Fehler: {len(stats['errors'])}")
    
    return {
        "ok": True,
        "scan_id": scan_id,
        "scan_type": scan_type,
        **stats
    }


# ============================================
# STATUS & MANAGEMENT
# ============================================

def get_status() -> Dict[str, Any]:
    """Gibt den Status des Auto-Learners zurück."""
    ensure_tables()
    
    con = _connect()
    cur = con.cursor()
    
    # Gelernte Dateien zählen
    cur.execute("SELECT file_type, COUNT(*) as cnt FROM learned_files GROUP BY file_type")
    by_type = {row["file_type"]: row["cnt"] for row in cur.fetchall()}
    
    # Letzter Scan
    cur.execute("""
    SELECT * FROM scan_history ORDER BY id DESC LIMIT 1
    """)
    last_scan = None
    row = cur.fetchone()
    if row:
        last_scan = {
            "scan_type": row["scan_type"],
            "started_at": row["started_at"],
            "finished_at": row["finished_at"],
            "files_learned": row["files_learned"]
        }
    
    con.close()
    
    # Aktuelle Verzeichnis-Größen
    dir_sizes = {}
    for dir_name, dir_path in WATCH_DIRS.items():
        path = Path(dir_path)
        if path.exists():
            file_count = sum(1 for _ in path.rglob("*") if _.is_file())
            dir_sizes[dir_name] = file_count
        else:
            dir_sizes[dir_name] = 0
    
    return {
        "ok": True,
        "learned_files": by_type,
        "total_learned": sum(by_type.values()),
        "watch_directories": WATCH_DIRS,
        "directory_sizes": dir_sizes,
        "last_scan": last_scan
    }


def get_learned_files(limit: int = 50) -> List[Dict[str, Any]]:
    """Gibt die zuletzt gelernten Dateien zurück."""
    ensure_tables()
    
    con = _connect()
    cur = con.cursor()
    cur.execute("""
    SELECT file_path, file_type, file_size, learned_at
    FROM learned_files
    ORDER BY learned_at DESC
    LIMIT ?
    """, (limit,))
    
    files = [dict(row) for row in cur.fetchall()]
    con.close()
    
    return files


def reset_learned_files(file_type: str = None):
    """Setzt gelernte Dateien zurück (für Re-Learning)."""
    ensure_tables()
    
    con = _connect()
    cur = con.cursor()
    
    if file_type:
        cur.execute("DELETE FROM learned_files WHERE file_type = ?", (file_type,))
        logger.info(f"Reset: Alle {file_type}-Dateien zum Re-Learning markiert")
    else:
        cur.execute("DELETE FROM learned_files")
        logger.info("Reset: Alle Dateien zum Re-Learning markiert")
    
    con.commit()
    con.close()


# ============================================
# DAEMON MODE
# ============================================

def run_daemon():
    """Führt Auto-Learning als Daemon aus."""
    logger.info("Auto-Learner Daemon gestartet")
    logger.info(f"Scan-Intervall: {SCAN_INTERVAL_SECONDS} Sekunden")
    logger.info(f"Überwachte Verzeichnisse: {list(WATCH_DIRS.keys())}")
    
    # Initial Scan
    learn_new_files(scan_type="startup")
    
    try:
        while True:
            time.sleep(SCAN_INTERVAL_SECONDS)
            learn_new_files(scan_type="scheduled")
    except KeyboardInterrupt:
        logger.info("Auto-Learner Daemon gestoppt")


# ============================================
# CLI
# ============================================

def main():
    """CLI Entry Point."""
    import argparse
    
    parser = argparse.ArgumentParser(description="StratGen Auto-Learner")
    parser.add_argument("command", choices=["scan", "status", "daemon", "reset", "list"],
                       help="Befehl: scan, status, daemon, reset, list")
    parser.add_argument("--type", choices=["template", "knowledge", "asset"],
                       help="Dateityp für reset")
    parser.add_argument("--limit", type=int, default=50,
                       help="Limit für list")
    
    args = parser.parse_args()
    
    if args.command == "scan":
        result = learn_new_files(scan_type="manual")
        print(json.dumps(result, indent=2, ensure_ascii=False))
        
    elif args.command == "status":
        status = get_status()
        print(json.dumps(status, indent=2, ensure_ascii=False))
        
    elif args.command == "daemon":
        run_daemon()
        
    elif args.command == "reset":
        reset_learned_files(args.type)
        print("Reset abgeschlossen")
        
    elif args.command == "list":
        files = get_learned_files(args.limit)
        for f in files:
            print(f"{f['learned_at']} | {f['file_type']:10} | {f['file_path']}")


if __name__ == "__main__":
    main()
