
from __future__ import annotations
import os, json, time, sqlite3
from typing import Dict, Any, List

# Optional: python-pptx – wenn nicht installiert, liefern wir eine saubere Fehlermeldung
try:
    from pptx import Presentation  # type: ignore
except Exception as e:  # pragma: no cover
    Presentation = None
    _pptx_import_error = e
else:
    _pptx_import_error = None

DB_PATH = os.path.join("data", "projects.sqlite")

def _connect():
    os.makedirs("data", exist_ok=True)
    con = sqlite3.connect(DB_PATH, timeout=30, isolation_level=None)
    cur = con.cursor()
    # WAL & Sync für weniger Sperren
    cur.execute("PRAGMA journal_mode=WAL")
    cur.execute("PRAGMA synchronous=NORMAL")
    return con

def _ensure_schema():
    con = _connect()
    cur = con.cursor()
    cur.executescript("""
    CREATE TABLE IF NOT EXISTS decks(
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      ts INTEGER NOT NULL,
      source TEXT,
      meta_json TEXT
    );
    CREATE TABLE IF NOT EXISTS deck_slides(
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      deck_id INTEGER NOT NULL,
      idx INTEGER NOT NULL DEFAULT 0,
      kind TEXT,
      title TEXT,
      bullets_json TEXT,
      notes TEXT
    );
    CREATE INDEX IF NOT EXISTS idx_deck_slides_deck ON deck_slides(deck_id);
    CREATE INDEX IF NOT EXISTS idx_deck_slides_kind ON deck_slides(kind);
    CREATE TABLE IF NOT EXISTS style_profiles(
      id INTEGER PRIMARY KEY AUTOINCREMENT,
      created_at INTEGER NOT NULL,
      name TEXT NOT NULL,
      profile_json TEXT NOT NULL
    );
    """)
    con.close()

def _retry_sql(fn, retries=5, wait=0.15):
    last = None
    for i in range(retries):
        try:
            return fn()
        except sqlite3.OperationalError as e:
            last = e
            if "locked" in str(e).lower():
                time.sleep(wait)
                continue
            raise
    if last:
        raise last

def _extract_slide(slide) -> Dict[str, Any]:
    # Titel
    title = ""
    if hasattr(slide, "shapes"):
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                # erster Platzhalter mit Titel ist meist shapes.title
                if getattr(shp, "text_frame", None) and getattr(getattr(shp, "text_frame", None), "text", None):
                    pass
            if getattr(slide, "shapes", None) and getattr(slide.shapes, "title", None):
                try:
                    title = slide.shapes.title.text or ""
                except Exception:
                    pass
    if not title:
        # fallback: erster Text
        try:
            for shp in slide.shapes:
                if getattr(shp, "has_text_frame", False):
                    tx = shp.text_frame.text or ""
                    if tx.strip():
                        title = tx.strip().split("\n",1)[0]
                        break
        except Exception:
            title = ""

    # Bullets
    bullets: List[str] = []
    try:
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False):
                tf = shp.text_frame
                # alle Absätze außer dem Titel
                for para in tf.paragraphs[1:] if slide.shapes.title and tf is getattr(slide.shapes.title, "text_frame", None) else tf.paragraphs:
                    txt = "".join(run.text for run in para.runs) or para.text or ""
                    if txt and txt.strip():
                        bullets.append(txt.strip())
    except Exception:
        pass

    # sehr einfache Heuristik für "kind"
    kind = "title" if not bullets and (title.lower() in ("title", "cover", "deck") or len(title.split())<=3) else "section"

    return {"title": title.strip(), "bullets": bullets, "kind": kind}

def ingest_pptx(path: str, source: str = "manual") -> Dict[str, Any]:
    """Liest ein PPTX ein, schreibt in decks/deck_slides und gibt deck_id + Anzahl zurück."""
    if _pptx_import_error is not None:
        return {
            "ok": False,
            "error": f"python-pptx missing: {_pptx_import_error}. Install with: .venv/bin/pip install python-pptx"
        }
    if not os.path.exists(path):
        return {"ok": False, "error": "file not found", "path": path}

    _ensure_schema()

    prs = Presentation(path)
    slides_data: List[Dict[str, Any]] = []
    for s in prs.slides:
        slides_data.append(_extract_slide(s))

    now = int(time.time())

    def _insert():
        con = _connect()
        cur = con.cursor()
        cur.execute("INSERT INTO decks(ts, source, meta_json) VALUES(?,?,?)", (now, source, json.dumps({"path": path})))
        deck_id = cur.lastrowid
        for idx, sl in enumerate(slides_data):
            cur.execute(
                "INSERT INTO deck_slides(deck_id, idx, kind, title, bullets_json, notes) VALUES(?,?,?,?,?,?)",
                (deck_id, idx, sl.get("kind"), sl.get("title"), json.dumps(sl.get("bullets") or []), "")
            )
        con.close()
        return deck_id

    deck_id = _retry_sql(_insert)
    return {"ok": True, "deck_id": deck_id, "slides": len(slides_data)}

def learn_styles(name: str = "default") -> Dict[str, Any]:
    """Bildet einen sehr einfachen Stil-„Fingerprint“ über alle Slides und legt ihn in style_profiles ab."""
    _ensure_schema()

    def _compute_and_store():
        con = _connect()
        cur = con.cursor()
        # Lade alle Slides
        rows = cur.execute("SELECT kind, bullets_json FROM deck_slides").fetchall()
        n = len(rows)
        avg_len = 0.0
        kind_count: Dict[str, int] = {}
        if n:
            total_bullets = 0
            for k, bj in rows:
                kind = (k or "").strip() or "unknown"
                kind_count[kind] = kind_count.get(kind, 0) + 1
                try:
                    bs = json.loads(bj) if bj else []
                except Exception:
                    bs = []
                total_bullets += len(bs)
            avg_len = total_bullets / n if n else 0.0

        # decks count
        decks = cur.execute("SELECT COUNT(1) FROM decks").fetchone()[0]

        profile = {"ok": True, "avg_len": round(avg_len, 2), "kind_share": kind_count, "decks": decks}
        cur.execute(
            "INSERT INTO style_profiles(created_at, name, profile_json) VALUES(?,?,?)",
            (int(time.time()), name, json.dumps(profile, ensure_ascii=False))
        )
        con.close()
        return {"ok": True, "style_id": None, "profile": profile}

    return _retry_sql(_compute_and_store)

def deck_preview(deck_id: int, limit: int = 5) -> Dict[str, Any]:
    """Gibt die ersten Slides (idx-Sortierung) eines Decks als JSON zurück."""
    _ensure_schema()

    def _load():
        con = _connect()
        cur = con.cursor()
        head = cur.execute("SELECT id, ts, source, meta_json FROM decks WHERE id=?", (deck_id,)).fetchone()
        if not head:
            con.close()
            return {"ok": False, "error": "deck not found", "deck_id": deck_id}
        rows = cur.execute(
            "SELECT idx, kind, title, bullets_json, notes FROM deck_slides WHERE deck_id=? ORDER BY idx ASC LIMIT ?",
            (deck_id, int(limit))
        ).fetchall()
        con.close()
        slides = []
        for idx, kind, title, bj, notes in rows:
            try:
                bullets = json.loads(bj) if bj else []
            except Exception:
                bullets = []
            slides.append({"idx": idx, "kind": kind, "title": title, "bullets": bullets, "notes": notes or ""})
        return {"ok": True, "deck_id": deck_id, "slides": slides}

    return _retry_sql(_load)
