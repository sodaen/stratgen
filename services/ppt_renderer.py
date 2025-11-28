from __future__ import annotations
import re

from pathlib import Path
from typing import Dict, List, Any, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from services.images import parse_img_tokens, render_image

# Visuals (Tabellen/Charts) – optional
try:
    from services.visuals import parse_visual_placeholders, render_table_png, render_chart_png
    HAS_VISUALS = True
except Exception:
    HAS_VISUALS = False

# -------------------- Helpers (Placeholders & Text) --------------------

def _ph(slide, allowed_types: tuple) -> Optional[Any]:
    for shp in slide.shapes:
        if getattr(shp, "is_placeholder", False):
            try:
                ph_type = shp.placeholder_format.type
            except Exception:
                continue
            if ph_type in allowed_types:
                return shp
    return None

def _title(slide):
    return _ph(slide, (PP_PLACEHOLDER.TITLE,))

def _subtitle(slide):
    return _ph(slide, (PP_PLACEHOLDER.SUBTITLE,))

def _body(slide):
    # BODY ist der Content-Platzhalter einer Inhaltsfolie
    return _ph(slide, (PP_PLACEHOLDER.BODY,))

def _set_text(shape, text: str, style: Dict[str, Any], size_pt: int = 20, color_rgb=None):
    if shape is None:
        return
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text or ""
    for r in p.runs:
        r.font.name = style.get("font_body", "Arial")
        r.font.size = Pt(size_pt)
        rgb = color_rgb or style.get("font_color_rgb")
        if rgb:
            try:
                r.font.color.rgb = RGBColor(*rgb)
            except Exception:
                pass

def _set_bullets(shape, lines: List[str], style: Dict[str, Any]):
    if shape is None:
        return
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return
    tf.clear()
    default_lines = ["• [[[TODO: Inhalte ergänzen]]]", "• [[[TODO]]]", "• [[[TODO]]]"]
    first = True
    for line in (lines or default_lines):
        if first:
            p = tf.paragraphs[0]
            p.text = line
            p.level = 0
            first = False
        else:
            p = tf.add_paragraph()
            p.text = line
            p.level = 0
        for r in p.runs:
            r.font.name = style.get("font_body", "Arial")
            r.font.size = Pt(style.get("size_body_pt", 16))
            rgb = style.get("font_color_rgb")
            if rgb:
                try:
                    r.font.color.rgb = RGBColor(*rgb)
                except Exception:
                    pass


def _textbox(prs, slide, left_in, top_in, width_in, height_in):
    return slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))

# -------------------- Footer & Citations --------------------

def _add_footer(prs, slide, txt: str, style: Dict[str, Any]):
    try:
        sw, sh = prs.slide_width.inches, prs.slide_height.inches
        left, top, width, height = 0.5, sh - 0.6, sw - 1.0, 0.35
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = txt
        for r in p.runs:
            r.font.name = style.get("font_body", "Arial")
            r.font.size = Pt(9)
            r.font.color.rgb = RGBColor(120, 120, 120)
    except Exception:
        pass

def _render_citations(prs, slide, style, citations: List[str]):
    if not citations:
        return
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    left, top, width, height = 0.6, sh - 1.2, sw - 1.2, 0.6
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "  ".join([f"[{i+1}] {c}" for i, c in enumerate(citations)])
    for r in p.runs:
        r.font.name = style.get("font_body", "Arial")
        r.font.size = Pt(9)
        r.font.color.rgb = RGBColor(90, 90, 90)

# -------------------- Image placement (Charts/Tables) --------------------

def _place_image(prs, slide, img_path: str, region: str = "right"):
    try:
        sw, sh = prs.slide_width.inches, prs.slide_height.inches
        if region == "right":
            left, top, width, height = sw * 0.55, 1.5, sw * 0.38, sh - 2.2
        elif region == "bottom":
            left, top, width, height = 1.0, sh * 0.58, sw - 2.0, sh * 0.32
        else:
            left, top, width, height = 1.0, sh * 0.55, sw - 2.0, sh * 0.35
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
    except Exception:
        pass

# -------------------- Template Mapping --------------------

def _detect_layouts(prs: Presentation) -> Dict[str, int]:
    """
    Versucht Layout-IDs anhand von Platzhalter-Texten zu erkennen.
    Erwartete Tokens im Layout:
      - Start:  #TITEL / #THEMA / #DATUM / #KUNDENLOGO
      - Section: #ZWISCHENFOLIE
      - Content: #INHALT  (optional #TOPIC / #SUBTOPIC)
      - End:     #ENDE
    Fallbacks: 0=start, 1=content, 2=section, 3=end (falls vorhanden)
    """
    found = {"start": None, "section": None, "content": None, "end": None}
    for idx, layout in enumerate(prs.slide_layouts):
        texts = []
        try:
            for shp in layout.shapes:
                if hasattr(shp, "text_frame") and shp.text_frame:
                    t = (shp.text or "").strip()
                    if t:
                        texts += [t]
        except Exception:
            pass
        blob = " | ".join(texts).upper()
        if any(tok in blob for tok in ["#TITEL", "#THEMA", "#DATUM", "#KUNDENLOGO"]) and found["start"] is None:
            found["start"] = idx
        if "#ZWISCHENFOLIE" in blob and found["section"] is None:
            found["section"] = idx
        if "#INHALT" in blob and found["content"] is None:
            found["content"] = idx
        if "#ENDE" in blob and found["end"] is None:
            found["end"] = idx

    # Fallbacks
    if found["start"] is None:   found["start"] = 0
    if found["content"] is None: found["content"] = min(1, len(prs.slide_layouts)-1)
    if found["section"] is None: found["section"] = min(2, len(prs.slide_layouts)-1)
    if found["end"] is None:     found["end"] = min(3, len(prs.slide_layouts)-1)
    return {k:int(v) for k,v in found.items()}

# -------------------- Core Rendering --------------------

def _render_start_slide(prs, layouts, customer_name, project_title, thema, datum, logo_path, style):
    sld = prs.slides.add_slide(prs.slide_layouts[layouts["start"]])
    # Versuche, Tokens zu befüllen:
    # 1) Titel / Thema / Datum in TITLE/SUBTITLE/BODY-Ersetzungen
    t = _title(sld)
    if t: _set_text(t, project_title or "", style, size_pt=style.get("size_title_pt", 30), color_rgb=style.get("color_title_rgb"))
    sub = _subtitle(sld)
    if sub: _set_text(sub, thema or "", style, size_pt=style.get("size_subtitle_pt", 20), color_rgb=style.get("color_subtitle_rgb"))
    b = _body(sld)
    if b and datum:
        _set_text(b, datum, style, size_pt=14)

    # Falls im Layout Token-Texte stehen, ersetzen wir sie auch:
    for shp in sld.shapes:
        if hasattr(shp, "text_frame") and shp.text_frame:
            txt = shp.text_frame.text or ""
            if "#TITEL" in txt:
                _set_text(shp, (project_title or "").strip(), style, size_pt=style.get("size_title_pt", 30))
            if "#THEMA" in txt:
                _set_text(shp, (thema or "").strip(), style, size_pt=style.get("size_subtitle_pt", 20))
            if "#DATUM" in txt:
                _set_text(shp, (datum or "").strip(), style, size_pt=14)
            if "#KUNDENLOGO" in txt:
                # Wenn kein Logo, setze Kundenname
                if logo_path and Path(logo_path).exists():
                    # Wir versuchen das Bild rechts unten zu platzieren
                    try:
                        sw, sh = prs.slide_width.inches, prs.slide_height.inches
                        sld.shapes.add_picture(logo_path, Inches(sw-2.5), Inches(0.5), width=Inches(2.0))
                        _set_text(shp, "", style, size_pt=12)
                    except Exception:
                        _set_text(shp, customer_name or "", style, size_pt=18)
                else:
                    _set_text(shp, customer_name or "", style, size_pt=18)
    _add_footer(prs, sld, "Modus: START · Platzhalter = [[[TODO]]]", style)
    return sld

def _render_section_slide(prs, layouts, topic, style):
    sld = prs.slides.add_slide(prs.slide_layouts[layouts["section"]])
    # Ersetze #ZWISCHENFOLIE oder setze Titel
    replaced = False
    for shp in sld.shapes:
        if hasattr(shp, "text_frame") and shp.text_frame:
            txt = shp.text_frame.text or ""
            if "#ZWISCHENFOLIE" in txt:
                _set_text(shp, topic, style, size_pt=style.get("size_title_pt", 28))
                replaced = True
    if not replaced:
        t = _title(sld) or _subtitle(sld) or _body(sld)
        _set_text(t, topic, style, size_pt=style.get("size_title_pt", 28))
    _add_footer(prs, sld, f"Modus: SECTION · {topic}", style)
    return sld

def _render_content(prs, slide, style, topic, subtopic, bullets_lines, citations: Optional[List[str]] = None):
    # Topic/Subtopic (wenn vorhanden)
    t = _title(slide)
    if t: _set_text(t, topic or "", style, size_pt=style.get("size_title_pt", 24), color_rgb=style.get("color_title_rgb"))
    st = _subtitle(slide)
    if st: _set_text(st, subtopic or "", style, size_pt=style.get("size_subtitle_pt", 16), color_rgb=style.get("color_subtitle_rgb"))

    body = _body(slide)
    if body is None:
        # Notfall-Textbox, falls kein BODY-Platzhalter existiert
        sw, sh = prs.slide_width.inches, prs.slide_height.inches
        body = _textbox(prs, slide, 0.8, 1.8, sw-1.6, sh-3.0)

    # Visuals & Images aus Bullets extrahieren
    # WARN-Markup extrahieren
    bullets_lines, warn_list = _extract_warnings(bullets_lines)

    # Charts/Tables
    if HAS_VISUALS:
        vis = parse_visual_placeholders(bullets_lines)
        tmp_lines = vis.get("bullets_clean") or bullets_lines
    else:
        vis = {"wants_table": False, "chart_type": None, "data": []}
        tmp_lines = bullets_lines
    # Images
    img_info = parse_img_tokens(tmp_lines)
    bullets_lines = img_info.get("bullets_clean") or tmp_lines
    img_specs = img_info.get("images") or []

    _set_bullets(body, bullets_lines, style)

    # Visual rechts platzieren
    try:
        if HAS_VISUALS:
            if vis.get("wants_table"):
                img = render_table_png(vis.get("data") or [], title=subtopic)
                _place_image(prs, slide, img, region="right")
            elif vis.get("chart_type"):
                img = render_chart_png(vis["chart_type"], vis.get("data") or [], title=subtopic)
                _place_image(prs, slide, img, region="right")
        # Danach optional Bild unten
        if img_specs:
            spec = img_specs[0]
            ipath = render_image(spec["keyword"], style=spec["style"], size=spec["size"])
            _place_image(prs, slide, ipath, region="bottom")
    except Exception:
        pass

    # Warnbanner am Ende der Content-Renderings
    if 'warn_list' in locals() and warn_list:
        _warning_banner(slide, ' | '.join(warn_list))

    # Citations als Fußzeile
    _render_citations(prs, slide, style, citations or [])
    _add_footer(prs, slide, f"Modus: CONTENT · Platzhalter = [[[TODO]]]", style)

def _render_sources_slide(prs, style, citations_list: List[str], layouts: Dict[str, int]):
    if not citations_list:
        return
    try:
        sld = prs.slides.add_slide(prs.slide_layouts[layouts["content"]])
    except Exception:
        sld = prs.slides.add_slide(prs.slide_layouts[0])

    t = _title(sld) or _subtitle(sld) or _body(sld)
    _set_text(t, "Quellen", style, size_pt=style.get("size_title_pt", 24))

    body = _body(sld)
    if body is None:
        sw, sh = prs.slide_width.inches, prs.slide_height.inches
        body = _textbox(prs, sld, 1.0, 2.2, sw-2.0, sh-3.2)

    lines = [f"• {c}" for c in citations_list]
    _set_bullets(body, lines, style)
    _add_footer(prs, sld, "Modus: SOURCES", style)

# -------------------- Public API --------------------

def generate_pptx_from_template(
    out_path: str,
    customer_name: str = "",
    project_title: str = "",
    thema: str = "",
    datum: str = "",
    logo_path: Optional[str] = None,
    agenda: Optional[List[Dict[str, Any]]] = None,
    content_map: Optional[Dict[str, Dict[str, Any]]] = None,
    mode_hint: str = "facts",
    style: Optional[Dict[str, Any]] = None,
) -> str:
    """
    Erwartet eine Vorlage unter styles/master.pptx
    Layouts:
      start:    Titelfolie (#TITEL/#THEMA/#DATUM/#KUNDENLOGO)
      section:  Kapitelwechsel (#ZWISCHENFOLIE)
      content:  Inhalt (#INHALT [+ #TOPIC/#SUBTOPIC optional])
      end:      Abschluss (#ENDE)
    """
    style = style or {"font_body":"Arial","size_title_pt":28,"size_subtitle_pt":16,"size_body_pt":16}

    tpl = Path("styles/master.pptx")
    if tpl.exists():
        prs = Presentation(tpl.as_posix())
    else:
        prs = Presentation()

    layouts = _detect_layouts(prs)

    # 1) Startfolie
    _render_start_slide(prs, layouts, customer_name, project_title, thema, datum, logo_path, style)

    # 2) Agenda-Folie (als Content-Layout)
    agenda = agenda or []
    agenda_titles = [f"{i+1}. {item.get('topic','')}" for i, item in enumerate(agenda) if item.get("topic")]
    sld_agenda = prs.slides.add_slide(prs.slide_layouts[layouts["content"]])
    _render_content(prs, sld_agenda, style, topic="Agenda", subtopic="", bullets_lines=[f"• {t}" for t in agenda_titles], citations=[])

    # Inhalte sammeln (auch für Quellen-Folie)
    all_citations: List[str] = []

    # 3) Kapitel & Inhaltsfolien
    for item in agenda:
        topic = item.get("topic","").strip() or "Kapitel"
        _render_section_slide(prs, layouts, topic, style)
        subs = item.get("subtopics") or []
        if not subs:
            subs = ["Seite 1","Seite 2","Seite 3"]

        for sub in subs:
            sld = prs.slides.add_slide(prs.slide_layouts[layouts["content"]])
            # Inhalte
            obj = None
            if content_map and topic in content_map:
                obj = content_map[topic].get(sub)
            if isinstance(obj, dict):
                bullets = obj.get("bullets") or ["• [[[TODO: Inhalte ergänzen]]]", "• [[[TODO]]]", "• [[[TODO]]]"]
                citations = obj.get("citations") or []
            else:
                text = obj if isinstance(obj, str) else ""
                bullets = [ln.strip() for ln in text.splitlines() if ln.strip()] or ["• [[[TODO: Inhalte ergänzen]]]", "• [[[TODO]]]", "• [[[TODO]]]"]
                citations = []

            _render_content(prs, sld, style, topic=topic, subtopic=sub, bullets_lines=bullets, citations=citations)
            for c in citations:
                if c not in all_citations:
                    all_citations.append(c)

    # 4) Quellenfolie vor Endfolie
    _render_sources_slide(prs, style, all_citations[:24], layouts)

    # 5) Endfolie (falls im Template vorhanden)
    try:
        sld_end = prs.slides.add_slide(prs.slide_layouts[layouts["end"]])
        for shp in sld_end.shapes:
            if hasattr(shp, "text_frame") and shp.text_frame and "#ENDE" in (shp.text_frame.text or ""):
                _set_text(shp, "Danke!", style, size_pt=style.get("size_title_pt", 28))
        _add_footer(prs, sld_end, "Modus: END", style)
    except Exception:
        pass

    # Warnbanner am Ende der Content-Renderings
    if 'warn_list' in locals() and warn_list:
        _warning_banner(slide, ' | '.join(warn_list))

    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out.as_posix())
    return out.as_posix()


def _extract_warnings(bullets_lines: List[str]):
    warn = []
    clean = []
    pat = re.compile(r"^\s*[•\-]?\s*\[\[\[WARN:(.*?)\]\]\]\s*$", re.I)
    for ln in bullets_lines or []:
        m = pat.match(ln or "")
        if m:
            warn.append(m.group(1).strip())
        else:
            clean.append(ln)
    return clean, warn

def _warning_banner(slide, text: str):
    if not text:
        return
    # schmaler, roter Banner oben
    left = Inches(0.5); top = Inches(0.2); width = Inches(9.0); height = Inches(0.6)
    box = slide.shapes.add_textbox(left, top, width, height)
    box.fill.solid()
    try:
        box.fill.fore_color.rgb = RGBColor(200, 40, 40)
    except Exception:
        pass

    # Warnbanner am Ende der Content-Renderings
    if 'warn_list' in locals() and warn_list:
        _warning_banner(slide, ' | '.join(warn_list))
    try:
        box.line.fill.background()
    except Exception:
        pass

    # Warnbanner am Ende der Content-Renderings
    if 'warn_list' in locals() and warn_list:
        _warning_banner(slide, ' | '.join(warn_list))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"⚠ {text}"
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.name = "Arial"
        r.font.size = Pt(16)
        try:
            r.font.color.rgb = RGBColor(255,255,255)
        except Exception:
            pass
