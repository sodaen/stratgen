# -*- coding: utf-8 -*-
"""
services/content_intelligence.py
================================
Features:
- Evidence Linker
- Slide Complexity Scorer
- Template Recommender
- Meeting Context Adapter
- Knowledge Gap Detector

Author: StratGen Agent V3.6
"""
from __future__ import annotations
import os
import re
import json
import sqlite3
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field, asdict
from datetime import datetime

# ============================================
# DATA CLASSES
# ============================================

@dataclass
class EvidenceLink:
    """Eine Verlinkung zwischen Behauptung und Beleg."""
    claim: str
    evidence: str
    source: str
    confidence: float
    source_type: str  # knowledge_base, template, uploaded


@dataclass
class ComplexityScore:
    """Komplexitätsbewertung eines Slides."""
    slide_index: int
    overall_score: float  # 0=einfach, 1=komplex
    word_count: int
    concept_count: int
    reading_time_seconds: int
    recommendation: str


@dataclass
class TemplateRecommendation:
    """Eine Template-Empfehlung."""
    template_name: str
    file_path: str
    match_score: float
    match_reasons: List[str]
    slide_count: int


@dataclass
class KnowledgeGap:
    """Eine identifizierte Wissenslücke."""
    topic: str
    gap_type: str  # missing_data, outdated, incomplete
    severity: str  # low, medium, high
    suggestion: str


# ============================================
# EVIDENCE LINKER
# ============================================

def find_evidence_for_claim(
    claim: str,
    knowledge_base_results: List[Dict[str, Any]] = None
) -> List[EvidenceLink]:
    """
    Findet Belege für eine Behauptung.
    
    Args:
        claim: Die zu belegende Behauptung
        knowledge_base_results: Suchergebnisse aus Knowledge Base
    
    Returns:
        Liste von EvidenceLink-Objekten
    """
    evidence_links = []
    
    if not knowledge_base_results:
        # Knowledge Base durchsuchen
        try:
            from services.knowledge_enhanced import search_knowledge_base
            knowledge_base_results = search_knowledge_base(claim, k=5)
        except ImportError:
            knowledge_base_results = []
    
    for result in knowledge_base_results:
        # Relevanz prüfen
        score = result.get("score", 0.5)
        
        if score > 0.4:
            evidence_links.append(EvidenceLink(
                claim=claim,
                evidence=result.get("content", result.get("snippet", ""))[:200],
                source=result.get("source", result.get("path", "Unknown")),
                confidence=score,
                source_type="knowledge_base"
            ))
    
    return evidence_links


def link_all_claims(slides: List[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Verlinkt alle Behauptungen in Slides mit Belegen.
    
    Args:
        slides: Liste der Slides
    
    Returns:
        Dictionary mit verlinkten Claims
    """
    linked_claims = []
    unlinked_claims = []
    
    for idx, slide in enumerate(slides):
        bullets = slide.get("bullets", [])
        
        for bullet in bullets:
            # Prüfe ob es eine Behauptung ist (enthält Zahlen, Aussagen, etc.)
            is_claim = bool(re.search(r'\d+%|\d+€|steiger|reduzier|verbess|führend|beste', bullet.lower()))
            
            if is_claim:
                evidence = find_evidence_for_claim(bullet)
                
                if evidence:
                    linked_claims.append({
                        "slide": idx,
                        "claim": bullet[:100],
                        "evidence_count": len(evidence),
                        "best_evidence": evidence[0].evidence[:100] if evidence else None,
                        "source": evidence[0].source if evidence else None
                    })
                else:
                    unlinked_claims.append({
                        "slide": idx,
                        "claim": bullet[:100],
                        "warning": "Kein Beleg gefunden"
                    })
    
    return {
        "ok": True,
        "total_claims": len(linked_claims) + len(unlinked_claims),
        "linked": len(linked_claims),
        "unlinked": len(unlinked_claims),
        "linked_claims": linked_claims,
        "unlinked_claims": unlinked_claims,
        "evidence_score": len(linked_claims) / max(1, len(linked_claims) + len(unlinked_claims))
    }


# ============================================
# SLIDE COMPLEXITY SCORER
# ============================================

def score_slide_complexity(slide: Dict[str, Any], slide_index: int = 0) -> ComplexityScore:
    """
    Bewertet die Komplexität eines Slides.
    
    Args:
        slide: Der Slide
        slide_index: Index des Slides
    
    Returns:
        ComplexityScore-Objekt
    """
    title = slide.get("title", "")
    bullets = slide.get("bullets", [])
    
    # Metriken berechnen
    all_text = f"{title} {' '.join(bullets)}"
    words = all_text.split()
    word_count = len(words)
    
    # Konzepte zählen (Großgeschriebene Wörter, Fachbegriffe)
    concepts = re.findall(r'\b[A-ZÄÖÜ][a-zäöüß]*(?:[A-ZÄÖÜ][a-zäöüß]*)*\b', all_text)
    concept_count = len(set(concepts))
    
    # Lesezeit (Durchschnitt: 200 Wörter/Minute für Präsentationen)
    reading_time = int(word_count / 200 * 60)
    
    # Komplexitätsscore berechnen
    complexity = 0.0
    
    # Wortanzahl (>100 = komplex)
    if word_count > 100:
        complexity += 0.3
    elif word_count > 50:
        complexity += 0.15
    
    # Bullet-Anzahl (>6 = komplex)
    if len(bullets) > 6:
        complexity += 0.3
    elif len(bullets) > 4:
        complexity += 0.15
    
    # Konzeptdichte
    if concept_count > 10:
        complexity += 0.2
    elif concept_count > 5:
        complexity += 0.1
    
    # Durchschnittliche Bullet-Länge
    if bullets:
        avg_bullet_len = sum(len(b) for b in bullets) / len(bullets)
        if avg_bullet_len > 80:
            complexity += 0.2
    
    complexity = min(1.0, complexity)
    
    # Empfehlung
    if complexity > 0.7:
        recommendation = "Slide ist zu komplex. Aufteilen oder vereinfachen."
    elif complexity > 0.5:
        recommendation = "Slide ist etwas dicht. Evtl. 1-2 Punkte entfernen."
    else:
        recommendation = "Slide hat gute Komplexität."
    
    return ComplexityScore(
        slide_index=slide_index,
        overall_score=round(complexity, 2),
        word_count=word_count,
        concept_count=concept_count,
        reading_time_seconds=reading_time,
        recommendation=recommendation
    )


def score_deck_complexity(slides: List[Dict[str, Any]]) -> Dict[str, Any]:
    """
    Bewertet die Komplexität des gesamten Decks.
    
    Args:
        slides: Liste der Slides
    
    Returns:
        Dictionary mit Komplexitätsanalyse
    """
    scores = [score_slide_complexity(slide, idx) for idx, slide in enumerate(slides)]
    
    avg_complexity = sum(s.overall_score for s in scores) / max(1, len(scores))
    complex_slides = [s for s in scores if s.overall_score > 0.6]
    
    total_reading_time = sum(s.reading_time_seconds for s in scores)
    
    return {
        "ok": True,
        "total_slides": len(slides),
        "average_complexity": round(avg_complexity, 2),
        "complex_slides_count": len(complex_slides),
        "total_reading_time_minutes": round(total_reading_time / 60, 1),
        "slides": [asdict(s) for s in scores],
        "recommendations": [s.recommendation for s in complex_slides]
    }


# ============================================
# TEMPLATE RECOMMENDER
# ============================================

def recommend_template(
    briefing: str,
    topic: str,
    industry: str = "",
    deck_size: str = "medium"
) -> List[TemplateRecommendation]:
    """
    Empfiehlt das beste Template basierend auf Briefing.
    
    Args:
        briefing: Das Projekt-Briefing
        topic: Hauptthema
        industry: Branche
        deck_size: Deckgröße
    
    Returns:
        Liste von TemplateRecommendation-Objekten
    """
    recommendations = []
    
    raw_dir = Path(os.getenv("STRATGEN_RAW_DIR", "data/raw"))
    if not raw_dir.exists():
        return recommendations
    
    keywords = set(re.findall(r'\b[A-ZÄÖÜa-zäöüß]{4,}\b', f"{briefing} {topic} {industry}".lower()))
    
    # Deck-Size zu Slide-Count
    size_ranges = {
        "short": (3, 8),
        "medium": (8, 15),
        "long": (15, 50)
    }
    target_range = size_ranges.get(deck_size, (8, 15))
    
    for pptx_file in raw_dir.glob("**/*.pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(str(pptx_file))
            slide_count = len(prs.slides)
            
            # Slide-Count Match
            size_match = target_range[0] <= slide_count <= target_range[1]
            
            # Keyword Match im Dateinamen
            filename_lower = pptx_file.stem.lower()
            filename_matches = sum(1 for kw in keywords if kw in filename_lower)
            
            # Content Match (Stichprobe)
            content_keywords = []
            for slide in list(prs.slides)[:5]:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content_keywords.extend(re.findall(r'\b[a-zäöüß]{4,}\b', shape.text.lower()))
            
            content_matches = sum(1 for kw in keywords if kw in content_keywords)
            
            # Score berechnen
            score = 0.0
            reasons = []
            
            if size_match:
                score += 0.3
                reasons.append(f"Passende Größe ({slide_count} Slides)")
            
            if filename_matches > 0:
                score += min(0.3, filename_matches * 0.1)
                reasons.append(f"{filename_matches} Keywords im Namen")
            
            if content_matches > 0:
                score += min(0.4, content_matches * 0.05)
                reasons.append(f"{content_matches} Keywords im Inhalt")
            
            if score > 0.2:
                recommendations.append(TemplateRecommendation(
                    template_name=pptx_file.stem,
                    file_path=str(pptx_file),
                    match_score=round(score, 2),
                    match_reasons=reasons,
                    slide_count=slide_count
                ))
                
        except Exception:
            continue
    
    # Nach Score sortieren
    recommendations.sort(key=lambda x: x.match_score, reverse=True)
    
    return recommendations[:5]


# ============================================
# MEETING CONTEXT ADAPTER
# ============================================

MEETING_CONTEXTS = {
    "pitch": {
        "max_slides": 10,
        "max_duration_minutes": 15,
        "focus": ["problem", "solution", "benefits", "cta"],
        "style": "überzeugend",
        "recommendations": [
            "Starker Hook am Anfang",
            "Klare Value Proposition",
            "Eindeutiger Call-to-Action"
        ]
    },
    "board": {
        "max_slides": 15,
        "max_duration_minutes": 30,
        "focus": ["executive_summary", "kpis", "risks", "recommendations"],
        "style": "analytisch",
        "recommendations": [
            "Executive Summary auf Slide 1",
            "Zahlen und KPIs prominent",
            "Risiken proaktiv adressieren"
        ]
    },
    "workshop": {
        "max_slides": 30,
        "max_duration_minutes": 120,
        "focus": ["agenda", "content", "exercises", "discussion"],
        "style": "interaktiv",
        "recommendations": [
            "Interaktive Elemente einbauen",
            "Diskussionspunkte markieren",
            "Pausen einplanen"
        ]
    },
    "webinar": {
        "max_slides": 25,
        "max_duration_minutes": 60,
        "focus": ["hook", "content", "engagement", "qa"],
        "style": "engaging",
        "recommendations": [
            "Visuelle Ankerpunkte alle 3-4 Slides",
            "Polling-Fragen einbauen",
            "Q&A Zeit einplanen"
        ]
    },
    "training": {
        "max_slides": 40,
        "max_duration_minutes": 180,
        "focus": ["learning_objectives", "content", "exercises", "summary"],
        "style": "didaktisch",
        "recommendations": [
            "Lernziele am Anfang",
            "Wiederholungen einbauen",
            "Übungen nach jedem Abschnitt"
        ]
    },
    "update": {
        "max_slides": 8,
        "max_duration_minutes": 10,
        "focus": ["status", "progress", "blockers", "next_steps"],
        "style": "kompakt",
        "recommendations": [
            "Nur Wesentliches",
            "Ampelsystem für Status",
            "Klare Action Items"
        ]
    }
}


def adapt_to_meeting_context(
    slides: List[Dict[str, Any]],
    meeting_type: str
) -> Dict[str, Any]:
    """
    Passt Deck an Meeting-Kontext an.
    
    Args:
        slides: Liste der Slides
        meeting_type: Art des Meetings
    
    Returns:
        Dictionary mit Anpassungsempfehlungen
    """
    context = MEETING_CONTEXTS.get(meeting_type.lower(), MEETING_CONTEXTS["pitch"])
    
    current_count = len(slides)
    target_count = context["max_slides"]
    
    adaptations = []
    
    # Slide-Anzahl
    if current_count > target_count:
        adaptations.append({
            "type": "reduce_slides",
            "message": f"Reduzieren von {current_count} auf max. {target_count} Slides",
            "action": f"Entferne {current_count - target_count} Slides oder konsolidiere"
        })
    
    # Fokus-Slides prüfen
    focus_types = context["focus"]
    slide_types = [s.get("type", "content") for s in slides]
    
    missing_focus = [f for f in focus_types if f not in slide_types]
    if missing_focus:
        adaptations.append({
            "type": "missing_focus",
            "message": f"Fehlende wichtige Slide-Typen: {', '.join(missing_focus)}",
            "action": "Diese Slides hinzufügen"
        })
    
    # Zeit-Schätzung
    estimated_time = current_count * 2  # ~2 Min pro Slide
    if estimated_time > context["max_duration_minutes"]:
        adaptations.append({
            "type": "too_long",
            "message": f"Geschätzte Zeit: {estimated_time} Min (max: {context['max_duration_minutes']})",
            "action": "Kürzen oder schneller präsentieren"
        })
    
    return {
        "ok": True,
        "meeting_type": meeting_type,
        "context": {
            "max_slides": context["max_slides"],
            "max_duration": context["max_duration_minutes"],
            "style": context["style"],
            "focus_areas": context["focus"]
        },
        "current_slides": current_count,
        "adaptations_needed": len(adaptations),
        "adaptations": adaptations,
        "recommendations": context["recommendations"]
    }


# ============================================
# KNOWLEDGE GAP DETECTOR
# ============================================

def detect_knowledge_gaps(
    briefing: str,
    topic: str,
    industry: str = ""
) -> List[KnowledgeGap]:
    """
    Erkennt Wissenslücken für ein Thema.
    
    Args:
        briefing: Das Projekt-Briefing
        topic: Hauptthema
        industry: Branche
    
    Returns:
        Liste von KnowledgeGap-Objekten
    """
    gaps = []
    
    # Knowledge Base durchsuchen
    try:
        from services.knowledge_enhanced import search_knowledge_base
        results = search_knowledge_base(f"{topic} {industry}", k=10)
        has_knowledge = len(results) > 0
        result_count = len(results)
    except ImportError:
        has_knowledge = False
        result_count = 0
    
    # Gap Detection basierend auf Suchergebnissen
    if result_count == 0:
        gaps.append(KnowledgeGap(
            topic=topic,
            gap_type="missing_data",
            severity="high",
            suggestion=f"Keine Daten zu '{topic}' gefunden. Laden Sie relevante Dokumente hoch."
        ))
    elif result_count < 3:
        gaps.append(KnowledgeGap(
            topic=topic,
            gap_type="incomplete",
            severity="medium",
            suggestion=f"Nur {result_count} Quellen zu '{topic}'. Mehr Daten würden bessere Ergebnisse liefern."
        ))
    
    # Branchenspezifische Gaps
    if industry and industry.lower() not in str(results).lower():
        gaps.append(KnowledgeGap(
            topic=f"{industry}-spezifische Daten",
            gap_type="missing_data",
            severity="medium",
            suggestion=f"Keine branchenspezifischen Daten für '{industry}' gefunden."
        ))
    
    # Zahlen/Statistiken Gap
    keywords_needing_data = ["roi", "kosten", "einsparung", "prozent", "statistik", "benchmark"]
    needs_numbers = any(kw in briefing.lower() for kw in keywords_needing_data)
    
    if needs_numbers and result_count < 5:
        gaps.append(KnowledgeGap(
            topic="Zahlen und Statistiken",
            gap_type="incomplete",
            severity="medium",
            suggestion="Das Briefing erfordert Zahlen. Stellen Sie sicher, dass relevante Statistiken in der Knowledge Base sind."
        ))
    
    return gaps


# ============================================
# API FUNCTIONS
# ============================================

def check_status() -> Dict[str, Any]:
    """Gibt den Status der Content Intelligence zurück."""
    return {
        "ok": True,
        "features": [
            "evidence_linker",
            "complexity_scorer", 
            "template_recommender",
            "meeting_adapter",
            "knowledge_gap_detector"
        ],
        "meeting_types": list(MEETING_CONTEXTS.keys())
    }
