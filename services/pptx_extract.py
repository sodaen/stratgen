from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List

try:
    from pptx import Presentation
except Exception:  # falls python-pptx nicht installiert ist
    Presentation = None


def extract_pptx_text(pptx_path: str | Path) -> Dict[str, Any]:
    """
    Sehr einfache Extraktion:
    - jede Folie durchlaufen
    - alle Textframes zusammensetzen
    - sehr grobe Rollen-Heuristik (title, agenda, section, closing)
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"pptx not found: {pptx_path}")

    if Presentation is None:
        # Fallback: wir können nix parsen, geben aber ein Minimalobjekt zurück
        return {
            "ok": False,
            "note": "python-pptx not available",
            "slides": [],
        }

    prs = Presentation(str(pptx_path))
    slides_out: List[Dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                # alle Absätze zusammensetzen
                paras = [p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()]
                if paras:
                    texts.extend(paras)
            elif hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    texts.append(txt)

        full_text = "\n".join(texts).strip()

        # ganz simple Heuristik:
        role = "section"
        lt = full_text.lower()
        if idx == 1:
            role = "title"
        elif "agenda" in lt or "inhalt" in lt or "contents" in lt:
            role = "agenda"
        elif "next steps" in lt or "nächste schritte" in lt:
            role = "closing"

        slides_out.append(
            {
                "index": idx,
                "role": role,
                "text": full_text,
            }
        )

    return {
        "ok": True,
        "slides": slides_out,
        "slide_count": len(slides_out),
    }


def write_flat_txt_from_extract(extract: Dict[str, Any], target: str | Path) -> Path:
    """
    Nimmt das Ergebnis von extract_pptx_text(...) und schreibt eine flache TXT,
    damit das bestehende knowledge.scan_dir(...) sie sieht.
    """
    target = Path(target)
    lines: List[str] = []
    lines.append("# PPTX EXTRACT")
    lines.append(f"Slides: {extract.get('slide_count', 0)}")
    lines.append("")

    for s in extract.get("slides", []):
        lines.append(f"--- Slide {s.get('index')} ({s.get('role')}):")
        txt = s.get("text") or ""
        lines.append(txt)
        lines.append("")

    target.write_text("\n".join(lines), encoding="utf-8")
    return target
