#!/usr/bin/env python3
"""
Patch: Ersetzt phase_render in agent_v3_api.py
"""
import re

filepath = "/home/sodaen/stratgen/backend/agent_v3_api.py"

with open(filepath, "r") as f:
    content = f.read()

old_render = '''def phase_render(
    req: AgentV3Request,
    slides: List[Dict[str, Any]],
    run_id: str
) -> Dict[str, Any]:
    """
    Phase 7: Rendert das Deck zu PPTX/PDF.
    """
    result = {
        "pptx_url": None,
        "pdf_url": None,
        "project_id": None,
    }
    
    if not req.export_pptx:
        return result
    
    # Projekt erstellen
    project_data = {
        "title": req.topic,
        "customer_name": req.customer_name,
        "topic": req.topic,
        "brief": req.brief,
        "outline": {
            "title": req.topic,
            "sections": slides
        },
        "meta": {
            "slide_plan": slides,
            "run_id": run_id,
        }
    }
    
    # Projekt speichern
    resp = _http_post("/projects/save", project_data)
    if resp and resp.get("project"):
        project_id = resp["project"].get("id")
        result["project_id"] = project_id
        
        # PPTX rendern
        if project_id:
            render_resp = _http_post(f"/pptx/render_from_project/{project_id}")
            if render_resp:
                pptx_path = render_resp.get("path") or render_resp.get("url")
                if pptx_path:
                    filename = Path(pptx_path).name
                    result["pptx_url"] = f"/exports/download/{filename}"
    
    return result'''

new_render = '''def phase_render(
    req: AgentV3Request,
    slides: List[Dict[str, Any]],
    run_id: str
) -> Dict[str, Any]:
    """
    Phase 7: Rendert das Deck zu PPTX/PDF.
    Baut PPTX direkt ohne externe API-Calls (vermeidet Worker-Sync-Problem).
    """
    result = {
        "pptx_url": None,
        "pdf_url": None,
        "project_id": run_id,
    }
    
    if not req.export_pptx:
        return result
    
    # Direkt PPTX bauen (ohne Umweg über /projects/save)
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
    except ImportError:
        return result
    
    prs = Presentation()
    
    for slide_data in slides:
        title = str(slide_data.get("title", "Slide")).strip() or "Slide"
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("notes", "")
        slide_type = slide_data.get("type", "content")
        
        # Layout wählen (0=Title, 1=Title+Content, 5=Blank)
        if slide_type == "title":
            layout = prs.slide_layouts[0]
        elif not bullets:
            layout = prs.slide_layouts[5]
        else:
            layout = prs.slide_layouts[1]
        
        slide = prs.slides.add_slide(layout)
        
        # Titel setzen
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Bullets setzen
        if bullets and len(slide.shapes.placeholders) > 1:
            try:
                body = slide.shapes.placeholders[1].text_frame
                body.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        body.text = str(bullet)
                    else:
                        p = body.add_paragraph()
                        p.text = str(bullet)
                        p.level = 0
            except Exception:
                pass
        
        # Speaker Notes
        if notes and slide.notes_slide:
            try:
                slide.notes_slide.notes_text_frame.text = str(notes)
            except Exception:
                pass
        
        # Chart einfügen wenn vorhanden
        chart_path = slide_data.get("chart")
        if chart_path and os.path.exists(chart_path):
            try:
                slide.shapes.add_picture(
                    chart_path,
                    Inches(5.5),
                    Inches(2.5),
                    width=Inches(4)
                )
            except Exception:
                pass
    
    # Speichern
    exports_dir = Path("data/exports")
    exports_dir.mkdir(parents=True, exist_ok=True)
    
    ts = int(time.time())
    safe_topic = "".join(c if c.isalnum() else "_" for c in req.topic[:30])
    out_name = f"v3_{safe_topic}_{ts}.pptx"
    out_path = exports_dir / out_name
    
    prs.save(str(out_path))
    
    result["pptx_url"] = f"/exports/download/{out_name}"
    result["pptx_path"] = str(out_path)
    
    return result'''

if old_render in content:
    content = content.replace(old_render, new_render)
    with open(filepath, "w") as f:
        f.write(content)
    print("✓ phase_render erfolgreich gepatcht!")
else:
    print("✗ Konnte alte phase_render nicht finden")
    print("  Prüfe ob die Funktion anders aussieht")
